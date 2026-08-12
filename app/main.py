from __future__ import annotations

import asyncio
import csv
import hashlib
import hmac
import json
import os
import re
import time
import uuid
import zipfile
from datetime import datetime, timezone
from io import BytesIO
from pathlib import Path
from typing import Any
from urllib.parse import urljoin, urlparse

import httpx
from fastapi import FastAPI, File, HTTPException, Request, UploadFile
from fastapi.responses import FileResponse, StreamingResponse
from fastapi.staticfiles import StaticFiles

from app import db
from app.schemas import (
    AgentCreate, AgentUpdate, ApprovalRequest, McpServerCreate, McpServerUpdate,
    ExpertInstallRequest, ExpertMemberRetryRequest, ExpertTeamCreate, ExpertTeamRunCreate,
    ExpertTeamUpdate, ExpertTemplateCreate, ExpertTemplateUpdate,
    ConversationSummaryUpdate, LoopCreate, LoopUpdate, MemoryCreate, MemoryUpdate, ModelConfigCreate, ModelConfigUpdate, RemoteInstall, SkillCreate, SkillFileUpdate, SkillPathInstall, SkillUpdate,
    CheckpointRestoreRequest, PolicyRuleCreate, PolicyRuleUpdate, TaskCommandRequest,
    TaskCreate, TaskResumeRequest, ToolInvokeRequest,
)
from app.seed import seed_agents
from app.services.agent_runtime import AgentRuntime, create_task_record
from app.services.context_service import ContextService, ExecutionScope, MemoryNotFoundError
from app.services.conversation_summary_service import ConversationSummaryConflictError
from app.services.expert_team_service import (
    ExpertConflictError, ExpertNotFoundError, ExpertPermissionError,
    ExpertTeamService, ExpertValidationError,
)
from app.services.event_bus import emit
from app.services.mcp_gateway import (
    ARTIFACT_DIR,
    McpGateway,
    ToolError,
    presentation_generation_status,
    resolve_artifact_path,
)
from app.services.loop_scheduler import (
    IdempotencyConflictError, LoopScheduler, create_webhook_event, next_schedule_at,
    serialize_loop, serialize_notification, serialize_run, serialize_trigger_event,
    validate_trigger_config,
)
from app.services.model_gateway import ModelGateway
from app.services.network_policy import (
    env_flag,
    outbound_network_enabled,
    require_outbound_network,
    validate_outbound_http_url,
)
from app.services.policy_engine import PolicyConfigurationError, PolicyEngine, PolicyRule
from app.services.secret_store import secret_store
from app.services.skill_registry import SkillRegistry, referenced_package_files
from app.services.task_state import StateNotFoundError, TaskStateError, TaskStateService

BASE_DIR = Path(__file__).resolve().parents[1]
WEB_DIR = BASE_DIR / "web"
UPLOAD_DIR = Path(os.getenv("APP_UPLOAD_DIR", str(BASE_DIR / "data" / "uploads")))

app = FastAPI(title="AgentNexus", version="0.1.0")

skill_registry = SkillRegistry()
mcp_gateway = McpGateway()
model_gateway = ModelGateway()
task_state = TaskStateService(auto_init=False)
policy_engine = PolicyEngine(
    http_enabled=outbound_network_enabled() and env_flag("APP_ALLOW_HTTP_POLICY"),
    http_allowlist=[
        item.strip()
        for item in os.getenv("APP_HTTP_POLICY_ALLOWLIST", "").split(",")
        if item.strip()
    ],
)
context_service = ContextService(auto_init=False)


def _is_env_name(value: str) -> bool:
    return bool(re.fullmatch(r"[A-Za-z_][A-Za-z0-9_]*", value or ""))


def _migrate_legacy_model_keys() -> None:
    """Move keys accidentally entered in the old env-name field into encrypted storage."""
    for row in db.query_all("SELECT id, api_key_env, api_key_ciphertext FROM model_configs"):
        legacy = str(row.get("api_key_env") or "").strip()
        if legacy and not _is_env_name(legacy) and not row.get("api_key_ciphertext"):
            db.execute(
                "UPDATE model_configs SET api_key_env = '', api_key_ciphertext = ?, updated_at = ? WHERE id = ?",
                (secret_store.encrypt(legacy), db.utc_now(), row["id"]),
            )


def _remote_install_flag() -> bool:
    return env_flag("APP_ALLOW_REMOTE_INSTALL")


def _remote_install_url(url: str) -> str:
    require_outbound_network("下载链接安装", error_type=ValueError)
    if not _remote_install_flag():
        raise ValueError("下载链接安装尚未开启，请由管理员设置 APP_ALLOW_REMOTE_INSTALL=true 后重启平台")
    if urlparse(url).scheme.lower() != "https":
        raise ValueError("下载链接必须使用 HTTPS")
    return validate_outbound_http_url(
        url,
        capability="下载链接安装",
        allowlist_env="APP_REMOTE_INSTALL_HOST_ALLOWLIST",
        require_allowlist=True,
        allow_query=True,
        error_type=ValueError,
    )


async def _download_remote_install(url: str, max_bytes: int) -> tuple[bytes, str]:
    checked = _remote_install_url(url)
    async with httpx.AsyncClient(timeout=30, follow_redirects=False) as client:
        response = await client.get(checked, headers={"User-Agent": "AgentNexus/0.1"})
        response.raise_for_status()
    if len(response.content) > max_bytes:
        raise ValueError("远程安装包超过大小限制")
    filename = Path(urlparse(str(response.url)).path).name or "download"
    return response.content, filename


def _install_skill_bytes(raw: bytes, filename: str) -> dict[str, Any]:
    if len(raw) > 2 * 1024 * 1024:
        raise ValueError("Skill 安装包不能超过 2MB")
    if filename.lower().endswith(".zip") or raw.startswith(b"PK\x03\x04"):
        with zipfile.ZipFile(BytesIO(raw)) as archive:
            files: dict[str, bytes] = {}
            for info in archive.infolist():
                if info.is_dir():
                    continue
                path = info.filename.replace("\\", "/")
                if path.startswith("/") or ".." in Path(path).parts:
                    raise ValueError("ZIP 包含不安全路径")
                if (info.external_attr >> 16) & 0o170000 == 0o120000:
                    raise ValueError("ZIP 不允许包含符号链接")
                files[path] = archive.read(info)
            candidates = [n for n in files if Path(n).name == "SKILL.md"]
            if not candidates:
                raise ValueError("ZIP 中未找到 SKILL.md")
            if len(candidates) > 1:
                raise ValueError("ZIP 中只能包含一个 Skill 包")
            fallback = Path(candidates[0]).parent.name or Path(filename).stem
            return skill_registry.install_package(files, fallback_id=fallback)
    else:
        content = raw.decode("utf-8")
        fallback = Path(filename).stem if filename.lower().endswith(".md") else "downloaded_skill"
    return skill_registry.install_content(content, fallback_id=fallback)


async def _install_skill_remote_url(url: str) -> dict[str, Any]:
    raw, filename = await _download_remote_install(url, 2 * 1024 * 1024)
    if not filename.lower().endswith(".zip") and not raw.startswith(b"PK\x03\x04"):
        content = raw.decode("utf-8")
        files: dict[str, bytes] = {"SKILL.md": raw}
        total = len(raw)
        for relative in referenced_package_files(content):
            stored_relative = relative
            try:
                child, _ = await _download_remote_install(urljoin(url, relative), 1024 * 1024)
            except (ValueError, httpx.HTTPError):
                alternate = str(Path(relative).with_name(Path(relative).name.lower())).replace("\\", "/")
                if alternate == relative:
                    continue
                try:
                    child, _ = await _download_remote_install(urljoin(url, alternate), 1024 * 1024)
                    stored_relative = alternate
                except (ValueError, httpx.HTTPError):
                    continue
            total += len(child)
            if total > 2 * 1024 * 1024:
                raise ValueError("Skill 包及引用文件合计超过 2MB")
            files[stored_relative] = child
        return skill_registry.install_package(files, fallback_id=Path(filename).stem)
    return _install_skill_bytes(raw, filename)


async def _install_mcp_remote_url(url: str) -> list[dict[str, Any]]:
    raw, _ = await _download_remote_install(url, 1024 * 1024)
    return mcp_gateway.import_config(json.loads(raw.decode("utf-8")))


runtime = AgentRuntime(
    skill_registry,
    mcp_gateway,
    model_gateway,
    skill_url_installer=_install_skill_remote_url,
    mcp_url_installer=_install_mcp_remote_url,
    task_state=task_state,
    policy_engine=policy_engine,
    context_service=context_service,
)
loop_scheduler = LoopScheduler(runtime)
expert_team_service = ExpertTeamService(runtime, task_state=task_state)
_runtime_tasks: set[asyncio.Task[Any]] = set()


def _schedule_runtime(task_id: str, run_id: str | None = None) -> asyncio.Task[Any]:
    background = asyncio.create_task(runtime.run_task(task_id, run_id=run_id))
    _runtime_tasks.add(background)
    background.add_done_callback(_runtime_tasks.discard)
    return background


def _schedule_team_run(team_run_id: str) -> asyncio.Task[Any]:
    background = asyncio.create_task(expert_team_service.run_team(team_run_id))
    _runtime_tasks.add(background)
    background.add_done_callback(_runtime_tasks.discard)
    return background


def _schedule_member_retry(
    team_run_id: str, member_run_id: str, scope: ExecutionScope
) -> asyncio.Task[Any]:
    background = asyncio.create_task(
        expert_team_service.retry_member(team_run_id, member_run_id, scope)
    )
    _runtime_tasks.add(background)
    background.add_done_callback(_runtime_tasks.discard)
    return background


def _reload_policy_rules() -> None:
    rules: list[dict[str, Any]] = []
    for row in db.query_all("SELECT rule_json FROM policy_rules WHERE enabled = 1 ORDER BY priority DESC, id"):
        value = db.json_loads(row.get("rule_json"), {})
        if isinstance(value, dict):
            rules.append(value)
    env_rules = os.getenv("APP_POLICY_RULES_JSON", "").strip()
    if env_rules:
        parsed = json.loads(env_rules)
        if not isinstance(parsed, list):
            raise ValueError("APP_POLICY_RULES_JSON 必须是规则数组")
        rules.extend(item for item in parsed if isinstance(item, dict))
    policy_engine.set_rules(rules)


def _fail_running_nodes(run_id: str, reason: str) -> None:
    for node in task_state.list_nodes(run_id):
        if node["status"] != "running":
            continue
        try:
            task_state.fail_node(
                node["id"],
                {"message": reason, "error_type": "ServiceRestart"},
                metadata={"interrupted": True},
            )
        except TaskStateError:
            continue


def _recover_interrupted_runs(
    exclude_task_ids: set[str] | None = None,
    preserve_task_ids: set[str] | None = None,
) -> list[dict[str, Any]]:
    """Recover ordinary Agent runs without stealing orchestrated work.

    ``exclude_task_ids`` are interrupted automation attempts and are closed as
    failed. ``preserve_task_ids`` belong to another durable orchestrator (for
    example a queued expert-team run); they are ignored here and left intact
    for that orchestrator to schedule.
    """
    excluded = exclude_task_ids or set()
    preserved = preserve_task_ids or set()
    protected = excluded | preserved
    for task_id in excluded:
        task = db.query_one("SELECT status FROM tasks WHERE id = ?", (task_id,))
        if task and task.get("status") not in {"completed", "failed", "cancelled"}:
            db.update_task_status(
                task_id, "failed",
                result={"error": "自动化尝试因平台服务重启而中断", "error_type": "ServiceRestart"},
            )
    recovered: list[dict[str, Any]] = []
    interrupted_runs = [
        *task_state.list_runs(status="running", limit=1000),
        *task_state.list_runs(status="paused", limit=1000),
    ]
    for waiting_run in task_state.list_runs(status="waiting_approval", limit=1000):
        waiting_task = db.query_one(
            "SELECT result_json FROM tasks WHERE id = ?", (waiting_run["task_id"],)
        ) or {}
        waiting_result = db.json_loads(waiting_task.get("result_json"), {})
        if waiting_result.get("pending_action") == "policy_approval":
            interrupted_runs.append(waiting_run)
    for old_run in interrupted_runs:
        task = db.query_one("SELECT id FROM tasks WHERE id = ?", (old_run["task_id"],))
        if not task:
            continue
        if old_run["task_id"] in preserved:
            continue
        _fail_running_nodes(old_run["id"], "平台服务重启，旧执行尝试已中断")
        if old_run["task_id"] in excluded:
            task_state.finish_run(
                old_run["id"],
                status="failed",
                error={"message": "自动化尝试因平台服务重启而中断", "error_type": "ServiceRestart"},
                metadata={"interrupted": True, "automation_run": True},
            )
            db.update_task_status(
                old_run["task_id"], "failed",
                result={"error": "自动化尝试因平台服务重启而中断", "error_type": "ServiceRestart"},
            )
            continue
        task_state.finish_run(
            old_run["id"],
            status="failed",
            error={"message": "平台服务重启，已创建恢复尝试", "error_type": "ServiceRestart"},
            metadata={"interrupted": True},
        )
        checkpoint = task_state.latest_checkpoint(old_run["id"], include_state=False)
        new_run = task_state.create_run(
            old_run["task_id"],
            resumed_from_checkpoint_id=checkpoint["id"] if checkpoint else None,
            metadata={"recovered_after_restart": True, "previous_run_id": old_run["id"]},
        )
        db.update_task_status(old_run["task_id"], "queued")
        emit_data = {
            "previous_run_id": old_run["id"],
            "run_id": new_run["id"],
            "checkpoint_id": checkpoint["id"] if checkpoint else "",
        }
        db.insert_event(
            old_run["task_id"],
            "recovery_scheduled",
            "已安排服务重启恢复",
            "将从最近安全检查点创建新的运行尝试。" if checkpoint else "未找到检查点，将从任务起点重新执行。",
            emit_data,
        )
        recovered.append(new_run)

    # Queued runs survive a restart unchanged. Legacy queued tasks without a
    # task_run receive one before they are scheduled.
    all_queued_runs = task_state.list_runs(status="queued", limit=1000)
    for queued in all_queued_runs:
        if queued["task_id"] in excluded:
            task_state.finish_run(
                queued["id"], status="cancelled",
                error={"message": "自动化尝试因平台服务重启而中断", "error_type": "ServiceRestart"},
                metadata={"interrupted": True, "automation_run": True},
            )
    queued_runs = [item for item in all_queued_runs if item["task_id"] not in protected]
    known_queued_tasks = {item["task_id"] for item in queued_runs}
    tasks_with_runs = {
        item["task_id"] for item in task_state.list_runs(limit=10_000)
    }
    for task in db.query_all("SELECT id FROM tasks WHERE status = 'running'"):
        if task["id"] in preserved:
            continue
        if task["id"] in excluded:
            db.update_task_status(
                task["id"], "failed",
                result={"error": "自动化尝试因平台服务重启而中断", "error_type": "ServiceRestart"},
            )
            continue
        if task["id"] in tasks_with_runs:
            continue
        db.update_task_status(task["id"], "queued")
        legacy_run = task_state.create_run(
            task["id"], metadata={"legacy_interrupted_task": True, "recovered_after_restart": True}
        )
        queued_runs.append(legacy_run)
        known_queued_tasks.add(task["id"])
        db.insert_event(
            task["id"],
            "recovery_scheduled",
            "已安排旧任务恢复",
            "检测到升级前遗留的运行中任务，将从任务起点重新执行。",
            {"run_id": legacy_run["id"]},
        )
    for task in db.query_all("SELECT id FROM tasks WHERE status = 'queued'"):
        if task["id"] not in protected and task["id"] not in known_queued_tasks:
            queued_runs.append(task_state.create_run(task["id"], metadata={"legacy_task": True}))
    by_id = {item["id"]: item for item in [*queued_runs, *recovered]}
    return list(by_id.values())


@app.on_event("startup")
async def on_startup() -> None:
    db.init_db()
    context_service.init_schema()
    task_state.init_schema()
    interrupted_loop_tasks = loop_scheduler.recover_interrupted_runs()
    queued_team_runs = expert_team_service.queued_runs_for_recovery()
    queued_team_task_ids = {item["parent_task_id"] for item in queued_team_runs}
    _migrate_legacy_model_keys()
    skill_registry.load_builtin_skills()
    mcp_gateway.seed_builtin_servers()
    seed_agents()
    _reload_policy_rules()
    loop_scheduler.start()
    for run in _recover_interrupted_runs(interrupted_loop_tasks, queued_team_task_ids):
        _schedule_runtime(run["task_id"], run["id"])
    for team_run in queued_team_runs:
        _schedule_team_run(team_run["id"])


@app.on_event("shutdown")
async def on_shutdown() -> None:
    await loop_scheduler.stop()
    pending = [item for item in _runtime_tasks if not item.done()]
    for item in pending:
        item.cancel()
    if pending:
        await asyncio.gather(*pending, return_exceptions=True)


@app.get("/api/health")
def health() -> dict[str, Any]:
    return {"ok": True, "name": "AgentNexus", "product": "AgentNexus", "display_name": "智枢"}


@app.get("/api/capabilities")
def capabilities() -> dict[str, Any]:
    presentation = presentation_generation_status()
    outbound_enabled = outbound_network_enabled()
    return {
        "outbound_network": {"supported": True, "enabled": outbound_enabled},
        "file_upload": {
            "supported": True,
            "max_mb": int(os.getenv("APP_MAX_UPLOAD_MB", "20")),
            "text_extraction": [
                "txt", "md", "csv", "json", "yaml", "code",
                "docx", "xlsx", "pptx", "pdf",
            ],
            "max_files_per_task": runtime.ATTACHMENT_MAX_FILES,
            "max_chars_per_file": runtime.ATTACHMENT_MAX_FILE_CHARS,
            "max_context_chars": runtime.ATTACHMENT_MAX_CONTEXT_CHARS,
        },
        "web_search": {
            "supported": True,
            "enabled": outbound_enabled and env_flag("APP_ALLOW_WEB_SEARCH"),
            "configured": bool(os.getenv("TAVILY_API_KEY") or os.getenv("BRAVE_SEARCH_API_KEY")),
            "provider": "tavily/brave",
        },
        "stdio_mcp": {"supported": True, "enabled": os.getenv("APP_ALLOW_STDIO_MCP", "false").lower() in {"1", "true", "yes"}},
        "remote_mcp": {"supported": True, "enabled": outbound_enabled and env_flag("APP_ALLOW_REMOTE_MCP")},
        "http_tools": {"supported": True, "enabled": outbound_enabled and env_flag("APP_ALLOW_HTTP_TOOLS")},
        "remote_install": {"supported": True, "enabled": outbound_enabled and _remote_install_flag()},
        "direct_api_key": {"supported": True, "encrypted": True, "storage": "local"},
        "models": ["deterministic", "openai", "openai_compatible"],
        "document_output": {
            "formats": ["markdown", "docx", "pdf", "xlsx", "csv", "html"],
            "optional_formats": ["pptx"],
            "pptx_configured": bool(presentation.get("configured")),
            "pptx_reason": str(presentation.get("reason") or ""),
        },
        "memory": {"supported": True, "scopes": ["organization", "workspace", "user", "agent", "conversation"], "revision_history": True, "conversation_summary": {"automatic": True, "viewable": True, "editable": True, "deletable": True}},
        "expert_teams": {
            "supported": True,
            "template_installation": True,
            "parallel_members": True,
            "isolated_member_context": True,
            "supervisor_aggregation": True,
            "single_member_retry": True,
        },
        "automation": {
            "supported": True,
            "triggers": ["manual", "interval", "cron", "once", "webhook"],
            "persistent_history": True,
            "signed_webhooks": True,
            "idempotency": True,
            "notifications": True,
            "structured_state_diff": True,
            "legacy_api": "/api/loops",
        },
        "policy_hooks": {
            "supported": True,
            "handlers": ["builtin_rule", "http"],
            "http_enabled": policy_engine.http_enabled,
            "arbitrary_shell": False,
        },
    }


def _api_scope(
    organization_id: str = "local-org",
    workspace_id: str = "default",
    user_id: str = "local-user",
    agent_id: str = "",
    conversation_id: str = "",
) -> ExecutionScope:
    try:
        return ExecutionScope(
            organization_id=organization_id,
            workspace_id=workspace_id,
            user_id=user_id,
            agent_id=agent_id,
            conversation_id=conversation_id,
        )
    except (TypeError, ValueError) as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc


def _expert_http_error(exc: Exception) -> HTTPException:
    if isinstance(exc, ExpertNotFoundError):
        return HTTPException(status_code=404, detail=str(exc))
    if isinstance(exc, ExpertPermissionError):
        return HTTPException(status_code=403, detail=str(exc))
    if isinstance(exc, ExpertConflictError):
        return HTTPException(status_code=409, detail=str(exc))
    return HTTPException(status_code=400, detail=str(exc))


def _public_expert_selection(
    team: dict[str, Any],
    scope: ExecutionScope,
    *,
    automatic: bool,
    recommendation: dict[str, Any] | None = None,
) -> dict[str, Any]:
    """Project an expert-team routing decision onto safe, user-facing fields."""

    supervisor_id = str(team.get("supervisor_agent_id") or "")
    supervisor = expert_team_service.get_agent(supervisor_id, scope) or {}
    members: list[dict[str, str]] = []
    for member in team.get("members") or []:
        agent_id = str(member.get("agent_id") or "")
        agent = expert_team_service.get_agent(agent_id, scope) or {}
        members.append(
            {
                "agent_id": agent_id,
                "agent_name": str(agent.get("name") or agent_id),
                "role": str(member.get("role") or "专家"),
            }
        )
    recommendation = recommendation or {}
    return {
        "selection_mode": "automatic" if automatic else "manual",
        "team_id": str(team.get("id") or ""),
        "team_name": str(team.get("name") or team.get("id") or "专家团"),
        "reason": str(
            recommendation.get("reason")
            or ("已根据当前目标自动匹配专家团" if automatic else "使用用户指定的专家团")
        ),
        "matched_terms": [str(item) for item in recommendation.get("matched_terms") or []][:8],
        "supervisor": {
            "agent_id": supervisor_id,
            "agent_name": str(supervisor.get("name") or supervisor_id),
        },
        "members": members,
    }


@app.get("/api/expert-templates")
def list_expert_templates(
    organization_id: str = "local-org",
    workspace_id: str = "default",
    user_id: str = "local-user",
    include_disabled: bool = False,
) -> list[dict[str, Any]]:
    return expert_team_service.list_templates(
        _api_scope(organization_id, workspace_id, user_id),
        include_disabled=include_disabled,
    )


@app.post("/api/expert-templates", status_code=201)
def create_expert_template(payload: ExpertTemplateCreate) -> dict[str, Any]:
    try:
        return expert_team_service.create_template(payload.model_dump())
    except (ExpertConflictError, ExpertValidationError) as exc:
        raise _expert_http_error(exc) from exc


@app.get("/api/expert-templates/{template_id}")
def get_expert_template(
    template_id: str,
    organization_id: str = "local-org",
    workspace_id: str = "default",
    user_id: str = "local-user",
) -> dict[str, Any]:
    item = expert_team_service.get_template(
        template_id, _api_scope(organization_id, workspace_id, user_id)
    )
    if not item:
        raise HTTPException(status_code=404, detail="专家模板不存在或当前作用域不可见")
    return item


@app.put("/api/expert-templates/{template_id}")
def update_expert_template(
    template_id: str,
    payload: ExpertTemplateUpdate,
    organization_id: str = "local-org",
    workspace_id: str = "default",
    user_id: str = "local-user",
) -> dict[str, Any]:
    try:
        return expert_team_service.update_template(
            template_id,
            _api_scope(organization_id, workspace_id, user_id),
            payload.model_dump(exclude_unset=True),
        )
    except (ExpertNotFoundError, ExpertPermissionError, ExpertConflictError, ExpertValidationError) as exc:
        raise _expert_http_error(exc) from exc


@app.delete("/api/expert-templates/{template_id}")
def delete_expert_template(
    template_id: str,
    organization_id: str = "local-org",
    workspace_id: str = "default",
    user_id: str = "local-user",
) -> dict[str, Any]:
    try:
        expert_team_service.delete_template(
            template_id, _api_scope(organization_id, workspace_id, user_id)
        )
        return {"ok": True, "id": template_id}
    except (ExpertNotFoundError, ExpertPermissionError, ExpertConflictError) as exc:
        raise _expert_http_error(exc) from exc


@app.post("/api/expert-templates/{template_id}/install", status_code=201)
def install_expert_template(
    template_id: str, payload: ExpertInstallRequest
) -> dict[str, Any]:
    scope = _api_scope(payload.organization_id, payload.workspace_id, payload.user_id)
    try:
        return expert_team_service.install_template(template_id, scope, payload.model_dump())
    except (ExpertNotFoundError, ExpertPermissionError, ExpertConflictError, ExpertValidationError) as exc:
        raise _expert_http_error(exc) from exc


@app.get("/api/expert-installations")
def list_expert_installations(
    organization_id: str = "local-org",
    workspace_id: str = "default",
    user_id: str = "local-user",
    include_disabled: bool = False,
) -> list[dict[str, Any]]:
    return expert_team_service.list_installations(
        _api_scope(organization_id, workspace_id, user_id),
        include_disabled=include_disabled,
    )


@app.delete("/api/expert-installations/{installation_id}")
def uninstall_expert(
    installation_id: str,
    organization_id: str = "local-org",
    workspace_id: str = "default",
    user_id: str = "local-user",
) -> dict[str, Any]:
    try:
        return expert_team_service.disable_installation(
            installation_id, _api_scope(organization_id, workspace_id, user_id)
        )
    except (ExpertNotFoundError, ExpertConflictError) as exc:
        raise _expert_http_error(exc) from exc


@app.get("/api/expert-teams")
def list_expert_teams(
    organization_id: str = "local-org",
    workspace_id: str = "default",
    user_id: str = "local-user",
    include_disabled: bool = False,
) -> list[dict[str, Any]]:
    return expert_team_service.list_teams(
        _api_scope(organization_id, workspace_id, user_id),
        include_disabled=include_disabled,
    )


@app.post("/api/expert-teams", status_code=201)
def create_expert_team(payload: ExpertTeamCreate) -> dict[str, Any]:
    try:
        return expert_team_service.create_team(payload.model_dump())
    except (ExpertConflictError, ExpertValidationError) as exc:
        raise _expert_http_error(exc) from exc


@app.get("/api/expert-teams/{team_id}")
def get_expert_team(
    team_id: str,
    organization_id: str = "local-org",
    workspace_id: str = "default",
    user_id: str = "local-user",
) -> dict[str, Any]:
    item = expert_team_service.get_team(
        team_id, _api_scope(organization_id, workspace_id, user_id)
    )
    if not item:
        raise HTTPException(status_code=404, detail="专家团不存在或当前作用域不可见")
    return item


@app.put("/api/expert-teams/{team_id}")
def update_expert_team(
    team_id: str,
    payload: ExpertTeamUpdate,
    organization_id: str = "local-org",
    workspace_id: str = "default",
    user_id: str = "local-user",
) -> dict[str, Any]:
    try:
        return expert_team_service.update_team(
            team_id,
            _api_scope(organization_id, workspace_id, user_id),
            payload.model_dump(exclude_unset=True),
        )
    except (ExpertNotFoundError, ExpertPermissionError, ExpertConflictError, ExpertValidationError) as exc:
        raise _expert_http_error(exc) from exc


@app.delete("/api/expert-teams/{team_id}")
def delete_expert_team(
    team_id: str,
    organization_id: str = "local-org",
    workspace_id: str = "default",
    user_id: str = "local-user",
) -> dict[str, Any]:
    try:
        expert_team_service.delete_team(
            team_id, _api_scope(organization_id, workspace_id, user_id)
        )
        return {"ok": True, "id": team_id}
    except (ExpertNotFoundError, ExpertPermissionError, ExpertConflictError) as exc:
        raise _expert_http_error(exc) from exc


@app.post("/api/expert-teams/{team_id}/runs", status_code=202)
async def run_expert_team(team_id: str, payload: ExpertTeamRunCreate) -> dict[str, Any]:
    if payload.model_id and payload.model_id != "deterministic":
        configured = db.query_one(
            "SELECT id FROM model_configs WHERE id = ? AND enabled = 1", (payload.model_id,)
        )
        if not configured:
            raise HTTPException(status_code=400, detail="所选模型不存在或未启用")
    scope = _api_scope(payload.organization_id, payload.workspace_id, payload.user_id)
    try:
        task, parent_run, team_run = expert_team_service.create_task_and_run(
            team_id, scope, message=payload.message, model_id=payload.model_id,
            conversation_id=payload.conversation_id,
        )
    except (ExpertNotFoundError, ExpertConflictError, ExpertValidationError) as exc:
        raise _expert_http_error(exc) from exc
    _schedule_team_run(team_run["id"])
    return {"accepted": True, "task": task, "run": parent_run, "team_run": team_run}


@app.get("/api/expert-teams/{team_id}/runs")
def list_expert_team_runs(
    team_id: str,
    organization_id: str = "local-org",
    workspace_id: str = "default",
    user_id: str = "local-user",
) -> list[dict[str, Any]]:
    try:
        return expert_team_service.list_team_runs(
            team_id, _api_scope(organization_id, workspace_id, user_id)
        )
    except ExpertNotFoundError as exc:
        raise _expert_http_error(exc) from exc


@app.get("/api/expert-team-runs/{team_run_id}")
def get_expert_team_run(
    team_run_id: str,
    organization_id: str = "local-org",
    workspace_id: str = "default",
    user_id: str = "local-user",
) -> dict[str, Any]:
    item = expert_team_service.get_team_run(
        team_run_id, _api_scope(organization_id, workspace_id, user_id)
    )
    if not item:
        raise HTTPException(status_code=404, detail="专家团运行不存在或当前作用域不可见")
    return item


@app.post(
    "/api/expert-team-runs/{team_run_id}/members/{member_run_id}/retry",
    status_code=202,
)
async def retry_expert_team_member(
    team_run_id: str,
    member_run_id: str,
    payload: ExpertMemberRetryRequest,
    organization_id: str = "local-org",
    workspace_id: str = "default",
    user_id: str = "local-user",
) -> dict[str, Any]:
    _ = payload
    scope = _api_scope(organization_id, workspace_id, user_id)
    try:
        expert_team_service.validate_member_retry(team_run_id, member_run_id, scope)
    except (ExpertNotFoundError, ExpertConflictError) as exc:
        raise _expert_http_error(exc) from exc
    _schedule_member_retry(team_run_id, member_run_id, scope)
    return {
        "accepted": True,
        "team_run_id": team_run_id,
        "member_run_id": member_run_id,
        "scope": {
            "organization_id": organization_id,
            "workspace_id": workspace_id,
            "user_id": user_id,
        },
    }


@app.get("/api/memories")
def list_memories(
    organization_id: str = "local-org",
    workspace_id: str = "default",
    user_id: str = "local-user",
    agent_id: str = "",
    conversation_id: str = "",
    scope_type: str | None = None,
    include_disabled: bool = True,
    include_expired: bool = True,
) -> list[dict[str, Any]]:
    try:
        return context_service.list_memories(
            _api_scope(organization_id, workspace_id, user_id, agent_id, conversation_id),
            scope_type=scope_type,
            include_disabled=include_disabled,
            include_expired=include_expired,
        )
    except (TypeError, ValueError) as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc


@app.post("/api/memories", status_code=201)
def create_memory(payload: MemoryCreate) -> dict[str, Any]:
    scope = _api_scope(
        payload.organization_id,
        payload.workspace_id,
        payload.user_id,
        payload.agent_id,
        payload.conversation_id,
    )
    try:
        return context_service.create_memory(
            scope,
            scope_type=payload.scope_type,
            kind=payload.kind,
            title=payload.title,
            content=payload.content,
            tags=payload.tags,
            source_type=payload.source_type,
            source_ref=payload.source_ref,
            trust_level=payload.trust_level,
            enabled=payload.enabled,
            expires_at=payload.expires_at,
            created_by=payload.user_id,
        )
    except (TypeError, ValueError) as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc


@app.get("/api/memories/{memory_id}")
def get_memory(
    memory_id: str,
    organization_id: str = "local-org",
    workspace_id: str = "default",
    user_id: str = "local-user",
    agent_id: str = "",
    conversation_id: str = "",
) -> dict[str, Any]:
    memory = context_service.get_memory(
        memory_id,
        _api_scope(organization_id, workspace_id, user_id, agent_id, conversation_id),
    )
    if not memory:
        raise HTTPException(status_code=404, detail="记忆不存在或不属于当前执行作用域")
    return memory


@app.put("/api/memories/{memory_id}")
def update_memory(
    memory_id: str,
    payload: MemoryUpdate,
    organization_id: str = "local-org",
    workspace_id: str = "default",
    user_id: str = "local-user",
    agent_id: str = "",
    conversation_id: str = "",
) -> dict[str, Any]:
    changes = payload.model_dump(exclude_unset=True)
    reason = str(changes.pop("reason", "updated"))
    try:
        return context_service.update_memory(
            memory_id,
            _api_scope(organization_id, workspace_id, user_id, agent_id, conversation_id),
            actor_id=user_id,
            reason=reason,
            **changes,
        )
    except MemoryNotFoundError as exc:
        raise HTTPException(status_code=404, detail=str(exc)) from exc
    except (TypeError, ValueError) as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc


@app.post("/api/memories/{memory_id}/enable")
def enable_memory(
    memory_id: str,
    organization_id: str = "local-org",
    workspace_id: str = "default",
    user_id: str = "local-user",
    agent_id: str = "",
    conversation_id: str = "",
) -> dict[str, Any]:
    try:
        return context_service.enable_memory(
            memory_id,
            _api_scope(organization_id, workspace_id, user_id, agent_id, conversation_id),
            actor_id=user_id,
        )
    except MemoryNotFoundError as exc:
        raise HTTPException(status_code=404, detail=str(exc)) from exc


@app.post("/api/memories/{memory_id}/disable")
def disable_memory(
    memory_id: str,
    organization_id: str = "local-org",
    workspace_id: str = "default",
    user_id: str = "local-user",
    agent_id: str = "",
    conversation_id: str = "",
) -> dict[str, Any]:
    try:
        return context_service.disable_memory(
            memory_id,
            _api_scope(organization_id, workspace_id, user_id, agent_id, conversation_id),
            actor_id=user_id,
        )
    except MemoryNotFoundError as exc:
        raise HTTPException(status_code=404, detail=str(exc)) from exc


@app.delete("/api/memories/{memory_id}")
def delete_memory(
    memory_id: str,
    organization_id: str = "local-org",
    workspace_id: str = "default",
    user_id: str = "local-user",
    agent_id: str = "",
    conversation_id: str = "",
) -> dict[str, Any]:
    try:
        return context_service.delete_memory(
            memory_id,
            _api_scope(organization_id, workspace_id, user_id, agent_id, conversation_id),
            actor_id=user_id,
            reason="api_delete",
        )
    except MemoryNotFoundError as exc:
        raise HTTPException(status_code=404, detail=str(exc)) from exc


@app.get("/api/memories/{memory_id}/revisions")
def list_memory_revisions(
    memory_id: str,
    organization_id: str = "local-org",
    workspace_id: str = "default",
    user_id: str = "local-user",
    agent_id: str = "",
    conversation_id: str = "",
) -> list[dict[str, Any]]:
    return context_service.list_revisions(
        memory_id,
        _api_scope(organization_id, workspace_id, user_id, agent_id, conversation_id),
    )


@app.get("/api/context/effective")
def get_effective_context(
    organization_id: str = "local-org",
    workspace_id: str = "default",
    user_id: str = "local-user",
    agent_id: str = "",
    conversation_id: str = "",
) -> dict[str, Any]:
    return context_service.get_effective_context(
        _api_scope(organization_id, workspace_id, user_id, agent_id, conversation_id)
    )


@app.get("/api/conversation-summaries")
def list_conversation_summaries(
    organization_id: str = "local-org",
    workspace_id: str = "default",
    user_id: str = "local-user",
    limit: int = 100,
) -> list[dict[str, Any]]:
    try:
        return runtime.conversation_summary_service.list(
            _api_scope(organization_id, workspace_id, user_id), limit=limit
        )
    except (TypeError, ValueError) as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc


@app.get("/api/conversation-summaries/{conversation_id}")
def get_conversation_summary(
    conversation_id: str,
    organization_id: str = "local-org",
    workspace_id: str = "default",
    user_id: str = "local-user",
) -> dict[str, Any]:
    try:
        item = runtime.conversation_summary_service.get(
            _api_scope(
                organization_id,
                workspace_id,
                user_id,
                conversation_id=conversation_id,
            ),
            conversation_id,
        )
    except (TypeError, ValueError) as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc
    if not item:
        raise HTTPException(status_code=404, detail="对话摘要不存在或不属于当前执行作用域")
    return item


@app.put("/api/conversation-summaries/{conversation_id}")
def update_conversation_summary(
    conversation_id: str,
    payload: ConversationSummaryUpdate,
    organization_id: str = "local-org",
    workspace_id: str = "default",
    user_id: str = "local-user",
) -> dict[str, Any]:
    scope = _api_scope(
        organization_id,
        workspace_id,
        user_id,
        conversation_id=conversation_id,
    )
    try:
        return runtime.conversation_summary_service.upsert(
            scope,
            conversation_id=conversation_id,
            summary=payload.summary,
            preserved_constraints=payload.preserved_constraints,
            through_task_id=payload.through_task_id,
            model_id=payload.model_id,
        )
    except ConversationSummaryConflictError as exc:
        raise HTTPException(status_code=409, detail="该对话 ID 已属于其他用户或工作区") from exc
    except (TypeError, ValueError) as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc


@app.delete("/api/conversation-summaries/{conversation_id}")
def delete_conversation_summary(
    conversation_id: str,
    organization_id: str = "local-org",
    workspace_id: str = "default",
    user_id: str = "local-user",
) -> dict[str, Any]:
    try:
        deleted = runtime.conversation_summary_service.delete(
            _api_scope(
                organization_id,
                workspace_id,
                user_id,
                conversation_id=conversation_id,
            ),
            conversation_id,
        )
    except (TypeError, ValueError) as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc
    if not deleted:
        raise HTTPException(status_code=404, detail="对话摘要不存在或不属于当前执行作用域")
    return {"conversation_id": conversation_id, "deleted": True}


def _policy_rule_to_api(value: dict[str, Any]) -> dict[str, Any]:
    return PolicyRule.from_dict(value).to_dict()


@app.get("/api/policies")
def list_policy_rules() -> list[dict[str, Any]]:
    result: list[dict[str, Any]] = []
    for row in db.query_all("SELECT rule_json FROM policy_rules ORDER BY priority DESC, id"):
        value = db.json_loads(row.get("rule_json"), {})
        if isinstance(value, dict):
            result.append(_policy_rule_to_api(value))
    return result


@app.get("/api/policies/{rule_id}")
def get_policy_rule(rule_id: str) -> dict[str, Any]:
    row = db.query_one("SELECT rule_json FROM policy_rules WHERE id = ?", (rule_id,))
    if not row:
        raise HTTPException(status_code=404, detail="策略规则不存在")
    return _policy_rule_to_api(db.json_loads(row.get("rule_json"), {}))


@app.post("/api/policies", status_code=201)
def create_policy_rule(payload: PolicyRuleCreate) -> dict[str, Any]:
    if db.query_one("SELECT id FROM policy_rules WHERE id = ?", (payload.id,)):
        raise HTTPException(status_code=409, detail="策略规则 ID 已存在")
    raw = payload.model_dump(exclude_none=True)
    raw["name"] = str(raw.get("name") or payload.id)
    try:
        rule = PolicyRule.from_dict(raw)
    except PolicyConfigurationError as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc
    now = db.utc_now()
    event_label = ",".join(rule.events)
    db.execute(
        """
        INSERT INTO policy_rules(id, name, event, scope, scope_id, priority, enabled, rule_json, created_at, updated_at)
        VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
        """,
        (
            rule.id,
            rule.name,
            event_label,
            rule.scope,
            rule.scope_id or "",
            rule.priority,
            1 if rule.enabled else 0,
            db.json_dumps(raw),
            now,
            now,
        ),
    )
    _reload_policy_rules()
    return rule.to_dict()


@app.put("/api/policies/{rule_id}")
def update_policy_rule(rule_id: str, payload: PolicyRuleUpdate) -> dict[str, Any]:
    row = db.query_one("SELECT rule_json FROM policy_rules WHERE id = ?", (rule_id,))
    if not row:
        raise HTTPException(status_code=404, detail="策略规则不存在")
    current = db.json_loads(row.get("rule_json"), {})
    incoming = payload.model_dump(exclude_unset=True)
    if "event" in incoming:
        current.pop("events", None)
    if "events" in incoming:
        current.pop("event", None)
    merged = {**current, **incoming, "id": rule_id}
    try:
        rule = PolicyRule.from_dict(merged)
    except PolicyConfigurationError as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc
    db.execute(
        """
        UPDATE policy_rules
        SET name = ?, event = ?, scope = ?, scope_id = ?, priority = ?, enabled = ?, rule_json = ?, updated_at = ?
        WHERE id = ?
        """,
        (
            rule.name,
            ",".join(rule.events),
            rule.scope,
            rule.scope_id or "",
            rule.priority,
            1 if rule.enabled else 0,
            db.json_dumps(merged),
            db.utc_now(),
            rule_id,
        ),
    )
    _reload_policy_rules()
    return rule.to_dict()


@app.delete("/api/policies/{rule_id}")
def delete_policy_rule(rule_id: str) -> dict[str, Any]:
    if not db.query_one("SELECT id FROM policy_rules WHERE id = ?", (rule_id,)):
        raise HTTPException(status_code=404, detail="策略规则不存在")
    db.execute("DELETE FROM policy_rules WHERE id = ?", (rule_id,))
    _reload_policy_rules()
    return {"ok": True, "id": rule_id}


@app.get("/api/skills")
def list_skills() -> list[dict[str, Any]]:
    return skill_registry.list_skills()


@app.get("/api/skills/{skill_id}")
def get_skill(skill_id: str) -> dict[str, Any]:
    skill = skill_registry.get_skill(skill_id)
    if not skill:
        raise HTTPException(status_code=404, detail="Skill not found")
    return skill


@app.get("/api/skills/{skill_id}/export")
def export_skill_package(skill_id: str) -> StreamingResponse:
    skill = skill_registry.get_skill(skill_id)
    if not skill:
        raise HTTPException(status_code=404, detail="Skill not found")
    buffer = BytesIO()
    with zipfile.ZipFile(buffer, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        for item in skill_registry.list_files(skill_id):
            row = db.query_one(
                "SELECT content FROM skill_files WHERE skill_id = ? AND path = ?",
                (skill_id, item["path"]),
            )
            raw = (row or {}).get("content") or b""
            if isinstance(raw, str):
                raw = raw.encode("utf-8")
            archive.writestr(item["path"], raw)
    buffer.seek(0)
    safe_name = re.sub(r"[^A-Za-z0-9._-]", "_", skill_id) or "skill"
    return StreamingResponse(
        buffer,
        media_type="application/zip",
        headers={"Content-Disposition": f'attachment; filename="{safe_name}.zip"'},
    )


@app.get("/api/skills/{skill_id}/files")
def list_skill_files(skill_id: str) -> list[dict[str, Any]]:
    if not skill_registry.get_skill(skill_id):
        raise HTTPException(status_code=404, detail="Skill not found")
    return skill_registry.list_files(skill_id)


@app.post("/api/skills/{skill_id}/files/upload")
async def upload_skill_file(skill_id: str, path: str, file: UploadFile = File(...)) -> dict[str, Any]:
    raw = await file.read(1024 * 1024 + 1)
    if len(raw) > 1024 * 1024:
        raise HTTPException(status_code=413, detail="单个 Skill 文件不能超过 1MB")
    try:
        return skill_registry.put_file(skill_id, path, raw)
    except ValueError as exc:
        status = 404 if str(exc) == "Skill not found" else 400
        raise HTTPException(status_code=status, detail=str(exc)) from exc


@app.get("/api/skills/{skill_id}/files/{file_path:path}")
def get_skill_file(skill_id: str, file_path: str) -> dict[str, Any]:
    try:
        item = skill_registry.get_file(skill_id, file_path)
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc
    if not item:
        raise HTTPException(status_code=404, detail="Skill file not found")
    return item


@app.put("/api/skills/{skill_id}/files/{file_path:path}")
def put_skill_file(skill_id: str, file_path: str, payload: SkillFileUpdate) -> dict[str, Any]:
    try:
        return skill_registry.put_file(skill_id, file_path, payload.content.encode("utf-8"))
    except ValueError as exc:
        status = 404 if str(exc) == "Skill not found" else 400
        raise HTTPException(status_code=status, detail=str(exc)) from exc


@app.delete("/api/skills/{skill_id}/files/{file_path:path}")
def delete_skill_file(skill_id: str, file_path: str) -> dict[str, Any]:
    try:
        skill_registry.delete_file(skill_id, file_path)
        return {"ok": True, "path": file_path}
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc


@app.post("/api/skills")
def create_skill(payload: SkillCreate) -> dict[str, Any]:
    if skill_registry.get_skill(payload.id):
        raise HTTPException(status_code=409, detail="Skill already exists")
    missing_mcps = [item for item in payload.required_mcps if not mcp_gateway.server_exists(item)]
    if missing_mcps:
        raise HTTPException(status_code=400, detail="Skill 引用了不存在的 MCP：" + "、".join(missing_mcps))
    return skill_registry.create_skill(payload.model_dump())


@app.post("/api/skills/install/upload")
async def install_skill_upload(file: UploadFile = File(...)) -> dict[str, Any]:
    raw = await file.read(2 * 1024 * 1024 + 1)
    if len(raw) > 2 * 1024 * 1024:
        raise HTTPException(status_code=413, detail="Skill 安装包不能超过 2MB")
    try:
        filename = file.filename or "SKILL.md"
        if filename != "SKILL.md" and not filename.lower().endswith((".md", ".zip")):
            raise ValueError("仅支持 SKILL.md 或 ZIP 安装包")
        return _install_skill_bytes(raw, filename)
    except (ValueError, UnicodeDecodeError, zipfile.BadZipFile) as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc


@app.post("/api/skills/install/url")
async def install_skill_url(payload: RemoteInstall) -> dict[str, Any]:
    try:
        return await _install_skill_remote_url(payload.url)
    except (ValueError, UnicodeDecodeError, zipfile.BadZipFile, httpx.HTTPError) as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc


@app.post("/api/skills/install/path")
def install_skill_path(payload: SkillPathInstall) -> dict[str, Any]:
    try:
        return skill_registry.install_from_path(payload.path, payload.enabled)
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc


@app.put("/api/skills/{skill_id}")
def update_skill(skill_id: str, payload: SkillUpdate) -> dict[str, Any]:
    if payload.required_mcps is not None:
        missing_mcps = [item for item in payload.required_mcps if not mcp_gateway.server_exists(item)]
        if missing_mcps:
            raise HTTPException(status_code=400, detail="Skill 引用了不存在的 MCP：" + "、".join(missing_mcps))
    updated = skill_registry.update_skill(skill_id, payload.model_dump(exclude_unset=True))
    if not updated:
        raise HTTPException(status_code=404, detail="Skill not found")
    return updated


@app.delete("/api/skills/{skill_id}")
def delete_skill(skill_id: str) -> dict[str, Any]:
    skill = skill_registry.get_skill(skill_id)
    if not skill:
        raise HTTPException(status_code=404, detail="Skill not found")
    if skill.get("category") == "builtin":
        raise HTTPException(status_code=400, detail="平台内置 Skill 不能卸载，可在编辑器中停用")
    for agent in list_agents():
        bound = [item for item in agent.get("skills", []) if item != skill_id]
        if bound != agent.get("skills", []):
            db.execute("UPDATE agents SET skills_json = ?, updated_at = ? WHERE id = ?", (db.json_dumps(bound), db.utc_now(), agent["id"]))
    skill_registry.delete_package(skill_id)
    db.execute("DELETE FROM skills WHERE id = ?", (skill_id,))
    return {"ok": True, "id": skill_id}


@app.get("/api/mcp")
def list_mcp_servers() -> list[dict[str, Any]]:
    return mcp_gateway.list_servers()


@app.post("/api/mcp")
def create_mcp_server(payload: McpServerCreate) -> dict[str, Any]:
    if mcp_gateway.server_exists(payload.id):
        raise HTTPException(status_code=409, detail="MCP server already exists")
    return mcp_gateway.create_server(payload.model_dump())


@app.post("/api/mcp/import")
async def import_mcp_config(file: UploadFile = File(...)) -> list[dict[str, Any]]:
    raw = await file.read(1024 * 1024 + 1)
    if len(raw) > 1024 * 1024:
        raise HTTPException(status_code=413, detail="MCP 配置文件不能超过 1MB")
    try:
        payload = json.loads(raw.decode("utf-8"))
        return mcp_gateway.import_config(payload)
    except (ValueError, UnicodeDecodeError, json.JSONDecodeError) as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc


@app.post("/api/mcp/install/url")
async def install_mcp_url(payload: RemoteInstall) -> list[dict[str, Any]]:
    try:
        return await _install_mcp_remote_url(payload.url)
    except (ValueError, UnicodeDecodeError, json.JSONDecodeError, httpx.HTTPError) as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc


@app.put("/api/mcp/{server_id}")
def update_mcp_server(server_id: str, payload: McpServerUpdate) -> dict[str, Any]:
    updated = mcp_gateway.update_server(server_id, payload.model_dump(exclude_unset=True))
    if not updated:
        raise HTTPException(status_code=404, detail="MCP server not found")
    return updated


@app.delete("/api/mcp/{server_id}")
def delete_mcp_server(server_id: str) -> dict[str, Any]:
    server = mcp_gateway.get_server(server_id)
    if not server:
        raise HTTPException(status_code=404, detail="MCP server not found")
    if server.get("kind") == "builtin":
        raise HTTPException(status_code=400, detail="平台内置 MCP 不能卸载，可在编辑器中停用")
    for agent in list_agents():
        bound = [item for item in agent.get("mcp_servers", []) if item != server_id]
        if bound != agent.get("mcp_servers", []):
            db.execute("UPDATE agents SET mcp_servers_json = ?, updated_at = ? WHERE id = ?", (db.json_dumps(bound), db.utc_now(), agent["id"]))
    db.execute("DELETE FROM mcp_servers WHERE id = ?", (server_id,))
    return {"ok": True, "id": server_id}


@app.get("/api/mcp/{server_id}")
def get_mcp_server(server_id: str) -> dict[str, Any]:
    server = mcp_gateway.get_server(server_id)
    if not server:
        raise HTTPException(status_code=404, detail="MCP server not found")
    return server


@app.get("/api/mcp/{server_id}/tools")
def list_mcp_tools(server_id: str) -> list[dict[str, Any]]:
    if not mcp_gateway.get_server(server_id):
        raise HTTPException(status_code=404, detail="MCP server not found")
    return mcp_gateway.list_tools(server_id)


@app.post("/api/mcp/{server_id}/discover")
async def discover_mcp_tools(server_id: str) -> list[dict[str, Any]]:
    try:
        return await mcp_gateway.discover_tools(server_id)
    except (ToolError, ImportError, OSError) as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc


def model_to_api(row: dict[str, Any]) -> dict[str, Any]:
    legacy_key = str(row.get("api_key_env") or "").strip()
    has_direct_key = bool(row.get("api_key_ciphertext")) or bool(legacy_key and not _is_env_name(legacy_key))
    safe = {k: v for k, v in row.items() if k not in {"api_key_ciphertext", "config_json"}}
    safe["api_key_env"] = legacy_key if _is_env_name(legacy_key) else ""
    return {**safe, "enabled": bool(row.get("enabled")), "config": db.json_loads(row.get("config_json"), {}), "has_api_key": has_direct_key, "api_key_mode": "direct" if has_direct_key else "env"}


@app.get("/api/models")
def list_models() -> list[dict[str, Any]]:
    deterministic = {"id": "deterministic", "name": "离线确定性模型", "provider": "deterministic", "model": "deterministic-offline", "base_url": "", "api_key_env": "", "enabled": True, "config": {}, "has_api_key": False, "api_key_mode": "env"}
    return [deterministic] + [model_to_api(r) for r in db.query_all("SELECT * FROM model_configs ORDER BY name")]


@app.post("/api/models")
def create_model(payload: ModelConfigCreate) -> dict[str, Any]:
    if payload.id == "deterministic" or db.query_one("SELECT id FROM model_configs WHERE id = ?", (payload.id,)):
        raise HTTPException(status_code=409, detail="Model config already exists")
    now = db.utc_now()
    if payload.api_key_mode not in {"env", "direct"}:
        raise HTTPException(status_code=400, detail="api_key_mode 必须是 env 或 direct")
    if payload.api_key_mode == "direct" and not payload.api_key:
        raise HTTPException(status_code=400, detail="直接密钥模式必须填写 API Key")
    if payload.api_key_mode == "env" and not _is_env_name(payload.api_key_env):
        raise HTTPException(status_code=400, detail="环境变量模式必须填写合法变量名，例如 OPENAI_API_KEY")
    encrypted = secret_store.encrypt(payload.api_key) if payload.api_key_mode == "direct" and payload.api_key else ""
    api_key_env = payload.api_key_env if payload.api_key_mode == "env" else ""
    db.execute(
        "INSERT INTO model_configs(id, name, provider, model, base_url, api_key_env, api_key_ciphertext, enabled, config_json, created_at, updated_at) VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)",
        (payload.id, payload.name, payload.provider, payload.model, payload.base_url, api_key_env, encrypted, 1 if payload.enabled else 0, db.json_dumps(payload.config), now, now),
    )
    return model_to_api(db.query_one("SELECT * FROM model_configs WHERE id = ?", (payload.id,)) or {})


@app.put("/api/models/{model_id}")
def update_model(model_id: str, payload: ModelConfigUpdate) -> dict[str, Any]:
    current = db.query_one("SELECT * FROM model_configs WHERE id = ?", (model_id,))
    if not current:
        raise HTTPException(status_code=404, detail="Model config not found")
    current_api = model_to_api(current)
    merged = {**current_api, **payload.model_dump(exclude_unset=True)}
    mode = payload.api_key_mode or current_api.get("api_key_mode", "env")
    if mode not in {"env", "direct"}:
        raise HTTPException(status_code=400, detail="api_key_mode 必须是 env 或 direct")
    encrypted = current.get("api_key_ciphertext") or ""
    if mode == "env":
        encrypted = ""
        if not _is_env_name(str(merged.get("api_key_env") or "")):
            raise HTTPException(status_code=400, detail="环境变量模式必须填写合法变量名，例如 OPENAI_API_KEY")
    elif payload.api_key:
        encrypted = secret_store.encrypt(payload.api_key)
    elif not encrypted:
        raise HTTPException(status_code=400, detail="直接密钥模式必须填写 API Key")
    db.execute(
        "UPDATE model_configs SET name = ?, provider = ?, model = ?, base_url = ?, api_key_env = ?, api_key_ciphertext = ?, enabled = ?, config_json = ?, updated_at = ? WHERE id = ?",
        (merged["name"], merged["provider"], merged["model"], merged.get("base_url", ""), merged.get("api_key_env", "") if mode == "env" else "", encrypted, 1 if merged.get("enabled") else 0, db.json_dumps(merged.get("config", {})), db.utc_now(), model_id),
    )
    return model_to_api(db.query_one("SELECT * FROM model_configs WHERE id = ?", (model_id,)) or {})


@app.delete("/api/models/{model_id}")
def delete_model(model_id: str) -> dict[str, Any]:
    if model_id == "deterministic":
        raise HTTPException(status_code=400, detail="内置离线模型不能删除")
    if not db.query_one("SELECT id FROM model_configs WHERE id = ?", (model_id,)):
        raise HTTPException(status_code=404, detail="Model config not found")
    db.execute("UPDATE agents SET model = 'deterministic', updated_at = ? WHERE model = ?", (db.utc_now(), model_id))
    db.execute("DELETE FROM model_configs WHERE id = ?", (model_id,))
    return {"ok": True, "id": model_id}


@app.post("/api/models/{model_id}/test")
async def test_model(model_id: str) -> dict[str, Any]:
    try:
        response = await model_gateway.summarize(
            "这是连接测试。请只回复 OK。",
            {"system_prompt": "你正在执行模型连接测试。"},
            model_config_id=model_id,
        )
        return {"ok": True, "model_id": model_id, "response": response[:1000]}
    except Exception as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc


@app.post("/api/mcp/{server_id}/tools/{tool_name}/invoke")
async def invoke_mcp_tool(server_id: str, tool_name: str, payload: ToolInvokeRequest) -> dict[str, Any]:
    definition = mcp_gateway.get_tool_definition(server_id, tool_name)
    if not definition:
        raise HTTPException(status_code=404, detail="MCP 工具不存在，请先同步工具清单")
    server_kind = str(definition.get("server_kind") or "")
    trusted_read_only = (
        server_kind == "builtin" and definition.get("effect") == "read"
    ) or (
        server_kind != "builtin"
        and definition.get("annotations", {}).get("readOnlyHint") is True
    )
    if not trusted_read_only:
        raise HTTPException(
            status_code=403,
            detail=(
                "页面测试调用只允许明确标注的只读工具；写入、破坏性或未标注工具"
                "请通过正式对话任务执行，以应用智能体权限、Policy 和人工审批。"
            ),
        )
    policy_context = {
        "organization_id": "local-org",
        "workspace_id": "default",
        "user_id": "local-user",
        "task_id": payload.task_id or "",
        "tool": {
            "server": server_id,
            "server_id": server_id,
            "name": tool_name,
            "tool_name": tool_name,
            "arguments": payload.arguments,
            "direct_test": True,
        },
    }
    try:
        evaluation = await policy_engine.evaluate("tool.before", policy_context)
        if evaluation.denied:
            raise HTTPException(status_code=403, detail=evaluation.summary)
        if evaluation.requires_approval:
            raise HTTPException(
                status_code=403,
                detail=(
                    f"{evaluation.summary}。页面测试调用不会创建审批；"
                    "请通过正式对话任务执行并完成审批。"
                ),
            )
        modified = evaluation.apply(policy_context)
        arguments = modified.get("tool", {}).get("arguments", payload.arguments)
        if not isinstance(arguments, dict):
            raise HTTPException(status_code=400, detail="Policy 修改后的工具参数必须是 JSON 对象")
        return await mcp_gateway.invoke_tool(server_id, tool_name, arguments, task_id=payload.task_id)
    except ToolError as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc


@app.get("/api/agents")
def list_agents() -> list[dict[str, Any]]:
    rows = db.query_all("SELECT * FROM agents ORDER BY name")
    return [
        {
            **r,
            "skills": db.json_loads(r.get("skills_json"), []),
            "mcp_servers": db.json_loads(r.get("mcp_servers_json"), []),
            "permissions": db.json_loads(r.get("permissions_json"), {}),
        }
        for r in rows
    ]


def _validate_agent_bindings(model_id: str, skill_ids: list[str], mcp_ids: list[str]) -> None:
    if model_id != "deterministic" and not db.query_one(
        "SELECT id FROM model_configs WHERE id = ? AND enabled = 1", (model_id,)
    ):
        raise HTTPException(status_code=400, detail="智能体选择的模型不存在或未启用")
    missing_skills = [item for item in dict.fromkeys(skill_ids) if not skill_registry.get_skill(item)]
    if missing_skills:
        raise HTTPException(status_code=400, detail="智能体引用了不存在的 Skill：" + "、".join(missing_skills))
    missing_mcps = [item for item in dict.fromkeys(mcp_ids) if not mcp_gateway.server_exists(item)]
    if missing_mcps:
        raise HTTPException(status_code=400, detail="智能体引用了不存在的 MCP：" + "、".join(missing_mcps))


@app.post("/api/agents")
def create_agent(payload: AgentCreate) -> dict[str, Any]:
    if db.query_one("SELECT id FROM agents WHERE id = ?", (payload.id,)):
        raise HTTPException(status_code=409, detail="Agent already exists")
    _validate_agent_bindings(payload.model, payload.skills, payload.mcp_servers)
    now = db.utc_now()
    db.execute(
        """
        INSERT INTO agents(id, name, description, model, system_prompt, skills_json, mcp_servers_json, permissions_json, created_at, updated_at)
        VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
        """,
        (
            payload.id, payload.name, payload.description, payload.model, payload.system_prompt,
            db.json_dumps(payload.skills), db.json_dumps(payload.mcp_servers), db.json_dumps(payload.permissions), now, now,
        ),
    )
    return next(a for a in list_agents() if a["id"] == payload.id)


@app.put("/api/agents/{agent_id}")
def update_agent(agent_id: str, payload: AgentUpdate) -> dict[str, Any]:
    current = db.query_one("SELECT * FROM agents WHERE id = ?", (agent_id,))
    if not current:
        raise HTTPException(status_code=404, detail="Agent not found")
    current_api = {
        **current,
        "skills": db.json_loads(current.get("skills_json"), []),
        "mcp_servers": db.json_loads(current.get("mcp_servers_json"), []),
        "permissions": db.json_loads(current.get("permissions_json"), {}),
    }
    incoming = payload.model_dump(exclude_unset=True)
    merged = {**current_api, **incoming}
    _validate_agent_bindings(merged["model"], merged.get("skills", []), merged.get("mcp_servers", []))
    db.execute(
        """
        UPDATE agents SET name = ?, description = ?, model = ?, system_prompt = ?, skills_json = ?, mcp_servers_json = ?, permissions_json = ?, updated_at = ?
        WHERE id = ?
        """,
        (
            merged["name"], merged["description"], merged["model"], merged["system_prompt"],
            db.json_dumps(merged.get("skills", [])), db.json_dumps(merged.get("mcp_servers", [])), db.json_dumps(merged.get("permissions", {})),
            db.utc_now(), agent_id,
        ),
    )
    return next(a for a in list_agents() if a["id"] == agent_id)


@app.get("/api/tasks")
def list_tasks() -> list[dict[str, Any]]:
    rows = db.query_all("SELECT * FROM tasks ORDER BY created_at DESC LIMIT 100")
    return [_public_task(r) for r in rows]


def _public_attachment(value: dict[str, Any]) -> dict[str, Any]:
    """Expose attachment metadata without leaking the server storage path."""
    allowed = {"id", "name", "content_type", "size", "created_at"}
    return {
        key: value.get(key)
        for key in allowed
        if value.get(key) not in (None, "")
    }


def _public_artifact(value: dict[str, Any]) -> dict[str, Any]:
    artifact_id = str(value.get("id") or "")
    allowed = {
        "id", "task_id", "run_id", "workspace_id", "name", "kind",
        "mime_type", "size", "sha256", "version", "created_at", "metadata",
    }
    result = {key: value.get(key) for key in allowed if value.get(key) not in (None, "")}
    if artifact_id:
        result["download_url"] = f"/api/artifacts/{artifact_id}/download"
        result["preview_url"] = f"/api/artifacts/{artifact_id}/preview"
    elif value.get("download_url"):
        result["download_url"] = value["download_url"]
    return result


def _sanitize_public_payload(value: Any) -> Any:
    if isinstance(value, list):
        return [_sanitize_public_payload(item) for item in value]
    if not isinstance(value, dict):
        return value
    looks_like_artifact = bool(
        str(value.get("id") or "").startswith("art_")
        or value.get("download_url")
        or (value.get("name") and value.get("kind") and ("path" in value or "relative_path" in value))
    )
    if looks_like_artifact:
        return _public_artifact(value)
    return {
        key: _sanitize_public_payload(item)
        for key, item in value.items()
        if key not in {"api_key", "api_key_ciphertext", "path", "relative_path"}
    }


def _public_event(value: dict[str, Any]) -> dict[str, Any]:
    """Expose parsed event data and omit its duplicate internal JSON column."""
    result = {key: item for key, item in value.items() if key != "data_json"}
    result["data"] = _sanitize_public_payload(
        db.json_loads(value.get("data_json"), {})
    )
    return result


_PRIVATE_TASK_EVENT_TYPES = frozenset({"analysis", "reasoning", "thought", "thinking"})
_TASK_STREAM_TERMINAL_STATUSES = frozenset({"completed", "failed", "cancelled"})


def _is_public_task_event(value: dict[str, Any]) -> bool:
    """Keep private model-working events out of every browser-facing event feed."""
    return str(value.get("type") or "").lower() not in _PRIVATE_TASK_EVENT_TYPES


def _task_event_cursor(
    request: Request,
    *,
    cursor: int | None,
    after_id: int | None,
) -> int:
    """Resolve an SSE resume cursor from the query string or Last-Event-ID."""
    explicit = cursor if cursor is not None else after_id
    raw: Any = explicit if explicit is not None else request.headers.get("last-event-id", "0")
    try:
        resolved = int(raw or 0)
    except (TypeError, ValueError) as exc:
        raise HTTPException(status_code=400, detail="事件游标必须是非负整数") from exc
    if resolved < 0:
        raise HTTPException(status_code=400, detail="事件游标必须是非负整数")
    return resolved


def _public_runtime_record(value: dict[str, Any]) -> dict[str, Any]:
    """Remove duplicate storage columns from TaskState API records."""
    internal_suffix_fields = {
        "result_json", "error_json", "metadata_json", "input_json",
        "output_json", "payload_json", "state_json",
        "last_restore_metadata_json", "command_type",
    }
    return {
        key: _sanitize_public_payload(item)
        for key, item in value.items()
        if key not in internal_suffix_fields
    }


def _public_task(
    value: dict[str, Any],
    *,
    include_result: bool = True,
    include_attachments: bool = True,
) -> dict[str, Any]:
    """Return the stable public task shape and remove internal JSON/path fields."""
    internal_fields = {"result_json", "artifacts_json", "attachments_json"}
    result = {key: item for key, item in value.items() if key not in internal_fields}
    if include_result:
        result["result"] = _sanitize_public_payload(
            db.json_loads(value.get("result_json"), {})
        )
    result["artifacts"] = [
        _public_artifact(item)
        for item in db.json_loads(value.get("artifacts_json"), [])
        if isinstance(item, dict)
    ]
    if include_attachments:
        result["attachments"] = [
            _public_attachment(item)
            for item in db.json_loads(value.get("attachments_json"), [])
            if isinstance(item, dict)
        ]
    return result


def _task_or_404(task_id: str) -> dict[str, Any]:
    task = db.query_one("SELECT * FROM tasks WHERE id = ?", (task_id,))
    if not task:
        raise HTTPException(status_code=404, detail="任务不存在")
    return task


def _runtime_projection(task_id: str) -> dict[str, Any]:
    _task_or_404(task_id)
    raw_runs = task_state.list_runs(task_id=task_id, limit=100)
    runs = [
        {
            key: item.get(key)
            for key in (
                "id", "task_id", "attempt", "status", "current_node_id",
                "resumed_from_checkpoint_id", "started_at", "finished_at",
                "created_at", "updated_at", "metadata",
            )
        }
        for item in raw_runs
    ]
    attempts = {item["id"]: item["attempt"] for item in runs}
    nodes: list[dict[str, Any]] = []
    for run in sorted(runs, key=lambda item: item["attempt"]):
        nodes.extend(
            {
                **{
                    key: node.get(key)
                    for key in (
                        "id", "run_id", "task_id", "node_key", "parent_node_id",
                        "title", "kind", "sequence", "status", "started_at",
                        "finished_at", "created_at", "updated_at", "metadata",
                    )
                },
                "output_summary": str((node.get("output") or {}).get("summary") or "")[:500],
                "error_summary": str((node.get("error") or {}).get("message") or "")[:500],
                "attempt": attempts.get(run["id"], 1),
            }
            for node in task_state.list_nodes(run["id"])
        )
    checkpoints = [
        {
            **{
                key: item.get(key)
                for key in (
                    "id", "task_id", "run_id", "node_id", "sequence", "reason",
                    "restored_at", "restore_count", "created_at", "metadata",
                    "last_restore_metadata",
                )
            },
            "label": item.get("reason") or f"检查点 {item.get('sequence', '')}",
            "restorable": True,
            "attempt": attempts.get(item.get("run_id", ""), 1),
        }
        for item in task_state.list_checkpoints(task_id=task_id, include_state=False, limit=200)
    ]
    commands = [
        {
            key: item.get(key)
            for key in (
                "id", "task_id", "run_id", "type", "payload", "status", "priority",
                "worker_id", "result", "error", "available_at", "created_at",
                "claimed_at", "completed_at", "updated_at",
            )
        }
        for item in task_state.list_commands(task_id=task_id, limit=200)
    ]
    active_run = next(
        (
            item
            for item in runs
            if item["status"] in {"queued", "running", "paused", "waiting_approval"}
        ),
        None,
    )
    return {
        "runs": runs,
        "nodes": nodes,
        "checkpoints": checkpoints,
        "commands": commands,
        "active_run": active_run,
    }


def _active_run_or_409(task_id: str) -> dict[str, Any]:
    runtime_state = _runtime_projection(task_id)
    active = runtime_state.get("active_run")
    if not isinstance(active, dict):
        raise HTTPException(status_code=409, detail="当前任务没有可控制的运行尝试")
    return active


def _ensure_no_active_run(task_id: str) -> None:
    active = _runtime_projection(task_id).get("active_run")
    if isinstance(active, dict):
        raise HTTPException(
            status_code=409,
            detail=f"任务已有 {active.get('status', 'active')} 状态的运行尝试，请先取消或等待结束",
        )


def _complete_control_command(
    task_id: str,
    command_type: str,
    *,
    payload: dict[str, Any],
    run_id: str | None = None,
    result: dict[str, Any] | None = None,
) -> dict[str, Any]:
    command = task_state.enqueue_command(
        task_id,
        command_type,
        payload=payload,
        run_id=run_id,
    )
    claimed = task_state.claim_command(
        "api-control",
        task_id=task_id,
        run_id=run_id,
        command_types=[command_type],
    )
    if claimed:
        task_state.complete_command(claimed["id"], result=result or {"accepted": True})
    return task_state.get_command(command["id"]) or command


async def _retry_task(task_id: str) -> dict[str, Any]:
    task = _task_or_404(task_id)
    _ensure_no_active_run(task_id)
    if task.get("status") not in {"failed", "cancelled", "completed"}:
        raise HTTPException(status_code=409, detail="只有已结束的任务才能重试")
    run = task_state.create_run(task_id, metadata={"trigger": "retry"})
    command = _complete_control_command(
        task_id,
        "retry",
        payload={},
        run_id=run["id"],
        result={"run_id": run["id"]},
    )
    db.execute(
        "UPDATE tasks SET status = 'queued', result_json = '{}', updated_at = ? WHERE id = ?",
        (db.utc_now(), task_id),
    )
    db.insert_event(
        task_id,
        "retry_scheduled",
        "已创建重试运行",
        f"这是第 {run['attempt']} 次运行尝试。",
        {"run_id": run["id"], "attempt": run["attempt"]},
    )
    _schedule_runtime(task_id, run["id"])
    return {"ok": True, "run": run, "command": command}


async def _resume_task_from_checkpoint(
    task_id: str,
    checkpoint_id: str | None,
    *,
    trigger: str,
) -> dict[str, Any]:
    _task_or_404(task_id)
    _ensure_no_active_run(task_id)
    checkpoint = (
        task_state.get_checkpoint(checkpoint_id, include_state=False)
        if checkpoint_id
        else (task_state.list_checkpoints(task_id=task_id, include_state=False, limit=1) or [None])[0]
    )
    if not checkpoint:
        raise HTTPException(status_code=404, detail="当前任务没有可恢复的检查点")
    if checkpoint.get("task_id") != task_id:
        raise HTTPException(status_code=400, detail="该检查点不属于当前任务")
    restored = task_state.restore_checkpoint(
        checkpoint["id"],
        restore_metadata={"requested_by": "api", "trigger": trigger},
    )
    run = task_state.create_run(
        task_id,
        resumed_from_checkpoint_id=checkpoint["id"],
        metadata={"trigger": trigger, "checkpoint_restore_audited": True},
    )
    command = _complete_control_command(
        task_id,
        "restore_checkpoint" if trigger == "restore_checkpoint" else "resume",
        payload={"checkpoint_id": checkpoint["id"]},
        run_id=run["id"],
        result={"run_id": run["id"], "checkpoint_id": checkpoint["id"]},
    )
    db.execute(
        "UPDATE tasks SET status = 'queued', result_json = '{}', updated_at = ? WHERE id = ?",
        (db.utc_now(), task_id),
    )
    db.insert_event(
        task_id,
        "resume_scheduled",
        "已创建恢复运行",
        "新的运行尝试将从所选安全检查点继续。",
        {"run_id": run["id"], "checkpoint_id": checkpoint["id"], "restore_count": restored["restore_count"]},
    )
    _schedule_runtime(task_id, run["id"])
    return {"ok": True, "run": run, "checkpoint": restored, "command": command}


@app.get("/api/tasks/{task_id}/runtime")
def get_task_runtime(task_id: str) -> dict[str, Any]:
    return _runtime_projection(task_id)


@app.post("/api/tasks/{task_id}/commands", status_code=202)
async def send_task_command(task_id: str, payload: TaskCommandRequest) -> dict[str, Any]:
    _task_or_404(task_id)
    command_type = payload.type.strip().lower()
    if command_type not in {"message", "cancel", "retry", "resume", "restore_checkpoint"}:
        raise HTTPException(status_code=400, detail=f"不支持的任务指令：{command_type}")
    if command_type == "retry":
        return await _retry_task(task_id)
    if command_type in {"resume", "restore_checkpoint"}:
        return await _resume_task_from_checkpoint(
            task_id,
            str(payload.payload.get("checkpoint_id") or "") or None,
            trigger=command_type,
        )
    active = _active_run_or_409(task_id)
    if command_type == "cancel":
        command = task_state.request_cancel(
            task_id,
            run_id=active["id"],
            reason=str(payload.payload.get("reason") or "用户请求取消"),
            requested_by="user",
        )
        return {"ok": True, "command": command, "run": active}
    message = str(payload.payload.get("message") or "").strip()
    if not message:
        raise HTTPException(status_code=400, detail="追加指令不能为空")
    command = task_state.enqueue_command(
        task_id,
        "message",
        run_id=active["id"],
        payload={"message": message},
        priority=20,
    )
    db.insert_event(
        task_id,
        "command_queued",
        "已加入运行中指令",
        message,
        {"command_id": command["id"], "run_id": active["id"]},
    )
    return {"ok": True, "command": command, "run": active}


@app.post("/api/tasks/{task_id}/cancel", status_code=202)
async def cancel_task(task_id: str, payload: dict[str, Any] | None = None) -> dict[str, Any]:
    active = _active_run_or_409(task_id)
    values = payload or {}
    command = task_state.request_cancel(
        task_id,
        run_id=active["id"],
        reason=str(values.get("reason") or "用户请求取消"),
        requested_by="user",
    )
    return {"ok": True, "command": command, "run": active}


@app.post("/api/tasks/{task_id}/retry", status_code=202)
async def retry_task(task_id: str) -> dict[str, Any]:
    return await _retry_task(task_id)


@app.post("/api/tasks/{task_id}/resume", status_code=202)
async def resume_task(task_id: str, payload: TaskResumeRequest | None = None) -> dict[str, Any]:
    return await _resume_task_from_checkpoint(
        task_id,
        payload.checkpoint_id if payload else None,
        trigger="resume",
    )


@app.post("/api/tasks/{task_id}/checkpoints/{checkpoint_id}/restore", status_code=202)
async def restore_task_checkpoint(
    task_id: str,
    checkpoint_id: str,
    payload: CheckpointRestoreRequest | None = None,
) -> dict[str, Any]:
    _ = payload
    return await _resume_task_from_checkpoint(
        task_id,
        checkpoint_id,
        trigger="restore_checkpoint",
    )


def _validate_loop_bindings(agent_id: str, model_id: str) -> None:
    if not db.query_one("SELECT id FROM agents WHERE id = ?", (agent_id,)):
        raise HTTPException(status_code=400, detail="所选智能体不存在")
    if model_id != "deterministic" and not db.query_one("SELECT id FROM model_configs WHERE id = ? AND enabled = 1", (model_id,)):
        raise HTTPException(status_code=400, detail="所选模型不存在或未启用")


def _validate_loop_trigger(config: dict[str, Any], webhook_secret_ciphertext: str) -> None:
    try:
        validate_trigger_config(
            str(config.get("trigger_type") or "interval"),
            cron_expression=str(config.get("cron_expression") or ""),
            once_at=str(config.get("once_at") or ""),
            webhook_secret_configured=bool(webhook_secret_ciphertext),
        )
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc


@app.get("/api/loops")
def list_loops(
    organization_id: str = "local-org",
    workspace_id: str = "default",
    user_id: str = "local-user",
) -> list[dict[str, Any]]:
    return [
        serialize_loop(row)
        for row in db.query_all(
            """SELECT * FROM loops WHERE organization_id = ? AND workspace_id = ? AND user_id = ?
               ORDER BY created_at DESC""",
            (organization_id, workspace_id, user_id),
        )
    ]


@app.post("/api/loops")
def create_loop(payload: LoopCreate) -> dict[str, Any]:
    _validate_loop_bindings(payload.agent_id, payload.model_id)
    loop_id = (payload.id or ("loop_" + uuid.uuid4().hex[:12])).strip()
    if not re.fullmatch(r"[A-Za-z0-9_-]{2,80}", loop_id):
        raise HTTPException(status_code=400, detail="Loop ID 只能包含字母、数字、下划线和连字符")
    if db.query_one("SELECT id FROM loops WHERE id = ?", (loop_id,)):
        raise HTTPException(status_code=409, detail="Loop ID 已存在")
    encrypted_secret = secret_store.encrypt(payload.webhook_secret) if payload.webhook_secret else ""
    config = payload.model_dump(exclude={"webhook_secret", "initial_state", "auto_start", "id"})
    _validate_loop_trigger(config, encrypted_secret)
    now = db.utc_now()
    status = "active" if payload.auto_start else "paused"
    next_run_at = ""
    if payload.auto_start:
        try:
            next_run_at = next_schedule_at(config, immediate_interval=True)
        except ValueError as exc:
            raise HTTPException(status_code=400, detail=str(exc)) from exc
    db.execute(
        """INSERT INTO loops(
               id, name, prompt, agent_id, model_id, trigger_type, interval_seconds,
               cron_expression, once_at, organization_id, workspace_id, user_id,
               webhook_secret_ciphertext, webhook_tolerance_seconds, status, max_runs,
               max_failures, max_attempts, retry_backoff_seconds, state_json, last_diff_json,
               next_run_at, created_at, updated_at
           ) VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, '{}', ?, ?, ?)""",
        (
            loop_id, payload.name, payload.prompt, payload.agent_id, payload.model_id,
            payload.trigger_type, payload.interval_seconds, payload.cron_expression, payload.once_at,
            payload.organization_id, payload.workspace_id, payload.user_id, encrypted_secret,
            payload.webhook_tolerance_seconds, status, payload.max_runs, payload.max_failures,
            payload.max_attempts, payload.retry_backoff_seconds, db.json_dumps(payload.initial_state),
            next_run_at, now, now,
        ),
    )
    return serialize_loop(db.query_one("SELECT * FROM loops WHERE id = ?", (loop_id,)) or {})


@app.get("/api/loops/{loop_id}")
def get_loop(loop_id: str) -> dict[str, Any]:
    row = db.query_one("SELECT * FROM loops WHERE id = ?", (loop_id,))
    if not row:
        raise HTTPException(status_code=404, detail="Loop not found")
    result = serialize_loop(row)
    result["runs"] = [
        serialize_run(item)
        for item in db.query_all(
            "SELECT * FROM loop_runs WHERE loop_id = ? ORDER BY run_number DESC, attempt DESC",
            (loop_id,),
        )
    ]
    return result


@app.put("/api/loops/{loop_id}")
def update_loop(loop_id: str, payload: LoopUpdate) -> dict[str, Any]:
    current = db.query_one("SELECT * FROM loops WHERE id = ?", (loop_id,))
    if not current:
        raise HTTPException(status_code=404, detail="Loop not found")
    incoming = payload.model_dump(exclude_unset=True)
    webhook_secret = incoming.pop("webhook_secret", None)
    state = incoming.pop("state", None)
    merged = {**current, **incoming}
    _validate_loop_bindings(merged["agent_id"], merged["model_id"])
    encrypted_secret = str(current.get("webhook_secret_ciphertext") or "")
    if webhook_secret is not None:
        encrypted_secret = secret_store.encrypt(webhook_secret)
    _validate_loop_trigger(merged, encrypted_secret)
    state_json = db.json_dumps(state) if state is not None else str(current.get("state_json") or "{}")
    next_run_at = str(current.get("next_run_at") or "")
    if current["status"] == "active":
        try:
            next_run_at = next_schedule_at(merged)
        except ValueError as exc:
            raise HTTPException(status_code=400, detail=str(exc)) from exc
    db.execute(
        """UPDATE loops SET name = ?, prompt = ?, agent_id = ?, model_id = ?, trigger_type = ?,
               interval_seconds = ?, cron_expression = ?, once_at = ?, webhook_secret_ciphertext = ?,
               webhook_tolerance_seconds = ?, max_runs = ?, max_failures = ?, max_attempts = ?,
               retry_backoff_seconds = ?, state_json = ?, next_run_at = ?, updated_at = ?
           WHERE id = ?""",
        (
            merged["name"], merged["prompt"], merged["agent_id"], merged["model_id"],
            merged.get("trigger_type") or "interval", merged["interval_seconds"],
            merged.get("cron_expression") or "", merged.get("once_at") or "", encrypted_secret,
            merged.get("webhook_tolerance_seconds") or 300, merged["max_runs"],
            merged["max_failures"], merged.get("max_attempts") or 1,
            merged.get("retry_backoff_seconds") or 0, state_json, next_run_at, db.utc_now(), loop_id,
        ),
    )
    return serialize_loop(db.query_one("SELECT * FROM loops WHERE id = ?", (loop_id,)) or {})


@app.delete("/api/loops/{loop_id}")
def delete_loop(loop_id: str) -> dict[str, bool]:
    current = db.query_one("SELECT status FROM loops WHERE id = ?", (loop_id,))
    if not current:
        raise HTTPException(status_code=404, detail="Loop not found")
    if current["status"] == "running":
        raise HTTPException(status_code=409, detail="循环任务运行中，请先等待本轮结束并暂停")
    db.execute("DELETE FROM loop_runs WHERE loop_id = ?", (loop_id,))
    db.execute("DELETE FROM automation_trigger_events WHERE loop_id = ?", (loop_id,))
    db.execute("DELETE FROM loops WHERE id = ?", (loop_id,))
    return {"ok": True}


@app.post("/api/loops/{loop_id}/start")
def start_loop(loop_id: str) -> dict[str, Any]:
    current = db.query_one("SELECT * FROM loops WHERE id = ?", (loop_id,))
    if not current:
        raise HTTPException(status_code=404, detail="Loop not found")
    if current["status"] == "running":
        raise HTTPException(status_code=409, detail="循环任务正在运行")
    if int(current["run_count"]) >= int(current["max_runs"]):
        raise HTTPException(status_code=409, detail="已达到最大轮数，请提高最大轮数后再启动")
    if (current.get("trigger_type") or "interval") == "once" and int(current["run_count"]):
        raise HTTPException(status_code=409, detail="一次性自动化已经执行过，请新建自动化")
    _validate_loop_trigger(current, str(current.get("webhook_secret_ciphertext") or ""))
    now = db.utc_now()
    try:
        next_run_at = next_schedule_at(current, immediate_interval=True)
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc
    db.execute(
        "UPDATE loops SET status = 'active', next_run_at = ?, updated_at = ? WHERE id = ?",
        (next_run_at, now, loop_id),
    )
    return serialize_loop(db.query_one("SELECT * FROM loops WHERE id = ?", (loop_id,)) or {})


@app.post("/api/loops/{loop_id}/pause")
def pause_loop(loop_id: str) -> dict[str, Any]:
    if not db.query_one("SELECT id FROM loops WHERE id = ?", (loop_id,)):
        raise HTTPException(status_code=404, detail="Loop not found")
    db.execute("UPDATE loops SET status = 'paused', next_run_at = '', updated_at = ? WHERE id = ?", (db.utc_now(), loop_id))
    return serialize_loop(db.query_one("SELECT * FROM loops WHERE id = ?", (loop_id,)) or {})


@app.post("/api/loops/{loop_id}/run", status_code=202)
async def run_loop_now(loop_id: str) -> dict[str, Any]:
    if not db.query_one("SELECT id FROM loops WHERE id = ?", (loop_id,)):
        raise HTTPException(status_code=404, detail="Loop not found")
    try:
        background = loop_scheduler.dispatch_once(loop_id)
        _runtime_tasks.add(background)
        background.add_done_callback(_runtime_tasks.discard)
        return {"accepted": True, "loop_id": loop_id, "status": "queued"}
    except RuntimeError as exc:
        raise HTTPException(status_code=409, detail=str(exc)) from exc


@app.get("/api/loops/{loop_id}/runs")
def list_loop_runs(loop_id: str) -> list[dict[str, Any]]:
    if not db.query_one("SELECT id FROM loops WHERE id = ?", (loop_id,)):
        raise HTTPException(status_code=404, detail="Loop not found")
    return [
        serialize_run(item)
        for item in db.query_all(
            "SELECT * FROM loop_runs WHERE loop_id = ? ORDER BY run_number DESC, attempt DESC",
            (loop_id,),
        )
    ]


@app.get("/api/loops/{loop_id}/trigger-events")
def list_loop_trigger_events(loop_id: str) -> list[dict[str, Any]]:
    if not db.query_one("SELECT id FROM loops WHERE id = ?", (loop_id,)):
        raise HTTPException(status_code=404, detail="Loop not found")
    return [
        serialize_trigger_event(item)
        for item in db.query_all(
            "SELECT * FROM automation_trigger_events WHERE loop_id = ? ORDER BY received_at DESC",
            (loop_id,),
        )
    ]


def _webhook_timestamp(value: str) -> float:
    try:
        return float(value)
    except ValueError:
        try:
            parsed = datetime.fromisoformat(value.replace("Z", "+00:00"))
            if not parsed.tzinfo:
                parsed = parsed.replace(tzinfo=timezone.utc)
            return parsed.timestamp()
        except ValueError as exc:
            raise HTTPException(status_code=401, detail="Webhook 时间戳无效") from exc


@app.post("/api/loops/{loop_id}/webhook", status_code=202)
async def trigger_loop_webhook(loop_id: str, request: Request) -> dict[str, Any]:
    loop = db.query_one("SELECT * FROM loops WHERE id = ?", (loop_id,))
    if not loop:
        raise HTTPException(status_code=404, detail="Loop not found")
    if (loop.get("trigger_type") or "interval") != "webhook":
        raise HTTPException(status_code=409, detail="这个自动化不是 Webhook 触发类型")
    if loop["status"] != "active":
        raise HTTPException(status_code=409, detail="Webhook 自动化已暂停")
    raw_body = await request.body()
    if len(raw_body) > 1024 * 1024:
        raise HTTPException(status_code=413, detail="Webhook 请求正文不能超过 1MB")
    timestamp = request.headers.get("x-automation-timestamp") or request.headers.get("x-webhook-timestamp") or ""
    provided_signature = request.headers.get("x-automation-signature") or request.headers.get("x-webhook-signature") or ""
    idempotency_key = request.headers.get("idempotency-key", "").strip()
    if not timestamp or not provided_signature:
        raise HTTPException(status_code=401, detail="缺少 Webhook 时间戳或签名")
    if not idempotency_key or len(idempotency_key) > 200:
        raise HTTPException(status_code=400, detail="必须提供 1-200 位 Idempotency-Key")
    tolerance = int(loop.get("webhook_tolerance_seconds") or 300)
    if abs(time.time() - _webhook_timestamp(timestamp)) > tolerance:
        raise HTTPException(status_code=401, detail="Webhook 时间戳超出允许时间窗")
    encrypted_secret = str(loop.get("webhook_secret_ciphertext") or "")
    if not encrypted_secret:
        raise HTTPException(status_code=409, detail="Webhook 签名密钥未配置")
    try:
        secret = secret_store.decrypt(encrypted_secret)
    except RuntimeError as exc:
        raise HTTPException(status_code=500, detail="Webhook 签名密钥无法解密") from exc
    expected = hmac.new(
        secret.encode("utf-8"), timestamp.encode("utf-8") + b"." + raw_body, hashlib.sha256
    ).hexdigest()
    supplied = provided_signature.removeprefix("sha256=").strip().lower()
    if not hmac.compare_digest(expected, supplied):
        raise HTTPException(status_code=401, detail="Webhook 签名验证失败")
    try:
        event, duplicate = create_webhook_event(loop, idempotency_key, raw_body)
    except IdempotencyConflictError as exc:
        raise HTTPException(status_code=409, detail=str(exc)) from exc
    queued = True
    if event.get("status") == "accepted" and not loop_scheduler.is_busy(loop_id):
        try:
            background = loop_scheduler.dispatch_once(
                loop_id, scheduled=True, trigger_event_id=event["id"], trigger_type="webhook"
            )
            _runtime_tasks.add(background)
            background.add_done_callback(_runtime_tasks.discard)
            queued = False
        except RuntimeError:
            queued = True
    return {"accepted": True, "duplicate": duplicate, "queued": queued, "event": event}


@app.get("/api/notifications")
def list_notifications(
    organization_id: str = "local-org",
    workspace_id: str = "default",
    user_id: str = "local-user",
    status: str = "",
    limit: int = 100,
) -> list[dict[str, Any]]:
    if status and status not in {"unread", "read"}:
        raise HTTPException(status_code=400, detail="通知状态只能是 unread 或 read")
    capped_limit = max(1, min(int(limit), 500))
    sql = """SELECT * FROM notifications
             WHERE organization_id = ? AND workspace_id = ? AND user_id = ?"""
    params: list[Any] = [organization_id, workspace_id, user_id]
    if status:
        sql += " AND status = ?"
        params.append(status)
    sql += " ORDER BY created_at DESC LIMIT ?"
    params.append(capped_limit)
    return [serialize_notification(item) for item in db.query_all(sql, params)]


@app.post("/api/notifications/{notification_id}/read")
def read_notification(notification_id: str) -> dict[str, Any]:
    current = db.query_one("SELECT * FROM notifications WHERE id = ?", (notification_id,))
    if not current:
        raise HTTPException(status_code=404, detail="Notification not found")
    now = db.utc_now()
    db.execute(
        "UPDATE notifications SET status = 'read', read_at = ? WHERE id = ?",
        (now, notification_id),
    )
    return serialize_notification(
        db.query_one("SELECT * FROM notifications WHERE id = ?", (notification_id,)) or {}
    )


@app.post("/api/tasks")
async def create_task(payload: TaskCreate) -> dict[str, Any]:
    attachments = []
    for upload_id in payload.attachment_ids[:10]:
        upload = db.query_one("SELECT * FROM uploads WHERE id = ?", (upload_id,))
        if upload:
            attachments.append(upload)
    if payload.model_id and payload.model_id != "deterministic":
        model = db.query_one("SELECT id FROM model_configs WHERE id = ? AND enabled = 1", (payload.model_id,))
        if not model:
            raise HTTPException(status_code=400, detail="所选模型不存在或未启用")
    if payload.executor_type == "agent" and not db.query_one(
        "SELECT id FROM agents WHERE id = ?", (payload.agent_id,)
    ):
        raise HTTPException(status_code=400, detail="所选智能体不存在")
    scope = _api_scope(payload.organization_id, payload.workspace, payload.user_id)
    executor_id = payload.executor_id or payload.agent_id
    selected_agent_id = payload.agent_id
    expert_selection: dict[str, Any] | None = None
    if payload.executor_type == "team":
        recommendation: dict[str, Any] | None = None
        if payload.executor_id:
            team = expert_team_service.get_team(payload.executor_id, scope)
        else:
            try:
                recommendation = expert_team_service.recommend_team(payload.message, scope)
                team = recommendation["team"]
            except (ExpertNotFoundError, ExpertValidationError) as exc:
                raise _expert_http_error(exc) from exc
        if not team or not team.get("enabled"):
            raise HTTPException(status_code=400, detail="所选专家团不存在、不可见或已停用")
        selected_agent_id = team["supervisor_agent_id"]
        executor_id = str(team["id"])
        expert_selection = _public_expert_selection(
            team,
            scope,
            automatic=not bool(payload.executor_id),
            recommendation=recommendation,
        )
    task = create_task_record(
        payload.message,
        selected_agent_id,
        payload.workspace,
        attachments=attachments,
        model_id=payload.model_id,
        conversation_id=payload.conversation_id,
        organization_id=payload.organization_id,
        user_id=payload.user_id,
        parent_task_id=payload.parent_task_id or "",
        executor_type=payload.executor_type,
        executor_id=executor_id,
    )
    run = task_state.create_run(
        task["id"],
        metadata={
            "trigger": "user",
            "agent_id": selected_agent_id,
            "workspace": payload.workspace,
            "organization_id": payload.organization_id,
            "user_id": payload.user_id,
            "executor_type": payload.executor_type,
            "executor_id": executor_id,
        },
    )
    if payload.executor_type == "team":
        if expert_selection:
            emit(
                task["id"],
                "expert_selection",
                "已选择参与专家",
                (
                    f"已自动匹配“{expert_selection['team_name']}”，"
                    if expert_selection["selection_mode"] == "automatic"
                    else f"使用已指定的“{expert_selection['team_name']}”，"
                )
                + f"由 {len(expert_selection['members'])} 位专家并行分析，再由主管汇总。",
                expert_selection,
            )
        try:
            team_run = expert_team_service.create_run_for_task(
                executor_id, task["id"], scope, parent_run_id=run["id"]
            )
        except (ExpertNotFoundError, ExpertConflictError, ExpertValidationError) as exc:
            raise _expert_http_error(exc) from exc
        _schedule_team_run(team_run["id"])
        return {
            **_public_task(task),
            "result": {},
            "artifacts": [],
            "run": _public_runtime_record(run),
            "team_run": _public_runtime_record(team_run),
            "expert_selection": expert_selection,
        }
    _schedule_runtime(task["id"], run["id"])
    return {
        **_public_task(task),
        "result": {},
        "artifacts": [],
        "run": _public_runtime_record(run),
    }


@app.get("/api/tasks/{task_id}")
def get_task(task_id: str) -> dict[str, Any]:
    task = db.query_one("SELECT * FROM tasks WHERE id = ?", (task_id,))
    if not task:
        raise HTTPException(status_code=404, detail="Task not found")
    events = db.query_all("SELECT * FROM task_events WHERE task_id = ? AND type != 'answer_delta' ORDER BY id", (task_id,))
    return {
        **_public_task(task),
        "events": [_public_event(e) for e in events if _is_public_task_event(e)],
        "runtime": _runtime_projection(task_id),
    }


@app.get("/api/conversations/{conversation_id}/messages")
def get_conversation_messages(conversation_id: str) -> dict[str, Any]:
    rows = db.query_all(
        "SELECT * FROM tasks WHERE conversation_id = ? ORDER BY created_at, id LIMIT 100",
        (conversation_id,),
    )
    messages: list[dict[str, Any]] = []
    for task in rows:
        messages.append({"role": "user", "content": task["message"], "task_id": task["id"]})
        answer = db.query_one(
            "SELECT id, content FROM task_events WHERE task_id = ? AND type IN ('answer', 'error') ORDER BY id DESC LIMIT 1",
            (task["id"],),
        )
        if answer:
            messages.append({"role": "assistant", "content": answer["content"], "task_id": task["id"], "event_id": answer["id"]})
    return {"conversation_id": conversation_id, "messages": messages}


@app.post("/api/uploads")
async def upload_file(file: UploadFile = File(...)) -> dict[str, Any]:
    max_bytes = int(os.getenv("APP_MAX_UPLOAD_MB", "20")) * 1024 * 1024
    raw = await file.read(max_bytes + 1)
    if len(raw) > max_bytes:
        raise HTTPException(status_code=413, detail="文件超过上传大小限制")
    original = Path(file.filename or "upload.bin").name
    safe_name = re.sub(r"[^A-Za-z0-9._\-\u4e00-\u9fff]", "_", original)[:180] or "upload.bin"
    upload_id = "upl_" + uuid.uuid4().hex[:12]
    UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
    path = UPLOAD_DIR / f"{upload_id}_{safe_name}"
    path.write_bytes(raw)
    record = {"id": upload_id, "name": original, "content_type": file.content_type or "application/octet-stream", "size": len(raw), "path": str(path), "created_at": db.utc_now()}
    db.execute("INSERT INTO uploads(id, name, content_type, size, path, created_at) VALUES (?, ?, ?, ?, ?, ?)", tuple(record.values()))
    return {k: v for k, v in record.items() if k != "path"}


@app.get("/api/tasks/{task_id}/events")
def get_task_events(task_id: str, after_id: int = 0) -> list[dict[str, Any]]:
    _task_or_404(task_id)
    if after_id < 0:
        raise HTTPException(status_code=400, detail="事件游标必须是非负整数")
    events = db.query_all("SELECT * FROM task_events WHERE task_id = ? AND id > ? ORDER BY id", (task_id, after_id))
    return [_public_event(e) for e in events if _is_public_task_event(e)]


@app.get("/api/tasks/{task_id}/events/stream")
async def stream_task_events(
    task_id: str,
    request: Request,
    cursor: int | None = None,
    after_id: int | None = None,
):
    _task_or_404(task_id)
    initial_cursor = _task_event_cursor(request, cursor=cursor, after_id=after_id)

    async def event_generator():
        last_id = initial_cursor
        idle_rounds = 0
        while True:
            if await request.is_disconnected():
                break
            events = db.query_all("SELECT * FROM task_events WHERE task_id = ? AND id > ? ORDER BY id", (task_id, last_id))
            for event in events:
                last_id = int(event["id"])
                if not _is_public_task_event(event):
                    continue
                payload = _public_event(event)
                yield f"id: {last_id}\nevent: task_event\ndata: {json.dumps(payload, ensure_ascii=False)}\n\n"
            task = db.query_one("SELECT status FROM tasks WHERE id = ?", (task_id,))
            status = str((task or {}).get("status") or "")
            should_close = status in _TASK_STREAM_TERMINAL_STATUSES or status == "waiting_approval"
            if should_close and not events:
                idle_rounds += 1
                if idle_rounds > 1:
                    payload = {
                        "task_id": task_id,
                        "status": status,
                        "terminal": status in _TASK_STREAM_TERMINAL_STATUSES,
                        "cursor": last_id,
                    }
                    yield f"event: task_status\ndata: {json.dumps(payload, ensure_ascii=False)}\n\n"
                    break
            else:
                idle_rounds = 0
            await asyncio.sleep(0.1)
    return StreamingResponse(
        event_generator(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache, no-transform",
            "X-Accel-Buffering": "no",
        },
    )


@app.post("/api/tasks/{task_id}/approve")
async def approve_task(task_id: str, payload: ApprovalRequest) -> dict[str, Any]:
    task = _task_or_404(task_id)
    if task.get("status") != "waiting_approval":
        raise HTTPException(status_code=409, detail="当前任务不处于等待审批状态")
    result = db.json_loads(task.get("result_json"), {})
    if result.get("pending_action") == "policy_approval":
        active = _active_run_or_409(task_id)
        command = task_state.enqueue_command(
            task_id,
            "approval",
            run_id=active["id"],
            payload={"approved": payload.approved, "note": payload.note},
            priority=90,
            deduplicate=True,
        )
        return {"ok": True, "command": command}
    background = asyncio.create_task(runtime.resume_after_approval(task_id, payload.approved, payload.note))
    _runtime_tasks.add(background)
    background.add_done_callback(_runtime_tasks.discard)
    return {"ok": True}


def _artifact_row_to_public(row: dict[str, Any]) -> dict[str, Any]:
    metadata = db.json_loads(row.get("metadata_json"), {})
    return _public_artifact({**row, "metadata": metadata})


def _artifact_path(row: dict[str, Any]) -> Path:
    relative = str(row.get("relative_path") or "")
    try:
        if relative:
            return resolve_artifact_path(relative)
        legacy = Path(str(row.get("path") or "")).resolve(strict=True)
        root = ARTIFACT_DIR.resolve(strict=True)
        derived = legacy.relative_to(root).as_posix()
        return resolve_artifact_path(derived)
    except (FileNotFoundError, OSError, RuntimeError, ToolError, ValueError) as exc:
        raise HTTPException(status_code=404, detail="产物文件不存在或不在受控目录中") from exc


def _artifact_or_404(artifact_id: str) -> dict[str, Any]:
    artifact = db.query_one("SELECT * FROM artifacts WHERE id = ?", (artifact_id,))
    if not artifact:
        raise HTTPException(status_code=404, detail="Artifact not found")
    return artifact


@app.get("/api/artifacts")
def list_artifacts(
    task_id: str = "",
    run_id: str = "",
    workspace_id: str = "",
    kind: str = "",
    limit: int = 200,
) -> list[dict[str, Any]]:
    limit = max(1, min(int(limit), 1000))
    clauses: list[str] = []
    params: list[Any] = []
    for column, value in (
        ("task_id", task_id),
        ("run_id", run_id),
        ("workspace_id", workspace_id),
        ("kind", kind),
    ):
        if value:
            clauses.append(f"{column} = ?")
            params.append(value)
    where = " WHERE " + " AND ".join(clauses) if clauses else ""
    rows = db.query_all(
        f"SELECT * FROM artifacts{where} ORDER BY created_at DESC LIMIT ?",  # noqa: S608 - fixed columns
        (*params, limit),
    )
    return [_artifact_row_to_public(row) for row in rows]


@app.get("/api/tasks/{task_id}/artifacts")
def list_task_artifacts(task_id: str, run_id: str = "") -> list[dict[str, Any]]:
    _task_or_404(task_id)
    return list_artifacts(task_id=task_id, run_id=run_id)


@app.get("/api/artifacts/{artifact_id}")
def get_artifact(artifact_id: str) -> dict[str, Any]:
    return _artifact_row_to_public(_artifact_or_404(artifact_id))


def _sanitise_html_preview(value: str) -> str:
    cleaned = re.sub(
        r"<\s*(script|iframe|object|embed|base|form|link|meta)[^>]*>[\s\S]*?<\s*/\s*\1\s*>",
        "",
        value,
        flags=re.IGNORECASE,
    )
    cleaned = re.sub(
        r"<\s*(script|iframe|object|embed|base|form|link|meta)\b[^>]*?/?>",
        "",
        cleaned,
        flags=re.IGNORECASE,
    )
    cleaned = re.sub(r"\s+on[a-z]+\s*=\s*(['\"]).*?\1", "", cleaned, flags=re.IGNORECASE | re.DOTALL)
    cleaned = re.sub(
        r"\s+(href|src)\s*=\s*(['\"])(?!#|data:)[\s\S]*?\2",
        "",
        cleaned,
        flags=re.IGNORECASE,
    )
    csp = "<meta http-equiv=\"Content-Security-Policy\" content=\"default-src 'none'; img-src data:; style-src 'unsafe-inline'; font-src data:\">"
    return csp + cleaned


@app.get("/api/artifacts/{artifact_id}/preview")
def preview_artifact(artifact_id: str) -> dict[str, Any]:
    artifact = _artifact_or_404(artifact_id)
    path = _artifact_path(artifact)
    kind = str(artifact.get("kind") or path.suffix.lstrip(".")).lower()
    public = _artifact_row_to_public(artifact)
    if path.stat().st_size > 25 * 1024 * 1024:
        return {"artifact": public, "preview_kind": "unavailable", "message": "文件超过 25MB，请下载后查看。"}
    try:
        if kind in {"markdown", "md"}:
            return {"artifact": public, "preview_kind": "markdown", "content": path.read_text(encoding="utf-8", errors="replace")[:500_000]}
        if kind == "html":
            raw = path.read_text(encoding="utf-8", errors="replace")[:500_000]
            return {"artifact": public, "preview_kind": "html", "content": _sanitise_html_preview(raw), "sandbox": ""}
        if kind in {"csv", "tsv"}:
            delimiter = "\t" if kind == "tsv" else ","
            with path.open("r", encoding="utf-8-sig", errors="replace", newline="") as stream:
                rows = [row[:30] for _, row in zip(range(101), csv.reader(stream, delimiter=delimiter))]
            return {"artifact": public, "preview_kind": "spreadsheet", "sheets": [{"name": path.stem, "rows": rows}]}
        if kind == "xlsx":
            from openpyxl import load_workbook

            workbook = load_workbook(path, read_only=True, data_only=True)
            sheets = []
            for sheet in workbook.worksheets[:10]:
                rows = [
                    ["" if cell is None else str(cell) for cell in row[:30]]
                    for _, row in zip(range(101), sheet.iter_rows(values_only=True))
                ]
                sheets.append({"name": sheet.title, "rows": rows})
            workbook.close()
            return {"artifact": public, "preview_kind": "spreadsheet", "sheets": sheets}
        if kind == "docx":
            from docx import Document

            document = Document(path)
            paragraphs = [item.text for item in document.paragraphs if item.text.strip()][:500]
            tables = [
                [[cell.text for cell in row.cells] for row in table.rows[:100]]
                for table in document.tables[:20]
            ]
            return {"artifact": public, "preview_kind": "document", "paragraphs": paragraphs, "tables": tables}
        if kind == "pptx":
            from pptx import Presentation

            presentation = Presentation(path)
            slides = []
            for index, slide in enumerate(presentation.slides, start=1):
                if index > 100:
                    break
                texts = [
                    str(shape.text).strip()
                    for shape in slide.shapes
                    if hasattr(shape, "text") and str(shape.text).strip()
                ]
                slides.append({"number": index, "title": texts[0] if texts else f"第 {index} 页", "texts": texts})
            return {"artifact": public, "preview_kind": "slides", "slides": slides}
        if kind == "pdf":
            return {"artifact": public, "preview_kind": "pdf", "url": public["download_url"] + "?inline=true"}
        if kind in {"txt", "text", "json", "yaml", "yml"}:
            return {"artifact": public, "preview_kind": "text", "content": path.read_text(encoding="utf-8", errors="replace")[:500_000]}
    except Exception as exc:
        return {"artifact": public, "preview_kind": "error", "message": f"预览生成失败：{exc}"}
    return {"artifact": public, "preview_kind": "unavailable", "message": "该格式暂不支持平台内预览，可下载原文件。"}


@app.get("/api/artifacts/{artifact_id}/download")
def download_artifact(artifact_id: str, inline: bool = False):
    artifact = _artifact_or_404(artifact_id)
    path = _artifact_path(artifact)
    media_type = str(artifact.get("mime_type") or "application/octet-stream")
    if inline:
        return FileResponse(path, media_type=media_type)
    return FileResponse(path, filename=artifact["name"], media_type=media_type)


app.mount("/", StaticFiles(directory=str(WEB_DIR), html=True), name="web")
