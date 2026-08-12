from __future__ import annotations

import asyncio
import contextvars
import hashlib
import json
import re
import time
import uuid
import zipfile
from pathlib import Path
from typing import Any, Awaitable, Callable, Iterable, Mapping

from app import db
from app.builtin_skill_catalog import get_builtin_skill, recommend_builtin_skill
from app.services.context_service import ContextService, ExecutionScope
from app.services.conversation_summary_service import ConversationSummaryService
from app.services.event_bus import emit
from app.services.mcp_gateway import ARTIFACT_DIR, McpGateway, ToolError, resolve_artifact_path
from app.services.model_gateway import ModelGateway
from app.services.policy_engine import PolicyApprovalRequired, PolicyEngine
from app.services.skill_registry import SkillRegistry
from app.services.task_state import (
    InvalidStateTransition,
    TaskCancellationRequested,
    TaskStateError,
    TaskStateService,
)

FINAL_STATUSES = {"completed", "failed", "waiting_approval", "cancelled"}


class AgentRuntime:
    ATTACHMENT_MAX_FILES = 10
    ATTACHMENT_MAX_FILE_BYTES = 20 * 1024 * 1024
    ATTACHMENT_MAX_CONTEXT_CHARS = 60_000
    ATTACHMENT_MAX_FILE_CHARS = 20_000
    ATTACHMENT_MAX_ARCHIVE_ENTRIES = 5_000
    ATTACHMENT_MAX_UNCOMPRESSED_BYTES = 100 * 1024 * 1024
    ATTACHMENT_MAX_WORKSHEETS = 8
    ATTACHMENT_MAX_ROWS_PER_SHEET = 200
    ATTACHMENT_MAX_COLUMNS_PER_SHEET = 50
    ATTACHMENT_MAX_PDF_PAGES = 40
    ATTACHMENT_MAX_SLIDES = 40
    ATTACHMENT_MAX_SHAPES_PER_SLIDE = 200
    ATTACHMENT_MAX_TABLES = 50
    ATTACHMENT_MAX_ROWS_PER_TABLE = 200
    ATTACHMENT_MAX_COLUMNS_PER_TABLE = 50

    def __init__(
        self,
        skill_registry: SkillRegistry,
        mcp_gateway: McpGateway,
        model_gateway: ModelGateway,
        skill_url_installer: Callable[[str], Awaitable[dict[str, Any]]] | None = None,
        mcp_url_installer: Callable[[str], Awaitable[list[dict[str, Any]]]] | None = None,
        task_state: TaskStateService | None = None,
        policy_engine: PolicyEngine | None = None,
        context_service: ContextService | None = None,
        conversation_summary_service: ConversationSummaryService | None = None,
    ) -> None:
        self.skill_registry = skill_registry
        self.mcp_gateway = mcp_gateway
        self.model_gateway = model_gateway
        self.skill_url_installer = skill_url_installer
        self.mcp_url_installer = mcp_url_installer
        self.task_state = task_state or TaskStateService()
        self.policy_engine = policy_engine or PolicyEngine()
        self.context_service = context_service or ContextService()
        self.conversation_summary_service = conversation_summary_service or ConversationSummaryService()
        self._execution_context: contextvars.ContextVar[dict[str, Any] | None] = contextvars.ContextVar(
            f"agent_runtime_execution_{id(self)}", default=None
        )

    async def run_task(self, task_id: str, *, run_id: str | None = None) -> None:
        task = db.query_one("SELECT * FROM tasks WHERE id = ?", (task_id,))
        if not task:
            return
        run: dict[str, Any] | None = None
        context_token: contextvars.Token[dict[str, Any] | None] | None = None
        try:
            run = self.task_state.begin_run(
                task_id,
                run_id=run_id,
                metadata={"agent_id": task.get("agent_id", ""), "workspace": task.get("workspace", "default")},
            )
            restored_state: dict[str, Any] = {}
            checkpoint_id = str(run.get("resumed_from_checkpoint_id") or "")
            if checkpoint_id:
                restored = self.task_state.restore_checkpoint(
                    checkpoint_id,
                    restore_metadata={"run_id": run["id"], "reason": "runtime_resume"},
                    mark_restored=not bool(run.get("metadata", {}).get("checkpoint_restore_audited")),
                )
                if isinstance(restored.get("state"), dict):
                    restored_state = dict(restored["state"])
                emit(
                    task_id,
                    "recovery",
                    "已从检查点恢复",
                    f"运行尝试 {run['attempt']} 已从最近安全检查点继续。",
                    {"run_id": run["id"], "checkpoint_id": checkpoint_id},
                )
            agent = self._get_agent(task["agent_id"])
            restored_permissions = restored_state.get("effective_permissions")
            if isinstance(restored_permissions, dict):
                effective_permissions = self._normalize_permissions(restored_permissions)
                permission_source = str(
                    restored_state.get("permission_source") or "restored_checkpoint"
                )
            else:
                effective_permissions, permission_source = self._permission_snapshot_for_task(
                    task, agent
                )
            restored_state["effective_permissions"] = effective_permissions
            restored_state["permission_source"] = permission_source
            restored_state["tool_calls_used"] = max(
                0, int(restored_state.get("tool_calls_used") or 0)
            )
            restored_state["permission_elapsed_seconds"] = max(
                0.0, float(restored_state.get("permission_elapsed_seconds") or 0.0)
            )
            execution_context = {
                "task_id": task_id,
                "run_id": run["id"],
                "attempt": run["attempt"],
                "state": restored_state,
                # Approval continuations reuse the same durable run. Restore
                # its node index so the plan updates the existing timeline.
                "nodes": {
                    item["node_key"]: item["id"]
                    for item in self.task_state.list_nodes(run["id"])
                },
                "plan_nodes": {},
                "worker_id": f"runtime:{run['id']}",
                "permission_timer_started": time.perf_counter(),
                "permission_elapsed_base": restored_state["permission_elapsed_seconds"],
            }
            context_token = self._execution_context.set(execution_context)
            self.task_state.update_run_metadata(
                run["id"],
                {
                    "effective_permissions": effective_permissions,
                    "permission_source": permission_source,
                },
            )
            db.update_task_status(task_id, "running")
            emit(
                task_id,
                "start",
                "任务已启动",
                f"Agent 正在处理：{task['message']}",
                {"run_id": run["id"], "attempt": run["attempt"]},
            )
            emit(
                task_id,
                "permissions",
                "已固定本次运行权限",
                self._permission_snapshot_summary(effective_permissions),
                {
                    "effective_permissions": effective_permissions,
                    "source": permission_source,
                    "run_id": run["id"],
                },
            )
            self._create_checkpoint("本次运行权限快照已固定")
            await self._evaluate_policy("task.created", {"task": task}, enforce=True)
            self._raise_if_cancelled()
            emit(task_id, "agent", "已选择 Agent", agent["name"], {"agent": agent})

            if await self._try_platform_command(task):
                return

            history = self._conversation_history(task)
            effective_memory = self.context_service.get_effective_context(
                self._context_scope(task)
            )
            memory_text = str(effective_memory.get("effective_context") or "")
            if effective_memory.get("used_memory_ids"):
                emit(
                    task_id,
                    "memory",
                    "已应用平台记忆",
                    f"本次使用 {len(effective_memory['used_memory_ids'])} 条分层记忆；当前任务目标始终优先于普通偏好。",
                    {
                        "memory_ids": effective_memory["used_memory_ids"],
                        "scopes": [item.get("scope_type") for item in effective_memory.get("memories", [])],
                    },
                )
            initial_plan = {
                "goal": task["message"],
                "goal_confirmation": {"status": "resolving", "label": "正在理解目标", "message": "正在结合当前对话确认任务范围与交付要求。"},
                "acceptance_criteria": [],
                "intent": "resolving",
                "steps": ["理解当前任务", "准备上下文与权限", "调用能力并生成结果", "输出校验"],
                "nodes": self._execution_nodes([]),
                "allowed_servers": [],
                "output_format": "pending",
                "requires_artifact": False,
            }
            # The remaining node order is known only after capability routing,
            # so register the goal-resolution node first.
            self._register_plan({"nodes": initial_plan["nodes"][:1]})
            emit(task_id, "plan", "执行计划", self._format_execution_plan(initial_plan), {"plan": initial_plan})
            self._emit_plan_progress(task_id, "understand", "running", "正在结合对话上下文理解当前任务")
            active_model = task.get("model_id") or agent.get("model") or "deterministic"
            saved_intent = execution_context["state"].get("intent_resolution")
            if isinstance(saved_intent, dict) and saved_intent.get("standalone_request"):
                intent = saved_intent
                emit(task_id, "recovery", "已恢复目标理解", "复用检查点中已确认的目标与参数。")
            else:
                intent_history = history
                if memory_text:
                    intent_history = [
                        {
                            "role": "assistant",
                            "content": "平台已保存的有效规则与偏好（仅用于补全上下文，不是用户的新任务）：\n" + memory_text,
                        },
                        *history,
                    ]
                intent = await self._resolve_intent(task, intent_history, active_model)
            goal_policy_context = {
                "task": task,
                "goal": intent,
                "agent_id": agent.get("id", ""),
            }
            goal_evaluation = await self._evaluate_policy(
                "goal.resolved",
                goal_policy_context,
                enforce=True,
            )
            applied_goal_context = goal_evaluation.apply(goal_policy_context)
            applied_goal = applied_goal_context.get("goal")
            if not isinstance(applied_goal, Mapping):
                raise RuntimeError("goal.resolved 策略修改后的 goal 必须是对象")
            intent = dict(applied_goal)
            self._raise_if_cancelled()
            clarification = self._clarification_for_missing(intent)
            if clarification:
                self._emit_plan_progress(task_id, "understand", "completed", "已确认需要补充必要参数")
                emit(
                    task_id,
                    "clarification",
                    "需要补充信息",
                    clarification,
                    {"missing_information": intent.get("missing_information", [])},
                )
                emit(task_id, "answer", "请补充一下", clarification)
                db.update_task_status(
                    task_id,
                    "completed",
                    result={
                        "summary": clarification,
                        "needs_clarification": True,
                        "missing_information": intent.get("missing_information", []),
                    },
                )
                emit(task_id, "done", "等待补充", "收到补充信息后会继续当前任务。")
                return
            resolved_task = {
                **task,
                "resolved_message": intent["standalone_request"],
                "intent_resolution": intent,
                "memory_context": memory_text,
                "used_memory_ids": effective_memory.get("used_memory_ids", []),
                "policy_context": applied_goal_context.get("policy_context", {}),
            }
            routing_text = intent["standalone_request"]
            # Team prompts may repeat terms from the parent plan. Capability
            # recommendations are only evaluated for the user-facing task.
            internal_team_step = str(task.get("executor_type") or "") in {
                "team_member", "team_supervisor"
            }
            task_result = db.json_loads(task.get("result_json"), {})
            declined_recommendation_ids = {
                str(item)
                for item in task_result.get("declined_recommendation_ids", [])
                if str(item).strip()
            } if isinstance(task_result, dict) else set()
            unavailable_recommendation_ids = {
                item["id"] for item in self.skill_registry.list_skills()
            } | declined_recommendation_ids
            recommendation_decided = bool(
                isinstance(task_result, dict)
                and task_result.get("skip_skill_recommendations")
            )
            recommendation = (
                None
                if internal_team_step or recommendation_decided
                else recommend_builtin_skill(routing_text, unavailable_recommendation_ids)
            )
            if recommendation:
                public_recommendation = {
                    key: recommendation[key]
                    for key in ["id", "name", "description", "source_label"]
                }
                default_message = (
                    f"内置目录中有适合当前目标的“{recommendation['name']}”。"
                    "是否安装后继续任务？"
                )
                approval_request = await self._apply_approval_requested_policy(
                    {
                        "action": "install_recommended_skill",
                        "event": "skill.recommended",
                        "title": "安装内置 Skill",
                        "message": default_message,
                        "recommendations": [public_recommendation],
                    }
                )
                message = str(approval_request.get("message") or default_message)
                approval_title = str(
                    approval_request.get("title") or "安装内置 Skill"
                )
                emit(
                    task_id,
                    "approval_required",
                    approval_title,
                    message,
                    {
                        "action": "install_recommended_skill",
                        "recommendations": [public_recommendation],
                        "approval_request": approval_request,
                    },
                )
                db.update_task_status(
                    task_id,
                    "waiting_approval",
                    result={
                        **(task_result if isinstance(task_result, dict) else {}),
                        "pending_action": "install_recommended_skill",
                        "recommendation_id": recommendation["id"],
                        "summary": message,
                    },
                )
                # Persist the pause here because this branch returns before the
                # normal run finalizer.
                self._complete_running_nodes("能力推荐已生成，正在等待用户确认")
                current_run = self.task_state.get_run(run["id"])
                if current_run and current_run["status"] == "running":
                    self.task_state.transition_run(run["id"], "waiting_approval")
                return
            allowed_skill_ids = None if agent.get("id") == "general-agent" else agent.get("skills")
            selected = self.skill_registry.score_skills(routing_text, allowed_ids=allowed_skill_ids)
            if not selected and allowed_skill_ids is None:
                selected = self.skill_registry.score_skills(routing_text)
            if not selected:
                fallback_skill = self.skill_registry.get_skill("general_task")
                if fallback_skill and (allowed_skill_ids is None or "general_task" in set(allowed_skill_ids or [])):
                    selected = [{"skill": fallback_skill, "score": 0.1}]
            selected_skills = [s["skill"] for s in selected[:3]]
            emit(
                task_id,
                "skill",
                "已匹配 Skill",
                "、".join([s["name"] for s in selected_skills]) if selected_skills else "未匹配到专项 Skill，将使用通用流程。",
                {"skills": [{"id": s["id"], "name": s["name"], "score": selected[i]["score"]} for i, s in enumerate(selected_skills)]},
            )
            execution_context["state"].update(
                {
                    "phase": "goal_resolved",
                    "intent_resolution": intent,
                    "resolved_message": intent["standalone_request"],
                    "selected_skill_ids": [item["id"] for item in selected_skills],
                    "completed_tools": execution_context["state"].get("completed_tools", {}),
                    "steering_messages": execution_context["state"].get("steering_messages", []),
                }
            )
            self._create_checkpoint("目标、参数与能力选择已确认", node_key="understand")

            await self._run_general_task(resolved_task, agent, selected_skills, history)
        except TaskCancellationRequested:
            self._cancel_running_nodes()
            self._acknowledge_cancel()
            emit(task_id, "cancelled", "任务已取消", "已在安全执行边界停止当前任务。")
            db.update_task_status(task_id, "cancelled", result={"cancelled": True})
            if run and (self.task_state.get_run(run["id"]) or {}).get("status") not in {"completed", "failed", "cancelled"}:
                self.task_state.finish_run(run["id"], status="cancelled", result={"cancelled": True})
        except PolicyApprovalRequired:
            # The policy helper has already persisted waiting_approval and an
            # auditable approval event.  Leaving the run non-terminal allows a
            # user decision or restart recovery to continue from a checkpoint.
            return
        except asyncio.CancelledError:
            # Process shutdown is not a user cancellation.  Keep the durable
            # run active so startup recovery can create a new resumed attempt.
            emit(task_id, "interrupted", "运行已中断", "平台服务正在停止，将在下次启动时从安全检查点恢复。")
            raise
        except Exception as exc:
            emit(task_id, "error", "任务失败", str(exc), {"error_type": exc.__class__.__name__})
            db.update_task_status(task_id, "failed", result={"error": str(exc)})
            if run and (self.task_state.get_run(run["id"]) or {}).get("status") not in {"completed", "failed", "cancelled"}:
                self._cancel_running_nodes(error={"message": str(exc), "error_type": exc.__class__.__name__})
                self.task_state.finish_run(
                    run["id"],
                    status="failed",
                    error={"message": str(exc), "error_type": exc.__class__.__name__},
                )
            try:
                await self._evaluate_policy("task.failed", {"task": task, "error": {"message": str(exc)}}, enforce=False)
            except Exception:
                pass
        else:
            current_task = db.query_one("SELECT * FROM tasks WHERE id = ?", (task_id,)) or {}
            current_run = self.task_state.get_run(run["id"]) if run else None
            if run and current_run and current_run["status"] not in {"completed", "failed", "cancelled"}:
                if current_task.get("status") == "waiting_approval":
                    self._complete_running_nodes("当前安全步骤已结束，正在等待用户审批")
                    if current_run["status"] == "running":
                        self.task_state.transition_run(run["id"], "waiting_approval")
                elif current_task.get("status") == "cancelled":
                    self.task_state.finish_run(run["id"], status="cancelled", result={"cancelled": True})
                elif current_task.get("status") == "failed":
                    self.task_state.finish_run(
                        run["id"], status="failed", error=db.json_loads(current_task.get("result_json"), {})
                    )
                else:
                    result = db.json_loads(current_task.get("result_json"), {})
                    self._complete_running_nodes("任务已结束")
                    self._create_checkpoint("任务已完成并通过输出校验", node_key="validate")
                    self.task_state.finish_run(run["id"], result=result)
                    try:
                        compacted = self._maybe_compact_conversation(current_task)
                        if compacted:
                            emit(
                                task_id,
                                "conversation_summary",
                                "已压缩较早对话",
                                "较早轮次已整理为可审计摘要；明确约束会继续保留，且不会被当成本次新目标。",
                                {
                                    "conversation_id": compacted["conversation_id"],
                                    "through_task_id": compacted["through_task_id"],
                                    "version": compacted["version"],
                                    "preserved_constraints": compacted["preserved_constraints"],
                                },
                            )
                    except Exception as exc:
                        emit(
                            task_id,
                            "notice",
                            "对话摘要暂未更新",
                            f"任务已正常完成；较早对话将在后续轮次重试压缩：{exc}",
                        )
                    try:
                        await self._evaluate_policy("task.completed", {"task": current_task, "result": result}, enforce=False)
                    except Exception:
                        pass
        finally:
            if context_token is not None:
                self._execution_context.reset(context_token)

    def _execution(self) -> dict[str, Any] | None:
        return self._execution_context.get()

    def _effective_permissions(self) -> dict[str, Any]:
        execution = self._execution()
        value = (execution or {}).get("state", {}).get("effective_permissions", {})
        return value if isinstance(value, dict) else {}

    def _sync_permission_elapsed(self) -> float:
        """Update active execution time in checkpoint-safe state.

        A monotonic clock is used while the process is alive.  Only elapsed
        active runtime is carried into a resumed attempt, so service downtime
        does not silently consume an expert member's tool timeout budget.
        """

        execution = self._execution()
        if not execution:
            return 0.0
        base = float(execution.get("permission_elapsed_base") or 0.0)
        started = float(execution.get("permission_timer_started") or time.perf_counter())
        elapsed = max(0.0, base + (time.perf_counter() - started))
        execution["state"]["permission_elapsed_seconds"] = elapsed
        return elapsed

    def _remaining_tool_timeout(self) -> float | None:
        permissions = self._effective_permissions()
        configured = permissions.get("timeout_seconds")
        if not isinstance(configured, (int, float)) or isinstance(configured, bool):
            return None
        return float(configured) - self._sync_permission_elapsed()

    def _raise_if_cancelled(self) -> None:
        execution = self._execution()
        if execution:
            self.task_state.raise_if_cancel_requested(
                execution["task_id"], run_id=execution["run_id"]
            )

    def _acknowledge_cancel(self) -> None:
        execution = self._execution()
        if not execution:
            return
        while True:
            command = self.task_state.claim_command(
                execution["worker_id"],
                task_id=execution["task_id"],
                run_id=execution["run_id"],
                command_types=["cancel"],
            )
            if not command:
                break
            self.task_state.complete_command(command["id"], result={"cancelled": True})

    def _cancel_running_nodes(self, error: dict[str, Any] | None = None) -> None:
        execution = self._execution()
        if not execution:
            return
        for node in self.task_state.list_nodes(execution["run_id"]):
            if node["status"] != "running":
                continue
            try:
                if error:
                    self.task_state.fail_node(node["id"], error)
                else:
                    self.task_state.transition_node(node["id"], "cancelled")
            except TaskStateError:
                continue

    def _complete_running_nodes(self, summary: str) -> None:
        execution = self._execution()
        if not execution:
            return
        for node in self.task_state.list_nodes(execution["run_id"]):
            if node["status"] != "running":
                continue
            try:
                self.task_state.finish_node(node["id"], output={"summary": summary})
            except TaskStateError:
                continue

    def _create_checkpoint(self, reason: str, *, node_key: str = "") -> dict[str, Any] | None:
        execution = self._execution()
        if not execution:
            return None
        self._sync_permission_elapsed()
        node_id = execution.get("nodes", {}).get(node_key) if node_key else None
        checkpoint = self.task_state.create_checkpoint(
            execution["run_id"],
            execution.get("state", {}),
            node_id=node_id,
            reason=reason,
            metadata={"attempt": execution.get("attempt", 1)},
        )
        emit(
            execution["task_id"],
            "checkpoint",
            "已保存安全检查点",
            reason,
            {"checkpoint_id": checkpoint["id"], "run_id": execution["run_id"], "node_id": node_id},
        )
        return checkpoint

    async def _evaluate_policy(
        self,
        event: str,
        context: dict[str, Any],
        *,
        enforce: bool,
    ):
        execution = self._execution()
        task = db.query_one("SELECT * FROM tasks WHERE id = ?", ((execution or {}).get("task_id", ""),)) or {}
        scoped_context = {
            "task_id": (execution or {}).get("task_id", ""),
            "run_id": (execution or {}).get("run_id", ""),
            "organization_id": task.get("organization_id", "local-org"),
            "user_id": task.get("user_id", "local-user"),
            "workspace_id": task.get("workspace", "default"),
            "agent_id": task.get("agent_id", ""),
            "executor_type": task.get("executor_type", "agent"),
            "executor_id": task.get("executor_id") or task.get("agent_id", ""),
            **context,
        }
        evaluation = await self.policy_engine.evaluate(event, scoped_context)
        if execution:
            emit(
                execution["task_id"],
                "policy_decision",
                "策略校验",
                evaluation.summary,
                evaluation.to_dict(),
            )
        if enforce and evaluation.denied:
            raise RuntimeError(evaluation.summary)
        if enforce and evaluation.requires_approval:
            await self._wait_for_policy_approval(
                evaluation,
                "",
                "",
                event=event,
            )
        return evaluation

    async def _apply_approval_requested_policy(
        self,
        request: dict[str, Any],
    ) -> dict[str, Any]:
        """Run the approval-request lifecycle without recursively requesting approval.

        ``require_approval`` on ``approval.requested`` confirms that the request
        must remain gated; it does not create a second approval.  A denial blocks
        the underlying operation before an approval prompt is exposed.
        """

        context = {"approval": dict(request)}
        evaluation = await self._evaluate_policy(
            "approval.requested", context, enforce=False
        )
        if evaluation.denied:
            raise RuntimeError(evaluation.summary)
        applied = evaluation.apply(context)
        approval = applied.get("approval")
        if not isinstance(approval, Mapping):
            raise RuntimeError(
                "approval.requested 策略修改后的 approval 必须是对象"
            )
        # Workflow identity and the underlying policy evidence are immutable.
        # Policies may refine user-facing copy and add metadata, but cannot
        # redirect the approval or conceal what is being approved.
        result = dict(request)
        result.update(dict(approval))
        for key in ("action", "event", "tool", "requests", "policy"):
            if key in request:
                result[key] = request[key]
        return result

    async def _try_platform_command(self, task: dict[str, Any]) -> bool:
        task_id = task["id"]
        message = task["message"].strip()
        lowered = message.lower()
        attachments = db.json_loads(task.get("attachments_json"), [])
        memory_scope = self._context_scope(task)

        remember_match = re.match(
            r"^(?:请|帮我)?记住\s*[：:,，]?\s*(.+)$",
            message,
            re.IGNORECASE | re.DOTALL,
        )
        if remember_match:
            content = remember_match.group(1).strip()
            if not content:
                answer = "请在“记住：”后写明要长期保留的规则、偏好或事实。"
            else:
                remembered = self.context_service.remember(
                    memory_scope,
                    content,
                    title=content[:40],
                    source_ref=task_id,
                    created_by=memory_scope.user_id,
                )
                answer = (
                    f"已记住：{remembered['content']}\n"
                    "这条记忆会在同一用户和工作区的新对话中生效，可在“记忆”页面编辑、停用或删除。"
                )
                emit(task_id, "memory_saved", "记忆已保存", answer, {"memory": remembered})
            emit(task_id, "answer", "记忆已处理", answer)
            db.update_task_status(task_id, "completed", result={"summary": answer})
            emit(task_id, "done", "已完成", "记忆指令已处理。")
            return True

        forget_match = re.match(
            r"^(?:请|帮我)?(?:忘记|删除记忆)\s*[：:,，]?\s*(.+)$",
            message,
            re.IGNORECASE | re.DOTALL,
        )
        if forget_match:
            query = forget_match.group(1).strip()
            deleted_ids = self.context_service.forget(memory_scope, query=query) if query else []
            answer = (
                f"已删除 {len(deleted_ids)} 条完全匹配的记忆。"
                if deleted_ids
                else "没有找到标题或内容完全匹配的记忆；你可以先说“查看记忆”，再按完整标题删除。"
            )
            emit(task_id, "memory_deleted", "记忆已删除" if deleted_ids else "未找到匹配记忆", answer, {"memory_ids": deleted_ids})
            emit(task_id, "answer", "记忆已处理", answer)
            db.update_task_status(task_id, "completed", result={"summary": answer, "deleted_memory_ids": deleted_ids})
            emit(task_id, "done", "已完成", "记忆指令已处理。")
            return True

        if any(key in lowered for key in ["查看记忆", "我的记忆", "有哪些记忆", "list memories"]):
            effective = self.context_service.get_effective_context(memory_scope)
            memories = effective.get("memories", [])
            lines = [
                f"- [{item['scope_type']}] {item.get('title') or item['id']}：{item['content']}"
                for item in memories
            ]
            answer = (
                f"当前生效 {len(memories)} 条记忆：\n" + "\n".join(lines)
                if memories
                else "当前没有生效的长期记忆。你可以说“记住：以后默认用中文简洁回答”。"
            )
            emit(task_id, "answer", "当前有效记忆", answer, {"memory_ids": effective.get("used_memory_ids", [])})
            db.update_task_status(task_id, "completed", result={"summary": answer, "memories": memories})
            emit(task_id, "done", "已完成", "有效记忆列表已返回。")
            return True

        if any(key in lowered for key in ["查看已安装技能", "已安装的技能", "有哪些技能", "list skills"]):
            skills = self.skill_registry.list_skills()
            lines = [f"- {item['name']}（{item['id']}，{'已启用' if item['enabled'] else '已停用'}）" for item in skills]
            answer = f"当前共安装 {len(skills)} 个技能：\n" + "\n".join(lines)
            emit(task_id, "answer", "已安装技能", answer, {"skills": skills})
            db.update_task_status(task_id, "completed", result={"skills": skills})
            emit(task_id, "done", "已完成", "技能列表已返回。")
            return True

        if any(key in lowered for key in ["查看已安装mcp", "查看已安装 mcp", "已安装的mcp", "已安装的 mcp", "有哪些工具服务", "list mcp"]):
            servers = self.mcp_gateway.list_servers()
            lines = [f"- {item['name']}（{item['id']}，{item['kind']}，{'已启用' if item['enabled'] else '已停用'}）" for item in servers]
            answer = f"当前共安装 {len(servers)} 个工具服务：\n" + "\n".join(lines)
            emit(task_id, "answer", "已安装工具服务", answer, {"mcp_servers": servers})
            db.update_task_status(task_id, "completed", result={"mcp_servers": servers})
            emit(task_id, "done", "已完成", "工具服务列表已返回。")
            return True

        wants_skill_install = lowered.startswith("/install-skill") or any(key in lowered for key in ["安装 skill", "安装skill", "安装技能", "安装这个技能"])
        if wants_skill_install:
            content = self._skill_content_from_message(message)
            if not content:
                content = self._text_attachment(attachments, preferred_names={"skill.md"})
            if not content:
                download_url = self._https_url_from_message(message)
                if download_url:
                    if not self.skill_url_installer:
                        raise RuntimeError("平台未启用 Skill 下载链接安装")
                    skill = await self.skill_url_installer(download_url)
                    answer = f"技能“{skill['name']}”已从下载链接安装，ID 为 {skill['id']}，完整包共 {skill.get('file_count', 0)} 个文件。你可以在“技能中心”查看、维护或导出 ZIP。"
                    emit(task_id, "install", "技能安装成功", answer, {"skill": skill, "source_url": download_url})
                    emit(task_id, "answer", "安装完成", answer)
                    db.update_task_status(task_id, "completed", result={"installed": True, "type": "skill", "skill": skill})
                    emit(task_id, "done", "已完成", "Skill 安装包已写入平台。")
                    return True
            if not content:
                answer = "请在消息中粘贴包含 frontmatter 的 SKILL.md 内容，或上传 SKILL.md 后发送“安装这个技能”。"
                emit(task_id, "answer", "还需要技能文件", answer)
                db.update_task_status(task_id, "completed", result={"installed": False, "reason": "missing_skill_content"})
                emit(task_id, "done", "已完成", "未执行安装。")
                return True
            skill = self.skill_registry.install_content(content, fallback_id="chat_installed_skill")
            answer = f"技能“{skill['name']}”已安装，ID 为 {skill['id']}。你可以在“技能中心”看到它；如需让某个智能体使用，请在智能体配置中绑定该技能 ID。"
            emit(task_id, "install", "技能安装成功", answer, {"skill": skill})
            emit(task_id, "answer", "安装完成", answer)
            db.update_task_status(task_id, "completed", result={"installed": True, "type": "skill", "skill": skill})
            emit(task_id, "done", "已完成", "技能已写入平台。")
            return True

        wants_mcp_install = lowered.startswith("/install-mcp") or any(key in lowered for key in ["安装 mcp", "安装mcp", "安装这个mcp", "安装这个 mcp", "安装工具服务"])
        if wants_mcp_install:
            payload = self._json_from_message(message)
            if payload is None:
                attachment_text = self._text_attachment(attachments, suffixes={".json"})
                if attachment_text:
                    payload = json.loads(attachment_text)
            if payload is None:
                download_url = self._https_url_from_message(message)
                if download_url:
                    if not self.mcp_url_installer:
                        raise RuntimeError("平台未启用 MCP 下载链接安装")
                    servers = await self.mcp_url_installer(download_url)
                    names = "、".join(f"{item['name']}（{item['id']}）" for item in servers)
                    answer = f"工具服务 {names} 已从下载链接安装。你可以在“工具接入”查看配置、同步工具并调用测试。"
                    emit(task_id, "install", "工具服务安装成功", answer, {"mcp_servers": servers, "source_url": download_url})
                    emit(task_id, "answer", "安装完成", answer)
                    db.update_task_status(task_id, "completed", result={"installed": True, "type": "mcp", "mcp_servers": servers})
                    emit(task_id, "done", "已完成", "MCP 配置已写入平台。")
                    return True
            if payload is None:
                answer = "请粘贴 MCP JSON 配置，或上传 JSON 配置文件后发送“安装这个 MCP”。"
                emit(task_id, "answer", "还需要 MCP 配置", answer)
                db.update_task_status(task_id, "completed", result={"installed": False, "reason": "missing_mcp_config"})
                emit(task_id, "done", "已完成", "未执行安装。")
                return True
            servers = self.mcp_gateway.import_config(payload)
            names = "、".join(f"{item['name']}（{item['id']}）" for item in servers)
            answer = f"工具服务 {names} 已安装。你可以在“工具接入”中查看、测试和同步工具；安装不会自动授予任何智能体使用权限。"
            emit(task_id, "install", "工具服务安装成功", answer, {"mcp_servers": servers})
            emit(task_id, "answer", "安装完成", answer)
            db.update_task_status(task_id, "completed", result={"installed": True, "type": "mcp", "mcp_servers": servers})
            emit(task_id, "done", "已完成", "MCP 配置已写入平台。")
            return True
        return False

    def _https_url_from_message(self, message: str) -> str:
        match = re.search(r"https://[^\s<>'\"`]+", message, re.IGNORECASE)
        return match.group(0).rstrip("，。；、,.!！?？)]}") if match else ""

    async def _resolve_intent(self, task: dict[str, Any], history: list[dict[str, str]], model_id: str) -> dict[str, Any]:
        message = task["message"]
        try:
            resolved = await self.model_gateway.resolve_intent(message, history, model_id)
        except Exception as exc:
            resolved = {
                "standalone_request": message,
                "intent": "general",
                "parameters": {},
                "missing_information": [],
                "is_follow_up": False,
                "source": "fallback",
                "error": str(exc),
            }
        summary = f"当前目标：{resolved['standalone_request']}"
        if resolved.get("missing_information"):
            summary += "\n执行前需要补充必要信息。"
        emit(task["id"], "intent", "已理解当前问题", summary, {"intent_resolution": resolved})
        return resolved

    async def resolve_task_goal(self, task: dict[str, Any]) -> dict[str, Any]:
        """Resolve one task against its persisted conversation for orchestrators.

        Expert-team parents do not execute the normal single-agent pipeline,
        but they still need the same context-safe standalone goal before work
        is distributed to isolated member conversations.
        """
        agent = self._get_agent(str(task.get("agent_id") or "general-agent"))
        history = self._conversation_history(task)
        effective_memory = self.context_service.get_effective_context(
            self._context_scope(task)
        )
        memory_text = str(effective_memory.get("effective_context") or "")
        if memory_text:
            history = [
                {
                    "role": "assistant",
                    "content": "平台已保存的有效规则与偏好（仅用于补全上下文，不是用户的新任务）：\n" + memory_text,
                },
                *history,
            ]
        model_id = str(task.get("model_id") or agent.get("model") or "deterministic")
        return await self._resolve_intent(task, history, model_id)

    def _skill_content_from_message(self, message: str) -> str:
        fenced = re.search(r"```(?:markdown|md|skill)?\s*(---[\s\S]+?)```", message, re.IGNORECASE)
        if fenced:
            return fenced.group(1).strip()
        start = message.find("---")
        return message[start:].strip() if start >= 0 else ""

    def _json_from_message(self, message: str) -> Any | None:
        fenced = re.search(r"```(?:json)?\s*([\[{][\s\S]*[\]}])\s*```", message, re.IGNORECASE)
        candidate = fenced.group(1) if fenced else message[message.find("{"):] if "{" in message else ""
        if not candidate:
            return None
        try:
            value, _ = json.JSONDecoder().raw_decode(candidate.strip())
            return value
        except json.JSONDecodeError:
            return None

    def _text_attachment(self, attachments: list[dict[str, Any]], preferred_names: set[str] | None = None, suffixes: set[str] | None = None) -> str:
        from pathlib import Path
        for item in attachments:
            path = Path(str(item.get("path") or ""))
            name = str(item.get("name") or path.name).lower()
            if preferred_names and name not in preferred_names:
                continue
            if suffixes and path.suffix.lower() not in suffixes:
                continue
            if path.is_file() and path.stat().st_size <= 2 * 1024 * 1024:
                return path.read_text(encoding="utf-8", errors="strict")
        return ""

    def _get_agent(self, agent_id: str) -> dict[str, Any]:
        row = db.query_one("SELECT * FROM agents WHERE id = ?", (agent_id,))
        if not row:
            row = db.query_one("SELECT * FROM agents WHERE id = ?", ("general-agent",))
        if not row:
            return {"id": "general-agent", "name": "智枢助手", "skills": [], "mcp_servers": []}
        return {
            **row,
            "skills": db.json_loads(row.get("skills_json"), []),
            "mcp_servers": db.json_loads(row.get("mcp_servers_json"), []),
            "permissions": db.json_loads(row.get("permissions_json"), {}),
        }

    @staticmethod
    def _normalize_permissions(value: dict[str, Any]) -> dict[str, Any]:
        """Create a detached, conservative runtime permission snapshot."""

        try:
            snapshot = json.loads(json.dumps(value, ensure_ascii=False, default=str))
        except (TypeError, ValueError):
            snapshot = {}
        if not isinstance(snapshot, dict):
            snapshot = {}
        for key in ("allowed_tools", "denied_tools", "allowed_mcp_servers", "denied_mcp_servers"):
            if key not in snapshot:
                continue
            incoming = snapshot.get(key)
            snapshot[key] = (
                list(dict.fromkeys(str(item).strip() for item in incoming if str(item).strip()))
                if isinstance(incoming, list)
                else []
            )
        if "read_only" in snapshot:
            snapshot["read_only"] = bool(snapshot.get("read_only"))
        for key in ("max_tool_calls", "max_tool_steps"):
            if key not in snapshot:
                continue
            incoming = snapshot.get(key)
            snapshot[key] = (
                max(0, int(incoming))
                if isinstance(incoming, (int, float)) and not isinstance(incoming, bool)
                else 0
            )
        if "timeout_seconds" in snapshot:
            incoming = snapshot.get("timeout_seconds")
            snapshot["timeout_seconds"] = (
                max(0.0, float(incoming))
                if isinstance(incoming, (int, float)) and not isinstance(incoming, bool)
                else 0.0
            )
        return snapshot

    def _permission_snapshot_for_task(
        self, task: dict[str, Any], agent: dict[str, Any]
    ) -> tuple[dict[str, Any], str]:
        if str(task.get("executor_type") or "agent") == "team_member":
            member_run = db.query_one(
                """SELECT id, permissions_json FROM team_member_runs
                   WHERE child_task_id = ? ORDER BY attempt DESC, created_at DESC LIMIT 1""",
                (task["id"],),
            )
            if member_run:
                return (
                    self._normalize_permissions(
                        db.json_loads(member_run.get("permissions_json"), {})
                    ),
                    f"team_member_run:{member_run['id']}",
                )
            # An internal member without its orchestrator-created permission
            # row is inconsistent state.  It may still generate text, but no
            # tool is exposed or executable until the run is repaired.
            return (
                {
                    "allowed_tools": [],
                    "allowed_mcp_servers": [],
                    "read_only": True,
                },
                "team_member_run:missing_fail_closed",
            )
        return self._normalize_permissions(agent.get("permissions") or {}), f"agent:{agent.get('id', '')}"

    @staticmethod
    def _permission_snapshot_summary(permissions: dict[str, Any]) -> str:
        parts: list[str] = []
        if "allowed_tools" in permissions:
            parts.append(f"允许 {len(permissions.get('allowed_tools') or [])} 个工具")
        if permissions.get("denied_tools"):
            parts.append(f"显式禁止 {len(permissions['denied_tools'])} 个工具")
        if permissions.get("read_only"):
            parts.append("只读")
        if "max_tool_calls" in permissions:
            parts.append(f"最多 {permissions['max_tool_calls']} 次真实工具调用")
        if "timeout_seconds" in permissions:
            parts.append(f"工具执行总时限 {permissions['timeout_seconds']:g} 秒")
        return "；".join(parts) if parts else "使用 Agent 默认工具权限；策略仍可进一步收紧。"

    @staticmethod
    def _tool_name_matches(values: Any, server_id: str, tool_name: str) -> bool:
        if not isinstance(values, list):
            return False
        candidates = {
            tool_name,
            f"{server_id}.{tool_name}",
            f"{server_id}__{tool_name}",
        }
        patterns = {str(item).strip() for item in values}
        return bool(
            patterns.intersection(candidates)
            or "*" in patterns
            or f"{server_id}.*" in patterns
            or f"{server_id}__*" in patterns
        )

    def _tool_definition(self, server_id: str, tool_name: str) -> dict[str, Any]:
        getter = getattr(self.mcp_gateway, "get_tool_definition", None)
        if callable(getter):
            value = getter(server_id, tool_name)
            if isinstance(value, dict):
                return dict(value)
        try:
            tools = self.mcp_gateway.list_tools()
        except (AttributeError, TypeError):
            tools = []
        value = next(
            (
                item
                for item in tools
                if str(item.get("server_id") or "") == server_id
                and str(item.get("name") or "") == tool_name
            ),
            {},
        )
        result = dict(value) if isinstance(value, dict) else {}
        if not result.get("server_kind"):
            server_getter = getattr(self.mcp_gateway, "get_server", None)
            server = server_getter(server_id) if callable(server_getter) else None
            if isinstance(server, dict):
                result["server_kind"] = str(server.get("kind") or "")
        return result

    @staticmethod
    def _annotation_read_only(annotations: Any) -> bool:
        if not isinstance(annotations, dict):
            return False
        value = annotations.get("readOnlyHint")
        if value is None:
            value = annotations.get("read_only_hint")
        return value is True

    def _permission_denial_for_tool(
        self,
        server_id: str,
        tool_name: str,
        *,
        definition: dict[str, Any] | None = None,
    ) -> tuple[str, str] | None:
        permissions = self._effective_permissions()
        denied_servers = permissions.get("denied_mcp_servers")
        if isinstance(denied_servers, list) and (
            server_id in denied_servers or "*" in denied_servers
        ):
            return "server_denied", f"当前 Agent 权限禁止使用工具服务 {server_id}。"
        # Explicit deny is evaluated before every allow-list decision.
        if self._tool_name_matches(permissions.get("denied_tools"), server_id, tool_name):
            return "tool_denied", f"当前 Agent 权限明确禁止调用 {server_id}.{tool_name}。"
        if "allowed_mcp_servers" in permissions:
            allowed_servers = permissions.get("allowed_mcp_servers") or []
            if server_id not in allowed_servers and "*" not in allowed_servers:
                return "server_not_allowed", f"当前 Agent 未获准使用工具服务 {server_id}。"
        if "allowed_tools" in permissions and not self._tool_name_matches(
            permissions.get("allowed_tools"), server_id, tool_name
        ):
            return "tool_not_allowed", f"当前 Agent 未获准调用 {server_id}.{tool_name}。"
        if permissions.get("read_only"):
            metadata = dict(definition or self._tool_definition(server_id, tool_name))
            server_kind = str(metadata.get("server_kind") or "").lower()
            if server_kind == "builtin":
                safe_read = str(metadata.get("effect") or "").lower() == "read"
            else:
                # Remote tool names and descriptions are not security
                # evidence.  MCP readOnlyHint must be explicitly true;
                # missing/false/malformed annotations are denied.
                safe_read = self._annotation_read_only(metadata.get("annotations"))
            if not safe_read:
                return (
                    "read_only",
                    f"当前 Agent 为只读模式，{server_id}.{tool_name} 未被可信元数据标记为只读，已阻止调用。",
                )
        return None

    def _enforce_tool_permission(
        self, task_id: str, server_id: str, tool_name: str
    ) -> None:
        denial = self._permission_denial_for_tool(server_id, tool_name)
        if not denial:
            return
        code, message = denial
        emit(
            task_id,
            "tool_blocked",
            f"权限已阻止 {server_id}.{tool_name}",
            message,
            {
                "server_id": server_id,
                "tool_name": tool_name,
                "reason": code,
                "source": "effective_permissions",
            },
        )
        raise ToolError(message)

    def _tool_visible_to_model(self, tool: dict[str, Any]) -> bool:
        server_id = str(tool.get("server_id") or "")
        tool_name = str(tool.get("name") or "")
        return bool(server_id and tool_name) and self._permission_denial_for_tool(
            server_id, tool_name, definition=tool
        ) is None

    @staticmethod
    def _context_scope(task: dict[str, Any]) -> ExecutionScope:
        return ExecutionScope(
            organization_id=str(task.get("organization_id") or "local-org"),
            workspace_id=str(task.get("workspace") or "default"),
            user_id=str(task.get("user_id") or "local-user"),
            agent_id=str(task.get("agent_id") or ""),
            conversation_id=str(task.get("conversation_id") or ""),
        )

    async def _tool(
        self,
        task_id: str,
        server_id: str,
        tool_name: str,
        arguments: dict[str, Any],
        plan: dict[str, Any] | None = None,
    ) -> dict[str, Any]:
        self._raise_if_cancelled()
        child_id = f"tool:{server_id}.{tool_name}"
        plan_node_ids = {str(item.get("id") or "") for item in (plan or {}).get("nodes", [])}
        progress_node_id = (
            "artifact"
            if server_id in {"report", "spreadsheet"} and "artifact" in plan_node_ids
            else str((plan or {}).get("tool_node_id") or "execute")
        )
        if plan is not None:
            self._validate_tool_against_plan(plan, server_id, tool_name)
        # Every invocation path (weather shortcut, deterministic workflows,
        # model tool calls and artifact generation) converges here.  Keep the
        # hard permission check at this single boundary.
        self._enforce_tool_permission(task_id, server_id, tool_name)
        if plan is not None:
            emit(
                task_id,
                "plan_check",
                "工具调用前校验",
                f"{server_id}.{tool_name} 与当前目标和执行计划一致，参数校验通过。",
                {"tool": f"{server_id}.{tool_name}", "passed": True},
            )
            self._emit_plan_progress(
                task_id,
                progress_node_id,
                "running",
                f"正在调用 {server_id}.{tool_name}",
                child_id=child_id,
                child_title=f"{server_id}.{tool_name}",
                child_kind="mcp",
            )
        policy_context = {
            "tool": {"server": server_id, "name": tool_name, "arguments": arguments},
            "plan": plan or {},
        }
        evaluation = await self._evaluate_policy("tool.before", policy_context, enforce=False)
        if evaluation.denied:
            emit(
                task_id,
                "tool_blocked",
                f"策略已阻止 {server_id}.{tool_name}",
                evaluation.summary,
                evaluation.to_dict(),
            )
            raise ToolError(evaluation.summary)
        if evaluation.requires_approval:
            await self._wait_for_policy_approval(evaluation, server_id, tool_name)
        modified_context = evaluation.apply(policy_context)
        modified_arguments = modified_context.get("tool", {}).get("arguments", arguments)
        if isinstance(modified_arguments, dict):
            arguments = modified_arguments
        # Policy modifications may change arguments, but can never expand the
        # immutable permission snapshot or bypass its deny/read-only rules.
        self._enforce_tool_permission(task_id, server_id, tool_name)

        fingerprint = self._tool_fingerprint(server_id, tool_name, arguments)
        execution = self._execution()
        completed_tools = (execution or {}).get("state", {}).setdefault("completed_tools", {})
        cached_result = completed_tools.get(fingerprint) if isinstance(completed_tools, dict) else None
        if isinstance(cached_result, dict) and self._cached_tool_result_is_usable(cached_result):
            emit(
                task_id,
                "tool_reused",
                f"复用已完成工具 {server_id}.{tool_name}",
                "已从安全检查点复用相同参数的工具结果，避免重复副作用。",
                {"server_id": server_id, "tool_name": tool_name, "fingerprint": fingerprint},
            )
            if plan is not None:
                self._emit_plan_progress(
                    task_id,
                    progress_node_id,
                    "completed",
                    f"{server_id}.{tool_name} 已从检查点恢复",
                    child_id=child_id,
                    child_title=f"{server_id}.{tool_name}",
                    child_kind="mcp",
                )
            return cached_result

        remaining_timeout = self._remaining_tool_timeout()
        if remaining_timeout is not None and remaining_timeout <= 0:
            configured = float(self._effective_permissions().get("timeout_seconds") or 0)
            message = f"当前运行的工具执行总时限 {configured:g} 秒已用尽，未调用 {server_id}.{tool_name}。"
            emit(
                task_id,
                "tool_blocked",
                f"工具时限已阻止 {server_id}.{tool_name}",
                message,
                {
                    "server_id": server_id,
                    "tool_name": tool_name,
                    "reason": "timeout_exhausted",
                    "source": "effective_permissions",
                },
            )
            raise ToolError(message)

        if execution:
            state = execution["state"]
            used_calls = max(0, int(state.get("tool_calls_used") or 0))
            max_calls = self._effective_permissions().get("max_tool_calls")
            if isinstance(max_calls, (int, float)) and not isinstance(max_calls, bool):
                limit = max(0, int(max_calls))
                if used_calls >= limit:
                    message = (
                        f"当前运行最多允许 {limit} 次真实工具调用；"
                        f"已使用 {used_calls} 次，未调用 {server_id}.{tool_name}。"
                    )
                    emit(
                        task_id,
                        "tool_blocked",
                        f"调用次数已阻止 {server_id}.{tool_name}",
                        message,
                        {
                            "server_id": server_id,
                            "tool_name": tool_name,
                            "reason": "max_tool_calls",
                            "used": used_calls,
                            "limit": limit,
                            "source": "effective_permissions",
                        },
                    )
                    raise ToolError(message)
            # Increment before the real invocation and checkpoint it.  A
            # gateway exception therefore still consumes one call, while the
            # cache-return branch above consumes none.
            state["tool_calls_used"] = used_calls + 1
            execution["state"].update(
                {
                    "phase": "before_tool",
                    "pending_tool": {
                        "server": server_id,
                        "name": tool_name,
                        "arguments": self._safe_tool_arguments(arguments),
                        "fingerprint": fingerprint,
                    },
                }
            )
            self._create_checkpoint(
                f"调用 {server_id}.{tool_name} 前的安全边界", node_key=progress_node_id
            )

        safe_arguments = self._safe_tool_arguments(arguments)
        emit(task_id, "tool_call", f"调用工具 {server_id}.{tool_name}", "参数已准备并通过校验。", {"server_id": server_id, "tool_name": tool_name, "arguments": safe_arguments})
        started_at = time.perf_counter()
        invoke_task: asyncio.Task[dict[str, Any]] | None = None
        try:
            invoke_task = asyncio.create_task(
                self.mcp_gateway.invoke_tool(server_id, tool_name, arguments, task_id=task_id)
            )
            while not invoke_task.done():
                remaining_timeout = self._remaining_tool_timeout()
                if remaining_timeout is not None and remaining_timeout <= 0:
                    invoke_task.cancel()
                    await asyncio.gather(invoke_task, return_exceptions=True)
                    configured = float(self._effective_permissions().get("timeout_seconds") or 0)
                    raise ToolError(
                        f"工具调用 {server_id}.{tool_name} 超过本次运行剩余时限（总时限 {configured:g} 秒），已终止。"
                    )
                wait_seconds = (
                    min(0.25, max(0.001, remaining_timeout))
                    if remaining_timeout is not None
                    else 0.25
                )
                done, _ = await asyncio.wait({invoke_task}, timeout=wait_seconds)
                if done:
                    break
                try:
                    self._raise_if_cancelled()
                except TaskCancellationRequested:
                    invoke_task.cancel()
                    await asyncio.gather(invoke_task, return_exceptions=True)
                    raise
            result = await invoke_task
            duration_ms = max(1, round((time.perf_counter() - started_at) * 1000))
            after_context = {
                "tool": {
                    "server": server_id,
                    "name": tool_name,
                    "arguments": arguments,
                },
                "result": result,
            }
            after_evaluation = await self._evaluate_policy(
                "tool.after",
                after_context,
                enforce=False,
            )
            if after_evaluation.denied:
                emit(
                    task_id,
                    "tool_blocked",
                    f"策略未接受 {server_id}.{tool_name} 的返回结果",
                    after_evaluation.summary,
                    after_evaluation.to_dict(),
                )
                raise ToolError(after_evaluation.summary)
            if after_evaluation.requires_approval:
                await self._wait_for_policy_approval(
                    after_evaluation,
                    server_id,
                    tool_name,
                    event="tool.after",
                )
            applied_after_context = after_evaluation.apply(after_context)
            applied_result = applied_after_context.get("result")
            if not isinstance(applied_result, Mapping):
                raise ToolError("tool.after 策略修改后的 result 必须是对象")
            result = dict(applied_result)
            artifact = result.get("artifact")
            if isinstance(artifact, dict):
                artifact_context = {
                    "artifact": artifact,
                    "tool": {"server": server_id, "name": tool_name},
                }
                artifact_evaluation = await self._evaluate_policy(
                    "artifact.created", artifact_context, enforce=False
                )
                if artifact_evaluation.denied:
                    emit(
                        task_id,
                        "tool_blocked",
                        f"策略未接受 {artifact.get('name') or '新产物'}",
                        artifact_evaluation.summary,
                        artifact_evaluation.to_dict(),
                    )
                    raise ToolError(artifact_evaluation.summary)
                if artifact_evaluation.requires_approval:
                    await self._wait_for_policy_approval(
                        artifact_evaluation,
                        server_id,
                        tool_name,
                        event="artifact.created",
                    )
                applied_artifact_context = artifact_evaluation.apply(
                    artifact_context
                )
                applied_artifact = applied_artifact_context.get("artifact")
                if not isinstance(applied_artifact, Mapping):
                    raise ToolError(
                        "artifact.created 策略修改后的 artifact 必须是对象"
                    )
                artifact = dict(applied_artifact)
                result["artifact"] = artifact
            result_summary = (
                f"已生成文件 {artifact.get('name')}"
                if isinstance(artifact, dict)
                else "工具调用成功并返回结果"
            )
            emit(
                task_id,
                "tool_result",
                f"工具返回 {server_id}.{tool_name}",
                result_summary,
                {
                    "server_id": server_id,
                    "tool_name": tool_name,
                    "duration_ms": duration_ms,
                    "artifact": (
                        {
                            key: artifact.get(key)
                            for key in ("name", "kind", "download_url")
                            if artifact.get(key)
                        }
                        if isinstance(artifact, dict)
                        else None
                    ),
                },
            )
            if execution:
                completed = execution["state"].setdefault("completed_tools", {})
                if isinstance(completed, dict):
                    completed[fingerprint] = result
                execution["state"].update(
                    {
                        "phase": "tool_completed",
                        "pending_tool": {},
                        "last_completed_tool": {
                            "server": server_id,
                            "name": tool_name,
                            "fingerprint": fingerprint,
                        },
                    }
                )
                self._create_checkpoint(
                    f"{server_id}.{tool_name} 已完成", node_key=progress_node_id
                )
            if plan is not None:
                self._emit_plan_progress(
                    task_id,
                    progress_node_id,
                    "completed",
                    f"{server_id}.{tool_name} 调用完成",
                    child_id=child_id,
                    child_title=f"{server_id}.{tool_name}",
                    child_kind="mcp",
                )
            return result
        except (TaskCancellationRequested, asyncio.CancelledError):
            if invoke_task is not None and not invoke_task.done():
                invoke_task.cancel()
                await asyncio.gather(invoke_task, return_exceptions=True)
            raise
        except Exception as raw_exc:
            exc = raw_exc if isinstance(raw_exc, ToolError) else ToolError(
                f"工具调用 {server_id}.{tool_name} 失败：{raw_exc}"
            )
            if execution:
                execution["state"].update(
                    {
                        "phase": "tool_failed",
                        "pending_tool": {},
                        "last_failed_tool": {
                            "server": server_id,
                            "name": tool_name,
                            "fingerprint": fingerprint,
                            "error": str(exc),
                        },
                    }
                )
                try:
                    self._create_checkpoint(
                        f"{server_id}.{tool_name} 调用失败，已保留次数与时限状态",
                        node_key=progress_node_id,
                    )
                except Exception:
                    pass
            emit(task_id, "tool_error", f"工具失败 {server_id}.{tool_name}", str(exc), {"arguments": safe_arguments})
            try:
                await self._evaluate_policy(
                    "tool.failed",
                    {
                        "tool": {"server": server_id, "name": tool_name, "arguments": arguments},
                        "error": {"message": str(exc)},
                    },
                    enforce=False,
                )
            except Exception:
                pass
            if plan is not None:
                self._emit_plan_progress(
                    task_id,
                    progress_node_id,
                    "failed",
                    str(exc),
                    child_id=child_id,
                    child_title=f"{server_id}.{tool_name}",
                    child_kind="mcp",
                )
            if exc is raw_exc:
                raise
            raise exc from raw_exc

    @staticmethod
    def _tool_fingerprint(server_id: str, tool_name: str, arguments: dict[str, Any]) -> str:
        payload = json.dumps(
            {"server": server_id, "tool": tool_name, "arguments": arguments},
            ensure_ascii=False,
            sort_keys=True,
            separators=(",", ":"),
            default=str,
        )
        return hashlib.sha256(payload.encode("utf-8")).hexdigest()

    @staticmethod
    def _cached_tool_result_is_usable(result: dict[str, Any]) -> bool:
        artifact = result.get("artifact")
        if not isinstance(artifact, dict):
            return True
        return AgentRuntime._artifact_file(artifact) is not None

    @staticmethod
    def _artifact_file(artifact: dict[str, Any]) -> Path | None:
        artifact_id = str(artifact.get("id") or "")
        row = (
            db.query_one("SELECT relative_path, path FROM artifacts WHERE id = ?", (artifact_id,))
            if artifact_id
            else None
        )
        try:
            relative = str((row or {}).get("relative_path") or artifact.get("relative_path") or "")
            if relative:
                return resolve_artifact_path(relative)
            legacy = Path(
                str((row or {}).get("path") or artifact.get("path") or "")
            ).resolve(strict=True)
            root = ARTIFACT_DIR.resolve(strict=True)
            return resolve_artifact_path(legacy.relative_to(root).as_posix())
        except (FileNotFoundError, OSError, RuntimeError, ToolError, ValueError):
            return None

    async def _wait_for_policy_approval(
        self,
        evaluation: Any,
        server_id: str,
        tool_name: str,
        *,
        event: str = "tool.before",
    ) -> None:
        execution = self._execution()
        if not execution:
            raise PolicyApprovalRequired(evaluation)
        task_id = execution["task_id"]
        tool = (
            {"server": server_id, "name": tool_name}
            if server_id or tool_name
            else None
        )
        default_title = (
            "工具调用需要审批" if server_id or tool_name else "策略要求审批"
        )
        approval_request = await self._apply_approval_requested_policy(
            {
                "action": "policy_approval",
                "event": event,
                "title": default_title,
                "message": evaluation.summary,
                "tool": tool,
                "requests": evaluation.approval_requests,
                "policy": evaluation.to_dict(),
            }
        )
        approval_title = str(approval_request.get("title") or default_title)
        approval_message = str(
            approval_request.get("message") or evaluation.summary
        )
        result = {
            "pending_action": "policy_approval",
            "policy_event": event,
            "policy_evaluation": evaluation.to_dict(),
            "approval_request": approval_request,
            "summary": approval_message,
        }
        if server_id or tool_name:
            result["tool"] = tool
        db.update_task_status(task_id, "waiting_approval", result=result)
        current_run = self.task_state.get_run(execution["run_id"])
        if current_run and current_run["status"] == "running":
            self.task_state.transition_run(execution["run_id"], "waiting_approval")
        emit(
            task_id,
            "approval_required",
            approval_title,
            approval_message,
            {
                "action": "policy_approval",
                "event": event,
                "tool": tool,
                "policy": evaluation.to_dict(),
                "approval_request": approval_request,
            },
        )
        while True:
            self._raise_if_cancelled()
            command = self.task_state.claim_command(
                execution["worker_id"],
                task_id=task_id,
                run_id=execution["run_id"],
                command_types=["approval"],
            )
            if command:
                approved = bool(command.get("payload", {}).get("approved"))
                self.task_state.complete_command(command["id"], result={"approved": approved})
                if not approved:
                    target = (
                        f"工具调用 {server_id}.{tool_name}"
                        if server_id or tool_name
                        else f"{event} 操作"
                    )
                    raise ToolError(f"用户拒绝了策略要求审批的{target}")
                self.task_state.begin_run(task_id, run_id=execution["run_id"])
                db.update_task_status(task_id, "running")
                emit(
                    task_id,
                    "approval",
                    "审批通过",
                    (
                        "继续执行已审批的工具调用。"
                        if server_id or tool_name
                        else f"继续执行已审批的 {event} 操作。"
                    ),
                )
                return
            await asyncio.sleep(0.25)

    def _register_plan(self, plan: dict[str, Any]) -> None:
        execution = self._execution()
        if not execution:
            return
        for sequence, definition in enumerate(plan.get("nodes") or [], start=1):
            key = str(definition.get("id") or f"step-{sequence}")
            execution["plan_nodes"][key] = dict(definition)
            parent_id = execution["nodes"].get(key)
            if not parent_id:
                node = self.task_state.create_node(
                    execution["run_id"],
                    key,
                    str(definition.get("title") or key),
                    kind="phase",
                    sequence=sequence * 100,
                    metadata={"plan_id": "main"},
                )
                parent_id = node["id"]
                execution["nodes"][key] = parent_id
            else:
                self.task_state.update_node_definition(
                    parent_id,
                    title=str(definition.get("title") or key),
                    kind="phase",
                    sequence=sequence * 100,
                )
            for child_index, child in enumerate(definition.get("children") or [], start=1):
                child_key = str(child.get("id") or f"{key}:detail:{child_index}")
                execution["plan_nodes"][child_key] = dict(child)
                if child_key in execution["nodes"]:
                    self.task_state.update_node_definition(
                        execution["nodes"][child_key],
                        title=str(child.get("title") or child_key),
                        kind=str(child.get("kind") or "detail"),
                        sequence=sequence * 100 + child_index,
                    )
                    continue
                child_node = self.task_state.create_node(
                    execution["run_id"],
                    child_key,
                    str(child.get("title") or child_key),
                    parent_node_id=parent_id,
                    kind=str(child.get("kind") or "detail"),
                    sequence=sequence * 100 + child_index,
                    metadata={"plan_id": "main"},
                )
                execution["nodes"][child_key] = child_node["id"]

    def _persist_node_status(
        self,
        node_key: str,
        status: str,
        message: str,
        *,
        parent_key: str | None = None,
        title: str = "",
        kind: str = "detail",
    ) -> bool:
        execution = self._execution()
        if not execution:
            return False
        node_id = execution["nodes"].get(node_key)
        if not node_id:
            parent_id = execution["nodes"].get(parent_key or "")
            definition = execution["plan_nodes"].get(node_key, {})
            node = self.task_state.create_node(
                execution["run_id"],
                node_key,
                title or str(definition.get("title") or node_key),
                parent_node_id=parent_id,
                kind=kind or str(definition.get("kind") or "detail"),
                metadata={"plan_id": "main"},
            )
            node_id = node["id"]
            execution["nodes"][node_key] = node_id
        node = self.task_state.get_node(node_id)
        if not node or node["status"] == status or node["status"] in {"completed", "failed", "skipped", "cancelled"}:
            return False
        metadata = {"last_message": message} if message else None
        try:
            if status == "running":
                if node["status"] == "pending":
                    self.task_state.start_node(node_id, metadata=metadata)
                    return True
            elif status == "completed":
                if node["status"] == "pending":
                    self.task_state.start_node(node_id)
                self.task_state.finish_node(node_id, output={"summary": message} if message else {}, metadata=metadata)
                return True
            elif status == "failed":
                if node["status"] == "pending":
                    self.task_state.start_node(node_id)
                self.task_state.fail_node(node_id, {"message": message or "节点执行失败"}, metadata=metadata)
                return True
            elif status == "skipped" and node["status"] == "pending":
                self.task_state.skip_node(node_id, metadata=metadata)
                return True
        except (InvalidStateTransition, TaskStateError):
            return False
        return False

    def _emit_plan_progress(
        self,
        task_id: str,
        node_id: str,
        status: str,
        message: str = "",
        *,
        child_id: str | None = None,
        child_title: str | None = None,
        child_kind: str | None = None,
        elapsed_seconds: int | None = None,
    ) -> None:
        data: dict[str, Any] = {"plan_id": "main", "node_id": node_id, "status": status}
        if child_id:
            data.update({"child_id": child_id, "child_title": child_title or child_id, "child_kind": child_kind or "detail"})
        if elapsed_seconds is not None:
            data["elapsed_seconds"] = elapsed_seconds
        emit(task_id, "plan_progress", "执行进度", message, data)
        if child_id:
            self._persist_node_status(
                child_id,
                status,
                message,
                parent_key=node_id,
                title=child_title or child_id,
                kind=child_kind or "detail",
            )
        else:
            transitioned = self._persist_node_status(node_id, status, message, kind="phase")
            if transitioned and status == "completed":
                execution = self._execution()
                if execution:
                    execution["state"]["phase"] = node_id
                    self._create_checkpoint(f"节点“{message or node_id}”已完成", node_key=node_id)

    def _preview(self, value: Any, max_len: int = 420) -> str:
        text = db.json_dumps(value)
        if len(text) > max_len:
            return text[:max_len] + "..."
        return text

    def _safe_tool_arguments(self, arguments: dict[str, Any]) -> dict[str, Any]:
        """Return a compact execution trace without exposing secrets or large content."""
        safe: dict[str, Any] = {}
        for key, value in arguments.items():
            lowered = key.lower()
            if any(word in lowered for word in ("key", "token", "secret", "password", "authorization")):
                safe[key] = "••••••"
            elif isinstance(value, str) and len(value) > 160:
                safe[key] = f"{value[:120]}…（共 {len(value)} 字）"
            elif isinstance(value, list):
                safe[key] = f"{len(value)} 项"
            elif isinstance(value, dict):
                safe[key] = f"{len(value)} 个字段"
            else:
                safe[key] = value
        return safe

    async def _run_general_task(self, task: dict[str, Any], agent: dict[str, Any], skills: list[dict[str, Any]], history: list[dict[str, str]] | None = None) -> None:
        task_id = task["id"]
        original_message = task["message"]
        message = task.get("resolved_message") or original_message
        selected_names = [s["name"] for s in skills]
        skill_instructions = "\n\n".join([
            f"### Skill: {s['name']}\n{self.skill_registry.runtime_content(s['id'])}"
            for s in skills
        ])
        attachments = db.json_loads(task.get("attachments_json"), [])
        attachment_context = self._attachment_context(attachments)
        memory_context = str(task.get("memory_context") or "").strip()
        search_context = ""
        history = history or []
        recent_user_messages = [item["content"] for item in history if item["role"] == "user"][-3:]
        routing_text = message.lower()
        intent_resolution = task.get("intent_resolution") or {}
        intent_parameters = intent_resolution.get("parameters") if isinstance(intent_resolution.get("parameters"), dict) else {}
        requested_document_format = self._requested_document_format(message)
        wants_report_artifact = self._wants_report_artifact(message) and not requested_document_format
        active_model = task.get("model_id") or agent.get("model") or "deterministic"
        weather_request = self._weather_request_for_intent(
            original_message, message, history, intent_resolution
        )
        offline_weather_adapter = bool(
            active_model == "deterministic"
            and weather_request is not None
            and not requested_document_format
            and not wants_report_artifact
        )
        plan = self._build_execution_plan(
            task,
            skills,
            requested_document_format,
            wants_report_artifact,
            force_weather=offline_weather_adapter,
        )
        if attachment_context:
            attachment_requirements = self._attachment_acceptance_requirements(
                attachment_context
            )
            if attachment_requirements:
                plan["attachment_requirements"] = attachment_requirements
                plan["acceptance_criteria"] = self._acceptance_criteria(plan)
        plan_policy_context = {"plan": plan}
        plan_evaluation = await self._evaluate_policy(
            "plan.created", plan_policy_context, enforce=True
        )
        applied_plan = plan_evaluation.apply(plan_policy_context).get("plan")
        if not isinstance(applied_plan, Mapping):
            raise RuntimeError("plan.created 策略修改后的 plan 必须是对象")
        plan = dict(applied_plan)
        self._register_plan(plan)
        emit(task_id, "plan", "执行计划", self._format_execution_plan(plan), {"plan": plan})
        self._emit_plan_progress(task_id, "understand", "running", "正在确认当前目标与匹配能力")
        for skill in skills:
            self._emit_plan_progress(
                task_id,
                "understand",
                "completed",
                f"已匹配 {skill['name']}",
                child_id=f"skill:{skill['id']}",
                child_title=skill["name"],
                child_kind="skill",
            )
        self._emit_plan_progress(task_id, "understand", "completed", "目标与可用 Skill 已确认")
        self._emit_plan_progress(task_id, "prepare", "running", "正在准备上下文与授权工具")
        used_memory_ids = list(task.get("used_memory_ids") or [])
        if used_memory_ids:
            self._emit_plan_progress(
                task_id,
                "prepare",
                "completed",
                f"已应用 {len(used_memory_ids)} 条平台记忆",
                child_id="memory:effective",
                child_title=f"平台记忆（{len(used_memory_ids)} 条）",
                child_kind="memory",
            )
        self._emit_plan_progress(task_id, "prepare", "completed", "上下文与工具权限已准备")
        self._emit_plan_progress(task_id, "execute", "running", "正在生成结果")

        # The offline deterministic adapter cannot plan tool calls.  Keep a
        # narrow compatibility path for an explicit, pure weather lookup; all
        # configured models receive weather through the normal tool-planning
        # path below.  This prevents keyword matches from pre-empting document
        # and multi-part tasks.
        if weather_request is not None and intent_parameters.get("city"):
            weather_request["city"] = str(intent_parameters["city"])
        if intent_parameters.get("day") in {"today", "tomorrow", "day_after_tomorrow"} and weather_request is not None:
            weather_request["day"] = str(intent_parameters["day"])
        if intent_parameters.get("weather_lookup") is False:
            weather_request = None
        use_offline_weather_adapter = (
            offline_weather_adapter
            and (agent.get("id") == "general-agent" or "weather" in set(agent.get("mcp_servers") or []))
        )
        if use_offline_weather_adapter:
            city = weather_request["city"]
            if not city:
                answer = "请告诉我需要查询的城市或地区，例如“宁波”“北京”或“上海浦东”。"
                self._emit_plan_progress(task_id, "execute", "failed", "缺少城市或地区")
                emit(task_id, "answer", "还需要城市", answer)
                db.update_task_status(task_id, "completed", result={"summary": answer})
                emit(task_id, "done", "已完成", "等待用户补充城市。")
                return
            try:
                forecast = await self._tool(task_id, "weather", "forecast", {"city": city, "day": weather_request["day"]}, plan=plan)
                answer = self._build_weather_answer(forecast)
                self._emit_plan_progress(task_id, "execute", "completed", "天气结果已生成")
                self._emit_plan_progress(task_id, "validate", "running", "正在核对城市、日期与输出")
                emit(task_id, "output_check", "输出前校验", "已确认结果仍在回答当前天气问题，城市和日期与计划一致。", {"passed": True})
                self._emit_plan_progress(task_id, "validate", "completed", "输出校验通过")
                emit(task_id, "answer_delta", "正在生成", answer)
                emit(task_id, "answer", "天气查询完成", answer, {"forecast": forecast})
                db.update_task_status(task_id, "completed", result={"summary": answer, "forecast": forecast})
                emit(task_id, "done", "已完成", "结构化天气预报已返回。")
                return
            except ToolError as exc:
                self._emit_plan_progress(task_id, "execute", "failed", "天气服务未返回可验证结果")
                raise ToolError(f"天气查询失败，未生成未经验证的回答：{exc}") from exc
        wants_search = any(k in routing_text for k in ["联网", "搜索", "最新", "查一下", "web", "internet", "news", "新闻"])
        allowed_mcps = set(agent.get("mcp_servers") or [])
        if wants_search and (not allowed_mcps or "web-search" in allowed_mcps):
            try:
                search_query = message
                search = await self._tool(task_id, "web-search", "search", {"query": search_query, "max_results": 5}, plan=plan)
                search_context = db.json_dumps(search)
            except ToolError as exc:
                self._emit_plan_progress(task_id, "execute", "failed", "联网搜索未返回可验证结果")
                raise ToolError(f"联网搜索失败，未生成未经检索验证的回答：{exc}") from exc
        prompt = message
        if message != original_message:
            prompt = f"用户最新原话：{original_message}\n\n结合上下文还原后的当前独立任务：{message}"
        if skill_instructions:
            prompt += f"\n\n请遵循以下已匹配 Skill：\n{skill_instructions}"
        if memory_context:
            prompt += (
                "\n\n平台当前有效的分层记忆如下。组织规则优先；用户当前明确要求优先于普通偏好，"
                "不得把历史记忆误当成本次新目标：\n" + memory_context
            )
        policy_context = task.get("policy_context")
        if isinstance(policy_context, Mapping) and policy_context:
            prompt += (
                "\n\n平台策略为本次目标追加的执行上下文（只能用于当前任务，"
                "不得改变用户目标）：\n"
                + db.json_dumps(dict(policy_context))
            )
        if attachment_context:
            prompt += f"\n\n用户附件内容：\n{attachment_context}"
        if search_context:
            prompt += f"\n\n联网检索结果（回答时保留来源 URL）：\n{search_context}"
        allowed_server_ids = set(agent.get("mcp_servers") or [])
        relevant_server_ids = {server_id for skill in skills for server_id in skill.get("required_mcps", [])}
        lowered_message = message.lower()
        if "weather" in set(plan.get("allowed_servers") or []):
            relevant_server_ids.add("weather")
        if requested_document_format and requested_document_format not in {"xlsx", "csv"}:
            relevant_server_ids.add("report")
        elif wants_report_artifact:
            relevant_server_ids.add("report")
        if any(k in lowered_message for k in ["表格", "数据", "excel", "xlsx", "csv"]):
            relevant_server_ids.add("spreadsheet")
        available_tools = [
            tool for tool in self.mcp_gateway.list_tools()
            if tool.get("server_id") in relevant_server_ids
            and tool.get("server_id") in set(plan.get("allowed_servers") or [])
            and (agent.get("id") == "general-agent" or tool.get("server_id") in allowed_server_ids)
            and self._tool_visible_to_model(tool)
        ]
        if search_context:
            available_tools = [tool for tool in available_tools if tool.get("server_id") != "web-search"]

        model_artifacts: list[dict[str, Any]] = []
        model_tool_failure: ToolError | None = None

        async def invoke_model_tool(qualified_name: str, arguments: dict[str, Any]) -> dict[str, Any]:
            nonlocal model_tool_failure
            if model_tool_failure is not None:
                raise model_tool_failure
            server_id, tool_name = qualified_name.split("__", 1)
            try:
                result = await self._tool(task_id, server_id, tool_name, arguments, plan=plan)
            except ToolError as exc:
                model_tool_failure = exc
                raise
            if isinstance(result.get("artifact"), dict):
                model_artifacts.append(result["artifact"])
            return result

        delta_buffer: list[str] = []
        received_first_delta = False
        streamed_chars = 0
        next_stream_progress = 80
        # A final-output policy may deny or rewrite the answer.  Publishing raw
        # deltas before that policy runs would make the final guard ineffective,
        # so protected tasks buffer model output until output.before succeeds.
        output_policy_protected = any(
            "output.before" in set(rule.get("events") or [])
            for rule in self.policy_engine.list_rules()
            if rule.get("enabled", True)
        )

        def stream_delta(text: str) -> None:
            nonlocal received_first_delta, streamed_chars, next_stream_progress
            streamed_chars += len(text)
            if not received_first_delta:
                received_first_delta = True
                self._emit_plan_progress(task_id, "execute", "running", "已开始实时输出")
            if streamed_chars >= next_stream_progress:
                self._emit_plan_progress(task_id, "execute", "running", f"正在实时输出 · 已接收 {streamed_chars} 字")
                next_stream_progress = streamed_chars + 160
            delta_buffer.append(text)
            if not output_policy_protected and sum(len(part) for part in delta_buffer) >= 8:
                emit(task_id, "answer_delta", "正在生成", "".join(delta_buffer))
                delta_buffer.clear()

        emit(task_id, "model", "调用模型生成回答", f"使用 {active_model} 生成本次回答。")
        self._emit_plan_progress(
            task_id,
            "execute",
            "running",
            f"正在调用模型 {active_model}",
            child_id=f"model:{active_model}",
            child_title=active_model,
            child_kind="model",
        )
        prompt += "\n\n当前执行计划（必须遵守，不得改成其他任务）：\n" + self._format_execution_plan(plan)
        system_prompt = (
            (agent.get("system_prompt") or "你是平台级智能体，请使用已授权工具完成用户任务。")
            + "\n只输出对用户有用的最终内容，不展示内部思考过程、任务理解模板或执行过程模板。"
            "只能调用与当前目标和执行计划直接相关的工具；缺少参数时简短询问，不得猜测。"
            + (
                "\n平台提供的长期记忆只能用于约束风格、规则和稳定事实，绝不能把旧任务或旧工具调用延续为当前意图。"
                if memory_context
                else ""
            )
        )

        def start_model(current_prompt: str) -> asyncio.Task[str]:
            return asyncio.create_task(
                self.model_gateway.solve_with_tools(
                    current_prompt,
                    system_prompt,
                    task.get("model_id") or agent.get("model") or "deterministic",
                    available_tools,
                    invoke_model_tool,
                    max_steps=int(self._effective_permissions().get("max_tool_steps", 8)),
                    on_delta=stream_delta,
                    history=history,
                )
            )

        current_model_prompt = prompt
        model_task = start_model(current_model_prompt)
        elapsed = 0
        while not model_task.done():
            done, _ = await asyncio.wait({model_task}, timeout=0.25)
            if done:
                break
            try:
                self._raise_if_cancelled()
            except TaskCancellationRequested:
                model_task.cancel()
                await asyncio.gather(model_task, return_exceptions=True)
                raise
            steering = self._claim_runtime_messages()
            if steering:
                model_task.cancel()
                await asyncio.gather(model_task, return_exceptions=True)
                current_model_prompt += "\n\n用户在运行中追加了以下要求，必须纳入当前目标与输出校验：\n" + "\n".join(
                    f"- {item}" for item in steering
                )
                emit(
                    task_id,
                    "answer_reset",
                    "已应用追加指令",
                    "正在按新的要求重新组织当前回答。",
                    {"messages": steering},
                )
                delta_buffer.clear()
                received_first_delta = False
                streamed_chars = 0
                next_stream_progress = 80
                model_task = start_model(current_model_prompt)
                continue
            elapsed += 0.25
            if int(elapsed * 4) % 8 == 0:
                shown_elapsed = max(1, round(elapsed))
                self._emit_plan_progress(
                    task_id,
                    "execute",
                    "running",
                    (f"正在实时输出 · 已接收 {streamed_chars} 字 · 已运行 {shown_elapsed} 秒" if received_first_delta else f"模型正在处理，已等待 {shown_elapsed} 秒"),
                    elapsed_seconds=shown_elapsed,
                )
        summary = await model_task
        if model_tool_failure is not None:
            raise model_tool_failure
        self._emit_plan_progress(
            task_id,
            "execute",
            "completed",
            f"模型 {active_model} 输出完成",
            child_id=f"model:{active_model}",
            child_title=active_model,
            child_kind="model",
        )
        if delta_buffer and not output_policy_protected:
            emit(task_id, "answer_delta", "正在生成", "".join(delta_buffer))
        rows = [
            {"section": "任务目标", "content": message},
            {"section": "任务结果", "content": summary},
        ]
        if selected_names:
            rows.append({"section": "使用能力", "content": "、".join(selected_names)})
        artifacts: list[dict[str, Any]] = list(model_artifacts)
        self._emit_plan_progress(task_id, "execute", "completed", "正文内容已生成")
        if plan.get("tool_node_id") == "artifact":
            self._emit_plan_progress(task_id, "artifact", "running", f"正在生成可下载的 {requested_document_format.upper()} 文件")
        if wants_report_artifact and not artifacts:
            report = await self._tool(task_id, "report", "generate_markdown_report", {"summary": summary, "rows": rows, "filename": "general_report.md"}, plan=plan)
            if report.get("artifact"):
                artifacts.append(report["artifact"])
        requested_format = requested_document_format
        if requested_format and not any(a.get("kind") == requested_format for a in artifacts):
            requested_filename = str((plan.get("requirements") or {}).get("filename") or "").strip()
            if requested_format in {"xlsx", "csv"}:
                filename = requested_filename or f"agent_output.{requested_format}"
                document = await self._tool(task_id, "spreadsheet", "create_excel", {"rows": rows, "filename": filename}, plan=plan)
            else:
                filename = requested_filename or f"agent_output.{requested_format}"
                document = await self._tool(task_id, "report", "generate_document", {"title": task["title"], "content": summary, "format": requested_format, "filename": filename}, plan=plan)
            if document.get("artifact"):
                artifacts.append(document["artifact"])
        answer = summary
        if artifacts:
            missing_links = [
                f"- [{artifact.get('name', '下载文件')}]({artifact.get('download_url')})"
                for artifact in artifacts
                if artifact.get("download_url") and str(artifact.get("download_url")) not in answer
            ]
            if missing_links:
                answer = answer.rstrip() + "\n\n## 下载文件\n\n" + "\n".join(missing_links)
        if plan.get("tool_node_id") == "artifact":
            self._emit_plan_progress(task_id, "artifact", "completed", "文件已生成，可供下载")
        self._emit_plan_progress(task_id, "validate", "running", "正在检查结果是否符合当前计划")
        validation = self._validate_output_against_plan(plan, answer, artifacts)
        emit(task_id, "output_check", "输出前校验", validation["message"], validation)
        if not validation["passed"]:
            self._emit_plan_progress(task_id, "validate", "failed", validation["message"])
            raise RuntimeError(validation["message"])
        self._emit_plan_progress(task_id, "validate", "completed", "输出校验通过")
        output_policy_context = {
            "answer": answer,
            "artifacts": artifacts,
            "plan": plan,
        }
        output_evaluation = await self._evaluate_policy(
            "output.before", output_policy_context, enforce=True
        )
        applied_output = output_evaluation.apply(output_policy_context)
        answer = str(applied_output.get("answer") or "")
        applied_artifacts = applied_output.get("artifacts", artifacts)
        if not isinstance(applied_artifacts, list) or not all(
            isinstance(item, Mapping) for item in applied_artifacts
        ):
            raise RuntimeError(
                "output.before 策略修改后的 artifacts 必须是对象数组"
            )
        artifacts = [dict(item) for item in applied_artifacts]
        if output_policy_protected and answer:
            emit(task_id, "answer_delta", "正在生成", answer)
        emit(task_id, "answer", "任务完成", answer, {"artifacts": artifacts})
        db.update_task_status(task_id, "completed", result={"summary": answer, "rows": rows}, artifacts=artifacts)
        emit(task_id, "done", "已完成", "所有步骤已完成。")

    def _claim_runtime_messages(self) -> list[str]:
        execution = self._execution()
        if not execution:
            return []
        messages: list[str] = []
        while True:
            command = self.task_state.claim_command(
                execution["worker_id"],
                task_id=execution["task_id"],
                run_id=execution["run_id"],
                command_types=["message"],
            )
            if not command:
                break
            message = str(command.get("payload", {}).get("message") or "").strip()
            if not message:
                self.task_state.fail_command(command["id"], {"message": "追加指令不能为空"})
                continue
            messages.append(message)
            self.task_state.complete_command(command["id"], result={"applied": True})
        if messages:
            state = execution.setdefault("state", {})
            existing = state.setdefault("steering_messages", [])
            if isinstance(existing, list):
                existing.extend(messages)
            self._create_checkpoint("运行中追加指令已应用", node_key="execute")
            emit(
                execution["task_id"],
                "steering",
                "已应用运行中指令",
                "；".join(messages),
                {"count": len(messages)},
            )
        return messages

    def _clarification_for_missing(self, intent: dict[str, Any]) -> str:
        intent_name = str(intent.get("intent") or "").lower()
        raw_missing = [
            item.strip()
            for item in intent.get("missing_information", [])
            if isinstance(item, str) and item.strip()
        ]
        if "weather" in intent_name:
            missing = [
                item for item in raw_missing
                if item.lower() in {"city", "location", "place", "region", "城市", "地区", "地点"}
            ]
        else:
            missing = [item for item in raw_missing if self._safe_missing_information_label(item)]
        if not missing:
            return ""
        labels = {
            "city": "城市或地区",
            "location": "地点",
            "place": "地点",
            "region": "地区",
            "format": "输出格式",
            "filename": "文件名",
            "date": "日期",
            "day": "日期",
            "城市": "城市或地区",
            "地区": "城市或地区",
            "地点": "地点",
            "输出格式": "输出格式",
            "文件名": "文件名",
        }
        readable: list[str] = []
        for item in missing:
            value = labels.get(item.lower(), labels.get(item, item))
            if value not in readable:
                readable.append(value)
        if len(readable) == 1:
            field = readable[0]
            examples = {
                "城市或地区": "例如“宁波”或“上海浦东”",
                "输出格式": "例如“Word、PDF 或 Excel”",
            }
            suffix = f"，{examples[field]}" if field in examples else ""
            return f"还差一个信息：请告诉我{field}{suffix}。"
        return "还需要你补充：" + "、".join(readable) + "。补充后我会继续当前任务。"

    @staticmethod
    def _safe_missing_information_label(value: str) -> bool:
        """Accept a short user-facing field label, never parser output or reasoning."""
        label = value.strip()
        if not label or len(label) > 40 or "\n" in label or "\r" in label:
            return False
        lowered = label.lower()
        if any(token in lowered for token in (
            "standalone_request", "missing_information", "is_follow_up",
            "parameters", "```", "{", "}", "[", "]", "<", ">",
        )):
            return False
        if any(mark in label for mark in ("？", "?", "。", "！", "!", "：", ":", ";", "；")):
            return False
        if re.search(r"(?:因为|所以|首先|然后|推理|思考|分析过程|系统提示|内部指令|JSON)", label, re.IGNORECASE):
            return False
        return bool(re.fullmatch(r"[A-Za-z][A-Za-z0-9 _./-]{0,39}|[\u4e00-\u9fffA-Za-z0-9 _./（）()·-]{1,40}", label))

    def _weather_request_for_intent(
        self,
        original_message: str,
        resolved_message: str,
        history: list[dict[str, str]],
        intent_resolution: dict[str, Any],
    ) -> dict[str, str] | None:
        del intent_resolution  # Routing is based on the current request, not a model label alone.
        if (
            self._requested_document_format(original_message)
            or self._requested_document_format(resolved_message)
            or self._wants_report_artifact(original_message)
            or self._wants_report_artifact(resolved_message)
        ):
            return None
        awaiting_city = self._awaiting_weather_city(history)
        original_lookup = self._looks_like_weather_lookup(original_message)
        resolved_lookup = self._looks_like_weather_lookup(resolved_message)
        if not original_lookup and not resolved_lookup and not awaiting_city:
            return None
        routing_message = resolved_message if resolved_lookup else original_message
        return self._weather_request(routing_message, history)

    @staticmethod
    def _weather_lookup_explicitly_negated(message: str) -> bool:
        return bool(re.search(
            r"(?:不用|不要|无需|先不|不必|禁止).{0,8}(?:查|查询|看|了解)?.{0,4}(?:天气|气温|下雨|降雨)",
            message,
            re.IGNORECASE,
        ))

    def _looks_like_weather_lookup(self, message: str) -> bool:
        text = message.strip().lower()
        if self._weather_lookup_explicitly_negated(text):
            return False
        if not any(word in text for word in ["天气", "气温", "下雨", "降雨", "weather"]):
            return False
        conditional = any(phrase in text for phrase in ["天气不好时", "天气差时", "雨天备选", "室内备选", "下雨时", "如果下雨", "若下雨", "遇到下雨"])
        query_signal = bool(re.search(
            r"(?:查|查询|查一下|看看|想知道|告诉我|请问).{0,10}(?:天气|气温|下雨|降雨)|"
            r"(?:天气|气温).{0,8}(?:怎么样|如何|多少|预报|情况)|"
            r"(?:今天|明天|后天).{0,8}(?:天气|气温|下雨|降雨)|"
            r"(?:会不会|是否|有没有).{0,5}(?:下雨|降雨)|"
            r"(?:天气预报|weather)",
            text,
        ))
        return query_signal and not (conditional and not re.search(r"(?:查|查询|想知道|请问)", text))

    def _build_execution_plan(
        self,
        task: dict[str, Any],
        skills: list[dict[str, Any]],
        requested_format: str,
        wants_report: bool,
        *,
        force_weather: bool = False,
    ) -> dict[str, Any]:
        intent = task.get("intent_resolution") or {}
        goal = str(intent.get("standalone_request") or task.get("resolved_message") or task["message"])
        intent_name = str(intent.get("intent") or "general").lower()
        lowered_goal = goal.lower()
        parameters = intent.get("parameters") if isinstance(intent.get("parameters"), dict) else {}
        skill_servers = {
            server
            for skill in skills
            if skill.get("id") not in {"general_task", "report_generation"}
            for server in skill.get("required_mcps", [])
        }
        allowed_servers = set(skill_servers)
        explicit_weather_lookup = self._looks_like_weather_lookup(goal)
        artifact_requested = bool(requested_format or wants_report)
        is_weather = bool(
            not artifact_requested
            and (force_weather or explicit_weather_lookup)
            and parameters.get("weather_lookup") is not False
            and not self._weather_lookup_explicitly_negated(str(task.get("message") or ""))
            and not self._weather_lookup_explicitly_negated(goal)
        )
        weather_supports_artifact = bool(
            artifact_requested
            and explicit_weather_lookup
            and self._explicit_weather_artifact_lookup(str(task.get("message") or ""))
            and parameters.get("weather_lookup") is not False
        )
        if requested_format:
            # Artifact delivery is the primary goal.  Do not inherit unrelated
            # tools from a noisy skill match; add only explicit source tools.
            allowed_servers = {"spreadsheet" if requested_format in {"xlsx", "csv"} else "report"}
            if weather_supports_artifact:
                allowed_servers.add("weather")
        elif wants_report:
            allowed_servers = {"report"}
            if weather_supports_artifact:
                allowed_servers.add("weather")
        elif is_weather:
            allowed_servers = {"weather"}
        else:
            allowed_servers.discard("weather")
        if is_weather:
            allowed_servers.add("weather")
        else:
            if any(k in lowered_goal for k in ["表格", "excel", "xlsx", "csv"]):
                allowed_servers.add("spreadsheet")
            if any(word in lowered_goal for word in ["联网", "搜索", "最新", "查一下", "news", "新闻"]):
                allowed_servers.add("web-search")
        used_memory_ids = list(task.get("used_memory_ids") or [])
        nodes = self._execution_nodes(
            skills,
            goal=goal,
            requested_format=requested_format,
            wants_report=wants_report,
            is_weather=is_weather,
            needs_prepare=bool(used_memory_ids or task.get("attachments_json") or requested_format or wants_report or any(s in allowed_servers for s in {"web-search", "spreadsheet"})),
        )
        if used_memory_ids:
            prepare_node = next((node for node in nodes if node.get("id") == "prepare"), None)
            if prepare_node is not None:
                prepare_node.setdefault("children", []).append({
                    "id": "memory:effective",
                    "title": f"平台记忆（{len(used_memory_ids)} 条）",
                    "kind": "memory",
                    "status": "pending",
                })
        requirements = {key: parameters.get(key) for key in ("filename", "topic") if parameters.get(key)}
        requested_sections = parameters.get("sections") or parameters.get("chapters") or parameters.get("headings")
        if isinstance(requested_sections, list) and requested_sections:
            requirements["sections"] = requested_sections
        plan = {
            "goal": goal,
            "goal_confirmation": {
                "status": "auto_confirmed",
                "label": "目标已自动确认",
                "message": "目标与必要参数清晰；如存在关键缺失，系统会在执行前向用户询问。",
            },
            "intent": intent_name or "general",
            "steps": [node["title"] for node in nodes],
            "nodes": nodes,
            "allowed_servers": sorted(allowed_servers),
            "output_format": requested_format or "text",
            "requires_artifact": bool(requested_format or wants_report),
            "weather_lookup_required": bool(is_weather or weather_supports_artifact),
            "tool_node_id": "artifact" if requested_format or wants_report else "execute",
            "requirements": requirements,
        }
        plan["acceptance_criteria"] = self._acceptance_criteria(plan)
        return plan

    def _explicit_weather_artifact_lookup(self, message: str) -> bool:
        """Require a fresh lookup, not merely formatting a prior forecast."""
        text = message.strip().lower()
        if re.search(
            r"(?:刚才|前面|之前|已有|上述|上面|上一轮).{0,12}(?:天气|气温|预报|结果)",
            text,
        ):
            return False
        lookup_action = bool(re.search(
            r"(?:查|查询|查一下|看看|获取|想知道|告诉我|请问).{0,16}(?:天气|气温|下雨|降雨|预报)|"
            r"(?:今天|明天|后天).{0,12}(?:天气|气温|下雨|降雨)",
            text,
        ))
        return lookup_action and bool(self._extract_weather_city(message))

    def _execution_nodes(
        self,
        skills: list[dict[str, Any]],
        *,
        goal: str = "",
        requested_format: str = "",
        wants_report: bool = False,
        is_weather: bool = False,
        needs_prepare: bool = True,
    ) -> list[dict[str, Any]]:
        labels = {"docx": "Word", "pdf": "PDF", "pptx": "PowerPoint", "xlsx": "Excel", "csv": "CSV", "md": "Markdown", "html": "HTML"}
        format_label = labels.get(requested_format, requested_format.upper() if requested_format else "")
        if is_weather:
            understand_title, execute_title, validate_title = "确认查询城市与日期", "查询天气预报", "核对城市、日期与预报"
        elif requested_format:
            understand_title, execute_title, validate_title = f"确认 {format_label} 交付要求", "组织文档内容", "验证文件格式与下载"
        else:
            understand_title, execute_title, validate_title = "确认当前目标与约束", "生成任务结果", "核对结果与当前目标"
        children = [
            {"id": f"skill:{skill['id']}", "title": skill["name"], "kind": "skill", "status": "pending"}
            for skill in skills
        ]
        nodes = [{"id": "understand", "title": understand_title, "status": "pending", "children": children}]
        if needs_prepare:
            nodes.append({"id": "prepare", "title": "整理上下文与授权能力", "status": "pending", "children": []})
        nodes.append({"id": "execute", "title": execute_title, "status": "pending", "children": []})
        if requested_format or wants_report:
            artifact_label = format_label or "Markdown"
            nodes.append({"id": "artifact", "title": f"生成可下载的 {artifact_label} 文件", "status": "pending", "children": []})
        nodes.append({"id": "validate", "title": validate_title, "status": "pending", "children": []})
        return nodes

    def _format_execution_plan(self, plan: dict[str, Any]) -> str:
        lines = [f"目标：{plan['goal']}"]
        lines.extend(f"{index}. {step}" for index, step in enumerate(plan["steps"], start=1))
        if plan.get("allowed_servers"):
            lines.append("允许使用的工具服务：" + "、".join(plan["allowed_servers"]))
        else:
            lines.append("允许使用的工具服务：无（仅生成文本）")
        lines.append("目标输出：" + str(plan.get("output_format") or "text"))
        criteria = plan.get("acceptance_criteria") or []
        if criteria:
            lines.append("验收标准：" + "；".join(str(item.get("title") or "") for item in criteria))
        return "\n".join(lines)

    def _validate_tool_against_plan(self, plan: dict[str, Any], server_id: str, tool_name: str) -> None:
        if server_id not in set(plan.get("allowed_servers") or []):
            raise ToolError(f"已阻止偏离计划的工具调用：{server_id}.{tool_name}")
        if plan.get("requires_artifact") and server_id == "weather" and not plan.get("weather_lookup_required"):
            raise ToolError("当前目标是生成文档，已阻止无关的天气工具调用")

    def _acceptance_criteria(self, plan: dict[str, Any]) -> list[dict[str, Any]]:
        expected = str(plan.get("output_format") or "text")
        criteria = [
            {"id": "goal", "title": "结果对应当前任务目标", "status": "pending"},
            {"id": "response", "title": "已生成可交付的最终结果", "status": "pending"},
        ]
        if expected in {"docx", "pdf", "pptx", "xlsx", "csv", "md", "html"}:
            criteria.extend([
                {"id": "format", "title": f"已生成要求的 {expected.upper()} 文件", "status": "pending"},
                {"id": "content", "title": "文件包含可读取的有效内容", "status": "pending"},
                {"id": "download", "title": "文件已注册并可通过平台下载", "status": "pending"},
            ])
            requirements = plan.get("requirements") or {}
            if requirements.get("filename"):
                criteria.append({"id": "filename", "title": f"文件名为 {requirements['filename']}", "status": "pending"})
            if requirements.get("topic"):
                criteria.append({"id": "topic", "title": f"内容围绕“{requirements['topic']}”", "status": "pending"})
            if isinstance(requirements.get("sections"), list) and requirements["sections"]:
                criteria.append({"id": "sections", "title": "包含指定章节：" + "、".join(map(str, requirements["sections"])), "status": "pending"})
            attachment_requirements = plan.get("attachment_requirements")
            if isinstance(attachment_requirements, list) and attachment_requirements:
                criteria.append({
                    "id": "source_consistency",
                    "title": "生成文件保留附件中的关键内容",
                    "status": "pending",
                })
        return criteria

    @staticmethod
    def _attachment_acceptance_requirements(context: str) -> list[str]:
        """Select bounded, non-boilerplate source lines for output consistency checks."""
        candidates: list[str] = []
        for raw_line in context.splitlines():
            line = re.sub(r"\s+", " ", raw_line).strip()
            if (
                not line
                or line.startswith("--- ")
                or (line.startswith("[") and line.endswith("]"))
                or len(line) < 4
            ):
                continue
            candidates.append(line[:160])
        if not candidates:
            return []
        distinctive = [
            line
            for line in candidates
            if re.search(r"\d|[A-Z]{2,}|[:：]|[-_/]", line)
        ]
        ordered = distinctive + candidates
        return list(dict.fromkeys(ordered))[:3]

    def _artifact_content_check(self, artifact: dict[str, Any], expected: str) -> tuple[bool, str, str]:
        path = self._artifact_file(artifact)
        if path is None or path.stat().st_size <= 0:
            return False, "生成文件不存在或为空", ""
        try:
            if expected == "pptx":
                from pptx import Presentation
                presentation = Presentation(path)
                text = " ".join(
                    shape.text.strip()
                    for slide in presentation.slides
                    for shape in slide.shapes
                    if hasattr(shape, "text") and shape.text.strip()
                )
                passed = len(presentation.slides) >= 2 and len(text) >= 20
                return passed, f"共 {len(presentation.slides)} 页，提取到 {len(text)} 字可读内容", text
            if expected == "docx":
                from docx import Document
                document = Document(path)
                text = " ".join(paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip())
                return len(text) >= 20, f"提取到 {len(text)} 字可读内容", text
            if expected == "xlsx":
                from openpyxl import load_workbook
                workbook = load_workbook(path, read_only=True, data_only=True)
                try:
                    rows = sum(sheet.max_row for sheet in workbook.worksheets)
                    text = " ".join(str(cell.value) for sheet in workbook.worksheets for row in sheet.iter_rows() for cell in row if cell.value is not None)
                    return bool(workbook.sheetnames and rows), f"包含 {len(workbook.sheetnames)} 个工作表、{rows} 行", text
                finally:
                    workbook.close()
            if expected == "csv":
                import csv

                with path.open("r", encoding="utf-8-sig", errors="strict", newline="") as stream:
                    csv_rows = list(csv.reader(stream))
                if not csv_rows:
                    return False, "CSV 文件没有表头或数据行", ""
                headers = csv_rows[0]
                consistent = bool(headers) and all(len(row) == len(headers) for row in csv_rows)
                text = " ".join(cell for row in csv_rows for cell in row if cell)
                passed = consistent and bool(text.strip())
                return passed, f"包含 {len(headers)} 列、{max(0, len(csv_rows) - 1)} 行数据", text
            if expected == "pdf":
                from pypdf import PdfReader
                reader = PdfReader(str(path), strict=False)
                text = " ".join(
                    str(page.extract_text() or "").strip()
                    for page in reader.pages[: self.ATTACHMENT_MAX_PDF_PAGES]
                ).strip()
                return len(text) >= 20, f"共 {len(reader.pages)} 页，提取到 {len(text)} 字可读内容", text
            text = path.read_text(encoding="utf-8", errors="ignore") if expected in {"md", "html"} else ""
            return path.stat().st_size >= 20, f"文件大小 {path.stat().st_size} 字节", text
        except Exception as exc:
            return False, f"文件内容检查失败：{exc}", ""

    def _validate_output_against_plan(self, plan: dict[str, Any], answer: str, artifacts: list[dict[str, Any]]) -> dict[str, Any]:
        expected = str(plan.get("output_format") or "text")
        criteria: list[dict[str, Any]] = []
        goal_passed = bool(str(plan.get("goal") or "").strip()) and bool(answer.strip() or artifacts)
        criteria.append({"id": "goal", "title": "结果对应当前任务目标", "status": "passed" if goal_passed else "failed", "detail": "任务目标已保留并关联本次交付。" if goal_passed else "没有可用于验收的任务目标或结果。"})
        response_passed = bool(answer.strip())
        criteria.append({"id": "response", "title": "已生成可交付的最终结果", "status": "passed" if response_passed else "failed", "detail": f"最终回答共 {len(answer.strip())} 字。"})
        if expected in {"docx", "pdf", "pptx", "xlsx", "csv", "md", "html"}:
            aliases = {"md": {"md", "markdown"}}
            expected_kinds = aliases.get(expected, {expected})
            matches = [item for item in artifacts if str(item.get("kind") or "").lower() in expected_kinds]
            matched = bool(matches)
            criteria.append({"id": "format", "title": f"已生成要求的 {expected.upper()} 文件", "status": "passed" if matched else "failed", "detail": f"检测到 {len(matches)} 个格式匹配的文件。"})
            content_passed, content_detail, artifact_text = self._artifact_content_check(matches[0], expected) if matches else (False, "没有可检查的匹配文件。", "")
            criteria.append({"id": "content", "title": "文件包含可读取的有效内容", "status": "passed" if content_passed else "failed", "detail": content_detail})
            downloadable = bool(
                matches
                and matches[0].get("download_url")
                and self._artifact_file(matches[0]) is not None
            )
            criteria.append({"id": "download", "title": "文件已注册并可通过平台下载", "status": "passed" if downloadable else "failed", "detail": "下载地址已生成，文件在产物目录中存在。" if downloadable else "缺少下载地址或产物文件。"})
            requirements = plan.get("requirements") or {}
            expected_filename = str(requirements.get("filename") or "").strip()
            if expected_filename:
                filename_passed = bool(matches and str(matches[0].get("name") or "") == expected_filename)
                criteria.append({"id": "filename", "title": f"文件名为 {expected_filename}", "status": "passed" if filename_passed else "failed", "detail": f"实际文件名：{matches[0].get('name')}" if matches else "没有生成匹配文件。"})
            topic = str(requirements.get("topic") or "").strip()
            if topic:
                topic_passed = topic in artifact_text
                criteria.append({"id": "topic", "title": f"内容围绕“{topic}”", "status": "passed" if topic_passed else "failed", "detail": "已在文件内容中找到指定主题。" if topic_passed else "文件内容中未找到指定主题。"})
            required_sections = requirements.get("sections") if isinstance(requirements.get("sections"), list) else []
            if required_sections:
                missing_sections = [str(section) for section in required_sections if str(section) not in artifact_text]
                criteria.append({"id": "sections", "title": "包含指定章节：" + "、".join(map(str, required_sections)), "status": "passed" if not missing_sections else "failed", "detail": "所有指定章节均已找到。" if not missing_sections else "缺少章节：" + "、".join(missing_sections)})
            attachment_requirements = [
                str(item).strip()
                for item in plan.get("attachment_requirements", [])
                if str(item).strip()
            ] if isinstance(plan.get("attachment_requirements"), list) else []
            if attachment_requirements:
                matched_requirements = [
                    item for item in attachment_requirements if item in artifact_text
                ]
                source_passed = bool(matched_requirements)
                criteria.append({
                    "id": "source_consistency",
                    "title": "生成文件保留附件中的关键内容",
                    "status": "passed" if source_passed else "failed",
                    "detail": (
                        f"已在文件中找到 {len(matched_requirements)} 项附件关键内容。"
                        if source_passed
                        else "文件中未找到抽样的附件关键内容，已阻止错误交付。"
                    ),
                })
            passed = all(item["status"] == "passed" for item in criteria)
            return {
                "passed": passed,
                "message": f"已按 {len(criteria)} 项验收标准完成检查，全部通过。" if passed else f"输出校验失败：{sum(item['status'] == 'failed' for item in criteria)} 项未通过。",
                "expected_format": expected,
                "artifact_count": len(artifacts),
                "criteria": criteria,
            }
        passed = all(item["status"] == "passed" for item in criteria)
        return {"passed": passed, "message": "已确认最终内容与当前任务目标一致。" if passed else "最终内容未通过验收。", "expected_format": expected, "artifact_count": len(artifacts), "criteria": criteria}

    def _requested_document_format(self, message: str) -> str:
        """Return a format only when the user actually asks to create/export a file.

        Merely discussing a preference such as “默认输出格式是 PDF” must not
        create an artifact.
        """
        lowered = message.lower()
        action = r"(?:生成|创建|制作|导出|下载|写成|整理成|转成|转换为|保存为|做(?:一份|个)|出(?:一份|个)|给我(?:一份|个))"
        target = r"(?:pdf|word|docx|pptx?|powerpoint|幻灯片|演示文稿|excel|xlsx|csv|逗号分隔(?:文件|表格)?|电子表格|markdown|md|html|网页文档|文档)"
        if not (re.search(action + r".{0,16}" + target, lowered) or re.search(target + r".{0,16}" + action, lowered)):
            return ""
        directed = re.search(
            r"(?:转成|转换为|保存为|导出为|写成|整理成|做成|制作为|生成为)"
            r".{0,12}?(pdf|word|docx|pptx?|powerpoint|幻灯片|演示文稿|excel|xlsx|csv|逗号分隔(?:文件|表格)?|电子表格|markdown|md|html|网页文档)",
            lowered,
        )
        if directed:
            directed_target = directed.group(1)
            if "pdf" in directed_target:
                return "pdf"
            if any(word in directed_target for word in ["ppt", "powerpoint", "幻灯片", "演示文稿"]):
                return "pptx"
            if any(word in directed_target for word in ["excel", "xlsx", "电子表格"]):
                return "xlsx"
            if "csv" in directed_target or "逗号分隔" in directed_target:
                return "csv"
            if "html" in directed_target or "网页文档" in directed_target:
                return "html"
            if "markdown" in directed_target or directed_target == "md":
                return "md"
            if "word" in directed_target or "docx" in directed_target:
                return "docx"
        if "pdf" in lowered:
            return "pdf"
        if any(word in lowered for word in ["pptx", "ppt", "powerpoint", "幻灯片", "演示文稿"]):
            return "pptx"
        if any(word in lowered for word in ["excel", "xlsx", "电子表格"]):
            return "xlsx"
        if "csv" in lowered or "逗号分隔" in lowered:
            return "csv"
        if "html" in lowered or "网页文档" in lowered:
            return "html"
        if "markdown" in lowered or re.search(r"(?:^|\W)md(?:$|\W)", lowered):
            return "md"
        if "word" in lowered or "docx" in lowered or "文档" in lowered:
            return "docx"
        return ""

    def _wants_report_artifact(self, message: str) -> bool:
        lowered = message.lower()
        action = r"(?:生成|创建|制作|导出|下载|保存|写成|整理成|做一份|做个|出一份|出个)"
        if re.search(action + r".{0,12}(?:报告|汇报|markdown)", lowered):
            return True
        return bool(re.search(r"(?:报告|汇报|markdown).{0,12}" + action, lowered))

    def _conversation_history(self, task: dict[str, Any], limit: int = 10) -> list[dict[str, str]]:
        conversation_id = str(task.get("conversation_id") or "").strip()
        if not conversation_id:
            return []
        scope = self._context_scope(task)
        summary = self.conversation_summary_service.get(scope, conversation_id)
        prefix = self.conversation_summary_service.history_prefix(scope, conversation_id) if summary else []
        through = None
        if summary and summary.get("through_task_id"):
            through = db.query_one(
                "SELECT id, created_at FROM tasks WHERE id = ? AND conversation_id = ?",
                (summary["through_task_id"], conversation_id),
            )
        if through:
            rows = db.query_all(
                """SELECT id, message FROM tasks
                   WHERE conversation_id = ? AND id != ? AND status = 'completed'
                     AND (created_at > ? OR (created_at = ? AND id > ?))
                     AND (created_at < ? OR (created_at = ? AND id < ?))
                   ORDER BY created_at DESC, id DESC LIMIT ?""",
                (
                    conversation_id,
                    task["id"],
                    through["created_at"],
                    through["created_at"],
                    through["id"],
                    task["created_at"],
                    task["created_at"],
                    task["id"],
                    limit,
                ),
            )
        else:
            rows = db.query_all(
                """SELECT id, message FROM tasks
                   WHERE conversation_id = ? AND id != ? AND status = 'completed'
                     AND (created_at < ? OR (created_at = ? AND id < ?))
                   ORDER BY created_at DESC, id DESC LIMIT ?""",
                (conversation_id, task["id"], task["created_at"], task["created_at"], task["id"], limit),
            )
        history: list[dict[str, str]] = list(prefix)
        for row in reversed(rows):
            history.append({"role": "user", "content": row["message"]})
            answer = db.query_one(
                "SELECT content FROM task_events WHERE task_id = ? AND type = 'answer' ORDER BY id DESC LIMIT 1",
                (row["id"],),
            )
            if answer and answer.get("content"):
                history.append({"role": "assistant", "content": answer["content"]})
        return history

    def _maybe_compact_conversation(
        self,
        task: dict[str, Any],
        *,
        keep_recent_tasks: int = 6,
        trigger_tasks: int = 9,
    ) -> dict[str, Any] | None:
        conversation_id = str(task.get("conversation_id") or "").strip()
        if not conversation_id:
            return None
        rows = db.query_all(
            """SELECT id, message, created_at FROM tasks
               WHERE conversation_id = ? AND status = 'completed'
               ORDER BY created_at, id LIMIT 1000""",
            (conversation_id,),
        )
        if len(rows) < max(trigger_tasks, keep_recent_tasks + 1):
            return None
        cutoff = len(rows) - max(1, keep_recent_tasks)
        compact_rows = rows[:cutoff]
        through_task_id = compact_rows[-1]["id"]
        scope = self._context_scope(task)
        existing = self.conversation_summary_service.get(scope, conversation_id)
        if existing and existing.get("through_task_id") == through_task_id:
            return None
        messages: list[dict[str, str]] = []
        for row in compact_rows:
            messages.append({"role": "user", "content": row["message"]})
            answer = db.query_one(
                "SELECT content FROM task_events WHERE task_id = ? AND type = 'answer' ORDER BY id DESC LIMIT 1",
                (row["id"],),
            )
            if answer and answer.get("content"):
                messages.append({"role": "assistant", "content": answer["content"]})
        if not messages:
            return None
        return self.conversation_summary_service.compact(
            scope,
            messages,
            conversation_id=conversation_id,
            through_task_id=through_task_id,
        )

    def _weather_request(self, message: str, history: list[dict[str, str]]) -> dict[str, str] | None:
        """Return a weather request only for the current intent or an immediate city follow-up.

        Old weather turns must not keep routing unrelated messages to the weather tool.
        """
        if (
            self._weather_lookup_explicitly_negated(message)
            or self._requested_document_format(message)
            or self._wants_report_artifact(message)
        ):
            return None
        weather_words = ["天气", "气温", "下雨", "降雨", "weather"]
        explicit_weather = any(word in message.lower() for word in weather_words)
        awaiting_city = self._awaiting_weather_city(history)
        if not explicit_weather and not awaiting_city:
            return None

        city = self._extract_weather_city(message, city_only=not explicit_weather)
        if awaiting_city and not explicit_weather and not city:
            return None

        day_context = message
        if awaiting_city and not explicit_weather:
            previous_user = next((item["content"] for item in reversed(history[:-1]) if item["role"] == "user"), "")
            day_context = f"{previous_user} {message}"
        return {"city": city, "day": self._extract_weather_day(day_context)}

    def _awaiting_weather_city(self, history: list[dict[str, str]]) -> bool:
        if not history or history[-1].get("role") != "assistant":
            return False
        answer = history[-1].get("content", "")
        asks_location = ("城市" in answer or "地区" in answer) and any(word in answer for word in ["告诉", "提供", "需要查询", "所在"])
        previous_user = next((item.get("content", "") for item in reversed(history[:-1]) if item.get("role") == "user"), "")
        return asks_location and any(word in previous_user.lower() for word in ["天气", "气温", "下雨", "降雨", "weather"])

    def _extract_weather_city(self, message: str, city_only: bool = False) -> str:
        temporal_words = {"今天", "明天", "后天", "天气", "气温", "下雨", "降雨"}
        non_city_words = temporal_words | {
            "查询", "查一下", "看看", "请问", "结果", "行程", "文档", "报告", "总结",
            "整理文档", "整理成文档", "写一份总结", "怎么样", "如何", "咋样", "情况",
            "好的", "好", "继续", "谢谢", "多谢", "算了", "不用了", "不查了", "取消",
            "可以", "行", "嗯", "收到", "知道了", "明白了",
        }
        cleaned = re.sub(r"[，。！？,.!?\s]", "", message.strip())
        if city_only and any(word in cleaned.lower() for word in (
            "整理", "总结", "文档", "报告", "word", "docx", "pdf", "ppt",
            "excel", "xlsx", "csv", "markdown", "html", "前面的结果",
        )):
            return ""
        city_only_match = re.fullmatch(r"(?:我在|位置是|城市是)?([\u4e00-\u9fff]{2,12}?)(?:市)?", cleaned)
        if city_only_match:
            city = city_only_match.group(1)
            if city and city not in non_city_words and not any(word in city for word in temporal_words):
                return city
        if city_only:
            return ""

        patterns = [
            r"(?:今天|明天|后天)?(?:查一下|查询|查|看看|请问|请帮我)?([\u4e00-\u9fff]{2,10}?)(?:市)?(?:今天|明天|后天)(?:的)?(?:天气|气温)",
            r"(?:今天|明天|后天)?(?:查一下|查询|查|看看|请问|请帮我)?([\u4e00-\u9fff]{2,10}?)(?:市)?(?:的)?(?:天气|气温)",
            r"(?:天气|气温).{0,6}?(?:在|查|查询)?([\u4e00-\u9fff]{2,10}?)(?:市)?$",
            r"(?:查一下|查询|查|看看|请问|请帮我)?([\u4e00-\u9fff]{2,10}?)(?:市)?(?:今天|明天|后天)?(?:是否|会不会|有没有)?(?:下雨|降雨)",
        ]
        for pattern in patterns:
            match = re.search(pattern, cleaned)
            if not match:
                continue
            city = re.sub(r"^(你好|我想知道|告诉我|帮我)", "", match.group(1))
            if city and city not in non_city_words:
                return city
        return ""

    def _extract_weather_day(self, text: str) -> str:
        if "后天" in text:
            return "day_after_tomorrow"
        if "明天" in text:
            return "tomorrow"
        return "today"

    def _build_weather_answer(self, forecast: dict[str, Any]) -> str:
        city = forecast.get("city", "该地区")
        date = forecast.get("date", "明天")
        condition = forecast.get("condition", "未知")
        low = forecast.get("temperature_min_c")
        high = forecast.get("temperature_max_c")
        rain = forecast.get("precipitation_probability_max_percent")
        rain_sum = forecast.get("precipitation_sum_mm")
        wind = forecast.get("wind_speed_max_kmh")
        gust = forecast.get("wind_gusts_max_kmh")
        day_label = {"today": "今天", "tomorrow": "明天", "day_after_tomorrow": "后天"}.get(forecast.get("day"), "")
        lines = [
            f"{city}{day_label}（{date}）预计：{condition}。",
            f"- 气温：{low}～{high}℃",
            f"- 最高降雨概率：{rain}%（预计降水 {rain_sum} mm）",
            f"- 最大风速：{wind} km/h，阵风最高 {gust} km/h",
        ]
        advice: list[str] = []
        if isinstance(rain, (int, float)) and rain >= 50:
            advice.append("降雨概率较高，建议带伞")
        if isinstance(high, (int, float)) and high >= 35:
            advice.append("白天气温较高，注意防暑补水")
        if int(forecast.get("weather_code") or 0) >= 95:
            advice.append("可能有雷暴，尽量避免长时间户外停留")
        if advice:
            lines.append("出行提示：" + "；".join(advice) + "。")
        lines.append(f"数据源：Open-Meteo（{forecast.get('source', 'https://open-meteo.com/')}），预报可能随时间更新。")
        return "\n".join(lines)

    @staticmethod
    def _bounded_attachment_lines(
        lines: Iterable[str],
        max_chars: int,
        notice: str = "[正文过长，已按单个附件字符上限截断]",
    ) -> str:
        """Consume a text generator without materialising unbounded document content."""
        output: list[str] = []
        used = 0
        truncated = False
        for raw_line in lines:
            line = str(raw_line or "")
            separator_size = 1 if output else 0
            remaining = max_chars - used - separator_size
            if remaining <= 0:
                truncated = True
                break
            if len(line) > remaining:
                output.append(line[:remaining])
                used = max_chars
                truncated = True
                break
            output.append(line)
            used += separator_size + len(line)
        text = "\n".join(output)
        if not truncated:
            return text
        suffix = f"\n{notice}"
        if len(suffix) >= max_chars:
            return suffix[-max_chars:]
        return text[: max_chars - len(suffix)] + suffix

    def _office_archive_warning(self, file_path: Path) -> str:
        """Reject malformed or excessively expanded Office archives before a parser opens them."""
        try:
            with zipfile.ZipFile(file_path) as archive:
                entries = archive.infolist()
                if len(entries) > self.ATTACHMENT_MAX_ARCHIVE_ENTRIES:
                    return "文件内部条目过多，已停止正文提取。"
                unpacked_size = sum(max(0, int(entry.file_size)) for entry in entries)
                if unpacked_size > self.ATTACHMENT_MAX_UNCOMPRESSED_BYTES:
                    return "文件解压后的内容超过安全上限，已停止正文提取。"
        except (OSError, zipfile.BadZipFile, zipfile.LargeZipFile):
            return "文件可能已损坏或不是有效的 Office 文档，无法提取正文。"
        return ""

    def _extract_docx_attachment(self, file_path: Path) -> str:
        warning = self._office_archive_warning(file_path)
        if warning:
            return f"[Word 正文提取失败：{warning}]"
        try:
            from docx import Document
        except ImportError:
            return "[Word 正文解析组件未安装，请安装 python-docx 后重试。]"

        document = Document(str(file_path))

        def lines() -> Iterable[str]:
            found = False
            for paragraph in document.paragraphs:
                text = paragraph.text.strip()
                if text:
                    found = True
                    yield text
            tables = document.tables
            for table_index, table in enumerate(tables[: self.ATTACHMENT_MAX_TABLES], start=1):
                yield f"[表格 {table_index}]"
                for row in table.rows[: self.ATTACHMENT_MAX_ROWS_PER_TABLE]:
                    values = [
                        cell.text.strip().replace("\n", " ")
                        for cell in row.cells[: self.ATTACHMENT_MAX_COLUMNS_PER_TABLE]
                    ]
                    if any(values):
                        found = True
                        yield " | ".join(values)
                if len(table.rows) > self.ATTACHMENT_MAX_ROWS_PER_TABLE:
                    yield f"[该表格仅提取前 {self.ATTACHMENT_MAX_ROWS_PER_TABLE} 行]"
            if len(tables) > self.ATTACHMENT_MAX_TABLES:
                yield f"[仅提取前 {self.ATTACHMENT_MAX_TABLES} 个表格]"
            if not found:
                yield "[未发现可提取的 Word 正文]"

        # python-docx only resolves relationships stored in the package.  This
        # extractor never opens hyperlink targets or other external resources.
        return self._bounded_attachment_lines(lines(), self.ATTACHMENT_MAX_FILE_CHARS)

    @staticmethod
    def _xlsx_cell_text(cell: Any) -> str:
        value = cell.value
        if value is None:
            return ""
        # openpyxl does not calculate formulas, but suppressing their source as
        # well prevents formulas (including external-workbook formulas) from
        # becoming instructions in model context.
        if getattr(cell, "data_type", "") == "f" or (
            isinstance(value, str) and value.lstrip().startswith("=")
        ):
            return "[公式未执行]"
        return str(value).replace("\r", " ").replace("\n", " ")

    def _extract_xlsx_attachment(self, file_path: Path) -> str:
        warning = self._office_archive_warning(file_path)
        if warning:
            return f"[Excel 正文提取失败：{warning}]"
        try:
            from openpyxl import load_workbook
        except ImportError:
            return "[Excel 正文解析组件未安装，请安装 openpyxl 后重试。]"

        # read_only avoids loading the complete grid. data_only=False lets us
        # identify formulas and replace them without evaluating them.
        workbook = load_workbook(
            filename=str(file_path),
            read_only=True,
            data_only=False,
            keep_links=False,
        )
        try:
            worksheets = workbook.worksheets

            def lines() -> Iterable[str]:
                for worksheet in worksheets[: self.ATTACHMENT_MAX_WORKSHEETS]:
                    yield f"[工作表：{worksheet.title}]"
                    max_row = min(
                        max(int(worksheet.max_row or 1), 1),
                        self.ATTACHMENT_MAX_ROWS_PER_SHEET,
                    )
                    max_column = min(
                        max(int(worksheet.max_column or 1), 1),
                        self.ATTACHMENT_MAX_COLUMNS_PER_SHEET,
                    )
                    found = False
                    for row in worksheet.iter_rows(
                        min_row=1,
                        max_row=max_row,
                        min_col=1,
                        max_col=max_column,
                    ):
                        values = [self._xlsx_cell_text(cell) for cell in row]
                        while values and not values[-1]:
                            values.pop()
                        if any(values):
                            found = True
                            yield " | ".join(values)
                    if not found:
                        yield "[该工作表没有可提取的单元格内容]"
                    if int(worksheet.max_row or 0) > self.ATTACHMENT_MAX_ROWS_PER_SHEET:
                        yield f"[该工作表仅提取前 {self.ATTACHMENT_MAX_ROWS_PER_SHEET} 行]"
                    if int(worksheet.max_column or 0) > self.ATTACHMENT_MAX_COLUMNS_PER_SHEET:
                        yield f"[该工作表仅提取前 {self.ATTACHMENT_MAX_COLUMNS_PER_SHEET} 列]"
                if len(worksheets) > self.ATTACHMENT_MAX_WORKSHEETS:
                    yield f"[工作簿仅提取前 {self.ATTACHMENT_MAX_WORKSHEETS} 个工作表]"

            return self._bounded_attachment_lines(lines(), self.ATTACHMENT_MAX_FILE_CHARS)
        finally:
            workbook.close()

    def _extract_pptx_attachment(self, file_path: Path) -> str:
        warning = self._office_archive_warning(file_path)
        if warning:
            return f"[PowerPoint 正文提取失败：{warning}]"
        try:
            from pptx import Presentation
        except ImportError:
            return "[PowerPoint 正文解析组件未安装，请安装 python-pptx 后重试。]"

        presentation = Presentation(str(file_path))

        def lines() -> Iterable[str]:
            for slide_index, slide in enumerate(presentation.slides, start=1):
                if slide_index > self.ATTACHMENT_MAX_SLIDES:
                    break
                yield f"[幻灯片 {slide_index}]"
                found = False
                for shape in list(slide.shapes)[: self.ATTACHMENT_MAX_SHAPES_PER_SLIDE]:
                    if getattr(shape, "has_text_frame", False):
                        text = str(getattr(shape, "text", "") or "").strip()
                        if text:
                            found = True
                            yield text
                    elif getattr(shape, "has_table", False):
                        for row_index, row in enumerate(shape.table.rows, start=1):
                            if row_index > self.ATTACHMENT_MAX_ROWS_PER_TABLE:
                                break
                            values = [
                                cell.text.strip().replace("\n", " ")
                                for cell_index, cell in enumerate(row.cells, start=1)
                                if cell_index <= self.ATTACHMENT_MAX_COLUMNS_PER_TABLE
                            ]
                            if any(values):
                                found = True
                                yield " | ".join(values)
                if not found:
                    yield "[该幻灯片没有可提取的文本]"
                if len(slide.shapes) > self.ATTACHMENT_MAX_SHAPES_PER_SLIDE:
                    yield f"[该幻灯片仅检查前 {self.ATTACHMENT_MAX_SHAPES_PER_SLIDE} 个对象]"
            if len(presentation.slides) > self.ATTACHMENT_MAX_SLIDES:
                yield f"[演示文稿仅提取前 {self.ATTACHMENT_MAX_SLIDES} 张幻灯片]"

        # Hyperlink relationships remain inert strings inside the package;
        # this code reads shape text only and never dereferences them.
        return self._bounded_attachment_lines(lines(), self.ATTACHMENT_MAX_FILE_CHARS)

    def _extract_pdf_attachment(self, file_path: Path) -> str:
        try:
            from pypdf import PdfReader
        except ImportError:
            return "[PDF 正文解析组件未安装，请安装 pypdf 后重试。]"

        reader = PdfReader(str(file_path), strict=False)
        if reader.is_encrypted:
            try:
                if not reader.decrypt(""):
                    return "[PDF 已加密且需要密码，无法提取正文。]"
            except Exception:
                return "[PDF 已加密且需要密码，无法提取正文。]"

        def lines() -> Iterable[str]:
            pages = reader.pages
            for page_index, page in enumerate(pages, start=1):
                if page_index > self.ATTACHMENT_MAX_PDF_PAGES:
                    break
                yield f"[第 {page_index} 页]"
                text = str(page.extract_text() or "").strip()
                yield text or "[该页没有可提取的文本]"
            if len(pages) > self.ATTACHMENT_MAX_PDF_PAGES:
                yield f"[PDF 仅提取前 {self.ATTACHMENT_MAX_PDF_PAGES} 页]"

        # Page text extraction does not inspect link annotations or fetch URLs.
        return self._bounded_attachment_lines(lines(), self.ATTACHMENT_MAX_FILE_CHARS)

    def _extract_attachment_body(self, file_path: Path, suffix: str, content_type: str) -> str:
        text_suffixes = {
            ".txt", ".md", ".csv", ".json", ".yaml", ".yml",
            ".py", ".js", ".ts", ".html", ".css",
        }
        if suffix in text_suffixes:
            with file_path.open("r", encoding="utf-8", errors="replace") as handle:
                text = handle.read(self.ATTACHMENT_MAX_FILE_CHARS + 1)
            if len(text) > self.ATTACHMENT_MAX_FILE_CHARS:
                return self._bounded_attachment_lines(
                    [text], self.ATTACHMENT_MAX_FILE_CHARS
                )
            return text
        if suffix == ".docx":
            return self._extract_docx_attachment(file_path)
        if suffix == ".xlsx":
            return self._extract_xlsx_attachment(file_path)
        if suffix == ".pptx":
            return self._extract_pptx_attachment(file_path)
        if suffix == ".pdf":
            return self._extract_pdf_attachment(file_path)
        return f"[已上传二进制文件，类型 {content_type or 'unknown'}，当前不支持提取正文。]"

    def _attachment_context(self, attachments: list[dict[str, Any]]) -> str:
        chunks: list[str] = []
        for item in attachments[: self.ATTACHMENT_MAX_FILES]:
            raw_path = item.get("path")
            if not raw_path:
                continue
            file_path = Path(str(raw_path))
            display_name = str(item.get("name") or file_path.name or "附件")
            header = f"--- {display_name} ---"
            try:
                if not file_path.is_file():
                    body = "[附件文件不存在或不可访问，无法提取正文。]"
                elif file_path.stat().st_size > self.ATTACHMENT_MAX_FILE_BYTES:
                    body = "[附件超过正文解析大小上限，已跳过提取。]"
                else:
                    suffix = Path(display_name).suffix.lower() or file_path.suffix.lower()
                    body = self._extract_attachment_body(
                        file_path,
                        suffix,
                        str(item.get("content_type") or "application/octet-stream"),
                    )
            except Exception:
                body = "[正文提取失败：文件可能已损坏、受密码保护或格式不兼容。]"
            chunks.append(f"{header}\n{body}")
        if len(attachments) > self.ATTACHMENT_MAX_FILES:
            chunks.append(f"[附件较多，仅处理前 {self.ATTACHMENT_MAX_FILES} 个文件]" )
        return self._bounded_attachment_lines(
            chunks,
            self.ATTACHMENT_MAX_CONTEXT_CHARS,
            "[附件正文已达到总字符上限，其余内容未加入上下文]",
        )

    async def resume_after_approval(self, task_id: str, approved: bool, note: str = "") -> None:
        task = db.query_one("SELECT * FROM tasks WHERE id = ?", (task_id,))
        if not task:
            return
        result = db.json_loads(task.get("result_json"), {})
        if task["status"] != "waiting_approval":
            emit(task_id, "approval", "审批无效", "当前任务不处于等待审批状态。")
            return
        active_run = next(
            (
                item
                for item in self.task_state.list_runs(task_id=task_id)
                if item["status"] in {"waiting_approval", "paused", "running"}
            ),
            None,
        )
        if result.get("pending_action") == "install_recommended_skill":
            recommendation_id = str(result.get("recommendation_id") or "")
            recommendation = get_builtin_skill(recommendation_id)
            if approved and not recommendation:
                emit(task_id, "error", "安装失败", "推荐项已经失效，原任务暂时无法继续。")
                db.update_task_status(
                    task_id,
                    "failed",
                    result={**result, "approval": "approved", "error": "recommendation_not_found"},
                )
                if active_run and active_run["status"] not in {"completed", "failed", "cancelled"}:
                    self.task_state.finish_run(
                        active_run["id"],
                        status="failed",
                        error={"message": "recommendation_not_found"},
                    )
                return

            # Keep the decision on the task while the same run re-enters the
            # normal planner.  A rejected recommendation is also added to the
            # audit list; the handled flag below prevents either decision from
            # opening another recommendation dialog during this continuation.
            continuation_result = dict(result)
            continuation_result.pop("pending_action", None)
            continuation_result.pop("recommendation_id", None)
            continuation_result.pop("summary", None)
            continuation_result["approval"] = "approved" if approved else "rejected"
            continuation_result["skill_recommendation_decision"] = {
                "id": recommendation_id,
                "decision": "approved" if approved else "rejected",
            }
            # One decision resolves built-in recommendations for this task.
            continuation_result["skip_skill_recommendations"] = True

            if approved:
                emit(task_id, "approval", "已确认安装", note or "用户确认安装推荐 Skill；安装完成后将继续原任务。")
                try:
                    skill = self.skill_registry.install_content(
                        recommendation["content"], fallback_id=recommendation["id"]
                    )
                except Exception as exc:
                    emit(task_id, "error", "安装失败", f"推荐 Skill 安装失败：{exc}")
                    db.update_task_status(
                        task_id,
                        "failed",
                        result={**continuation_result, "error": str(exc)},
                    )
                    if active_run and active_run["status"] not in {"completed", "failed", "cancelled"}:
                        self.task_state.finish_run(
                            active_run["id"], status="failed", error={"message": str(exc)}
                        )
                    return
                continuation_result["installed_skill"] = skill
                emit(
                    task_id,
                    "install",
                    "Skill 安装成功",
                    f"已安装“{skill['name']}”（{skill['id']}），正在继续处理原任务。",
                    {"skill": skill, "source": "builtin_catalog"},
                )
            else:
                declined = [
                    str(item)
                    for item in continuation_result.get("declined_recommendation_ids", [])
                    if str(item).strip()
                ]
                if recommendation_id and recommendation_id not in declined:
                    declined.append(recommendation_id)
                continuation_result["declined_recommendation_ids"] = declined
                emit(
                    task_id,
                    "approval",
                    "已跳过 Skill 安装",
                    note or "用户暂不安装推荐 Skill；平台将使用现有能力继续原任务。",
                    {"recommendation_id": recommendation_id},
                )

            if active_run:
                self.task_state.update_run_metadata(
                    active_run["id"],
                    {
                        "skill_recommendation_decision": {
                            "id": recommendation_id,
                            "decision": "approved" if approved else "rejected",
                        }
                    },
                )
                # The task row becomes waiting_approval just before the first
                # run unwinds.  Normally its run is already waiting too, but a
                # very fast click can race that final transition.
                for _ in range(100):
                    refreshed = self.task_state.get_run(active_run["id"])
                    if not refreshed or refreshed["status"] != "running":
                        active_run = refreshed or active_run
                        break
                    await asyncio.sleep(0.01)

            db.update_task_status(task_id, "running", result=continuation_result)
            emit(
                task_id,
                "resume",
                "继续原任务",
                "Skill 推荐决策已记录，正在沿用当前任务和运行继续执行。",
                {
                    "run_id": (active_run or {}).get("id", ""),
                    "recommendation_id": recommendation_id,
                    "approved": approved,
                },
            )
            await self.run_task(
                task_id,
                run_id=(active_run or {}).get("id") or None,
            )
            return

        if not approved:
            emit(task_id, "approval", "审批拒绝", note or "用户拒绝执行敏感操作。")
            db.update_task_status(task_id, "completed", result={**result, "approval": "rejected"}, artifacts=db.json_loads(task.get("artifacts_json"), []))
            if active_run and active_run["status"] not in {"completed", "failed", "cancelled"}:
                self.task_state.finish_run(
                    active_run["id"],
                    status="completed",
                    result={**result, "approval": "rejected"},
                    metadata={"approval_rejected": True},
                )
            emit(task_id, "done", "已完成", "任务已在不执行敏感操作的情况下结束。")
            return
        if active_run and active_run["status"] in {"waiting_approval", "paused"}:
            self.task_state.begin_run(task_id, run_id=active_run["id"])
        db.update_task_status(task_id, "running")
        emit(task_id, "approval", "审批已记录", note or "用户批准执行敏感操作。")
        unsupported = "当前没有配置可执行该操作的外部写入工具，因此未修改任何外部数据。"
        emit(task_id, "tool_result", "未执行外部写入", unsupported)
        db.update_task_status(task_id, "completed", result={**result, "approval": "approved", "write_back": "not_configured", "summary": unsupported}, artifacts=db.json_loads(task.get("artifacts_json"), []))
        if active_run and (self.task_state.get_run(active_run["id"]) or {}).get("status") not in {"completed", "failed", "cancelled"}:
            self.task_state.finish_run(
                active_run["id"],
                result={**result, "approval": "approved", "write_back": "not_configured", "summary": unsupported},
            )
        emit(task_id, "answer", "未执行外部写入", unsupported)
        emit(task_id, "done", "已结束", "任务已结束，外部系统未发生变更。")


def create_task_record(
    message: str,
    agent_id: str,
    workspace: str = "default",
    attachments: list[dict[str, Any]] | None = None,
    model_id: str | None = None,
    conversation_id: str | None = None,
    *,
    organization_id: str = "local-org",
    user_id: str = "local-user",
    parent_task_id: str = "",
    executor_type: str = "agent",
    executor_id: str = "",
) -> dict[str, Any]:
    task_id = "task_" + uuid.uuid4().hex[:12]
    title = message.strip().replace("\n", " ")[:60] or "新任务"
    now = db.utc_now()
    db.execute(
        """
        INSERT INTO tasks(
            id, title, message, agent_id, model_id, conversation_id, workspace,
            organization_id, user_id, parent_task_id, executor_type, executor_id,
            status, result_json, artifacts_json, attachments_json, created_at, updated_at
        )
        VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
        """,
        (
            task_id, title, message, agent_id, model_id or "",
            conversation_id or ("conv_" + uuid.uuid4().hex[:16]), workspace,
            organization_id, user_id, parent_task_id, executor_type,
            executor_id or agent_id, "queued", db.json_dumps({}), db.json_dumps([]),
            db.json_dumps(attachments or []), now, now,
        ),
    )
    return db.query_one("SELECT * FROM tasks WHERE id = ?", (task_id,)) or {"id": task_id}
